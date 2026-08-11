#!/usr/bin/env python3
"""
TargetInfo -- 靶点调研报告生成器（PyQt5 桌面版）
移植自 target.html，消除 CORS 限制，集成更多专业数据库。

数据源：
  文献: PubMed, OpenAlex, Semantic Scholar
  临床: ClinicalTrials.gov, ISRCTN, ChiCTR
  靶点: Open Targets, UniProt, ClinVar, ClinPGx, KEGG
  蛋白: PDB, AlphaFold DB, STRING DB
  表达: Human Protein Atlas, GTEx
  药物: Open Targets, ChEMBL, PubChem, KEGG
  专利: Google Patents, Patent9, drugfuture, USPTO, Lens.org, Espacenet, MCP
  AI  : DeepSeek / 小米 MiMo / 智谱 GLM / 自定义 OpenAI 兼容接口
"""

import asyncio
import io
import json
import os
import re
import sqlite3
import subprocess
import sys
import tempfile
import time
import xml.etree.ElementTree as ET
from datetime import datetime, timedelta
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from urllib.parse import urlencode, quote, unquote

import httpx
from html import escape
from dotenv import load_dotenv
from pydantic import BaseModel, Field

# PyQt5
from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QPushButton, QLineEdit, QLabel, QProgressBar, QGroupBox,
    QTextEdit, QTextBrowser, QScrollArea, QFrame, QMessageBox,
    QGridLayout, QFileDialog, QSizePolicy, QSpacerItem, QCheckBox,
    QComboBox, QTableWidget, QTableWidgetItem, QHeaderView,
    QAbstractItemView, QSplitter, QTabWidget,
)
from PyQt5.QtCore import Qt, QThread, pyqtSignal, QUrl, QTimer, QSettings
from PyQt5.QtGui import QFont, QDesktopServices

# PPT
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as PPT_RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Word
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

# Charts
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.font_manager import fontManager

# ─── Configuration ──────────────────────────────────────────────────────────

load_dotenv()

PUBMED_EMAIL = os.getenv('PUBMED_EMAIL', 'user@example.com')
NCBI_API_KEY = os.getenv('NCBI_API_KEY', '')  # optional, raises rate limit
OPENALEX_KEY = os.getenv('OPENALEX_KEY', '')

# API Base URLs
PUBMED_SEARCH_URL = 'https://eutils.ncbi.nlm.nih.gov/entrez/eutils/esearch.fcgi'
PUBMED_FETCH_URL = 'https://eutils.ncbi.nlm.nih.gov/entrez/eutils/efetch.fcgi'
CT_API_BASE = 'https://clinicaltrials.gov/api/v2'
OPEN_TARGETS_URL = 'https://api.platform.opentargets.org/api/v4/graphql'
UNIPROT_URL = 'https://rest.uniprot.org/uniprotkb/search'
RCSB_SEARCH_URL = 'https://search.rcsb.org/rcsbsearch/v2/query'
RCSB_GRAPHQL_URL = 'https://data.rcsb.org/graphql'
ALPHAFOLD_URL = 'https://alphafold.ebi.ac.uk/api'
HPA_URL = 'https://www.proteinatlas.org/api/search_download.php'
CHEMBL_SEARCH_URL = 'https://www.ebi.ac.uk/chembl/api/data/target/search'
CHEMBL_ACTIVITY_URL = 'https://www.ebi.ac.uk/chembl/api/data/activity'
STRING_API_URL = 'https://string-db.org/api'
PUBCHEM_URL = 'https://pubchem.ncbi.nlm.nih.gov/rest/pug'
KEGG_URL = 'https://rest.kegg.jp'
DGIDB_URL = 'https://dgidb.org/api/v2'
CLINPGX_URL = 'https://api.clinpgx.org/v1'  # PharmGKB 已于 2025 年更名为 ClinPGx，旧 api.pharmgkb.org 已关停
OPENALEX_URL = 'https://api.openalex.org'
SEMANTIC_SCHOLAR_URL = 'https://api.semanticscholar.org/graph/v1'
GTEX_URL = 'https://gtexportal.org/api/v2'

ISRCTN_URL = 'https://www.isrctn.com/search'
CHICTR_BASES = ('https://trialsearch.chictr.org.cn', 'https://www.chictr.org.cn')
PATENT9_URL = 'https://www.patent9.com'
DRUGFUTURE_URL = 'http://www2.drugfuture.com/cnpat'
CLINVAR_URL = 'https://eutils.ncbi.nlm.nih.gov/entrez/eutils/esearch.fcgi'
CLINVAR_FETCH_URL = 'https://eutils.ncbi.nlm.nih.gov/entrez/eutils/efetch.fcgi'

# ─── Retry helper ──────────────────────────────────────────────────────────


async def _request_with_retry(
    client: httpx.AsyncClient,
    method: str,
    url: str,
    max_retries: int = 3,
    base_delay: float = 1.5,
    **kwargs,
) -> httpx.Response:
    last_exc: Optional[Exception] = None
    for attempt in range(max_retries):
        try:
            if method == 'GET':
                return await client.get(url, **kwargs)
            elif method == 'POST':
                return await client.post(url, **kwargs)
            else:
                raise ValueError(f'Unsupported method: {method}')
        except (
            httpx.RemoteProtocolError,
            httpx.ConnectError,
            httpx.TimeoutException,
            httpx.ReadError,
            httpx.WriteError,
        ) as e:
            last_exc = e
            if attempt < max_retries - 1:
                await asyncio.sleep(base_delay * (2**attempt))
            continue
    raise last_exc or RuntimeError('Request failed after retries')


async def _request_json(
    client: httpx.AsyncClient,
    method: str,
    url: str,
    **kwargs,
) -> Any:
    resp = await _request_with_retry(client, method, url, **kwargs)
    resp.raise_for_status()
    return resp.json()


async def _request_text(
    client: httpx.AsyncClient,
    method: str,
    url: str,
    **kwargs,
) -> str:
    resp = await _request_with_retry(client, method, url, **kwargs)
    resp.raise_for_status()
    return resp.text


# ─── Concurrency limiter ────────────────────────────────────────────────────


class Limiter:
    def __init__(self, max_concurrent: int = 4):
        self._sem = asyncio.Semaphore(max_concurrent)

    async def run(self, coro):
        async with self._sem:
            return await coro


_limiter = Limiter(8)

# ─── Utility functions ─────────────────────────────────────────────────────


def _expand_query_terms(term: str) -> List[str]:
    seen = set()
    out = []
    def add(t):
        t = t.strip()
        if t and t not in seen:
            seen.add(t)
            out.append(t)
    add(term)
    add(term.upper())
    add(term.lower())
    add(term.capitalize())
    plain = re.sub(r'[\-–—.\' _]+', '', term)
    if plain != term:
        add(plain)
        add(plain.upper())
        add(plain.lower())
    h1 = re.sub(r'([a-zA-Z])(\d)', r'\1-\2', term)
    h2 = re.sub(r'(\d)([a-zA-Z])', r'\1-\2', term)
    if h1 != term:
        add(h1)
    if h2 != term:
        add(h2)
    for src, dst in [('2', 'II'), ('1', 'I')]:
        rv = term.replace(src, dst)
        if rv != term:
            add(rv)
            add(rv.upper())
            add(rv.lower())
    return out


def parse_target(name: str) -> Tuple[str, str, str]:
    raw = (name or '').strip()
    m = re.match(r'^([A-Za-z][A-Za-z0-9\-]*)[\s+\-._]+([A-Za-z]{1,3}\d{1,4}[A-Za-z]{0,2})$', raw)
    if m and re.match(r'^[A-Za-z]{1,3}\d', m.group(2)):
        return m.group(1), m.group(2), raw
    return raw, '', raw


def _clip(s: str, n: int = 2000) -> str:
    s = str(s or '')
    return s if len(s) <= n else s[:n] + '\n…（内容过长已截断）'


def slug(name: str) -> str:
    return re.sub(r'[\\/:*?"<>|]', '_', str(name or 'target'))


def _count_by_year(items: List, year_attr: str = 'year') -> Dict[str, int]:
    years = {}
    for item in items:
        y = getattr(item, year_attr, None) or ''
        if y and y.isdigit():
            years[y] = years.get(y, 0) + 1
    return dict(sorted(years.items()))


def _count_by_phase(trials) -> Dict[str, int]:
    phases = {}
    for t in trials:
        p = t.phase or '未明确'
        phases[p] = phases.get(p, 0) + 1
    return phases


def _collect_conditions(trials) -> str:
    conds = set()
    for t in trials:
        for c in t.conditions:
            conds.add(c)
    return '、'.join(list(conds)[:8]) if conds else '—'


def _fmt_date(d: dict) -> str:
    return '-'.join(str(d[k]) for k in ['year', 'month', 'day'] if d.get(k))


def _fmt_date_v2(d: dict) -> str:
    date_str = d.get('date', '')
    if date_str:
        return date_str[:10]
    return _fmt_date(d)


# ─── Data Models ────────────────────────────────────────────────────────────


class Paper(BaseModel):
    pmid: str = ''
    title: str = ''
    authors: List[str] = []
    journal: str = ''
    year: str = ''
    doi: str = ''
    abstract: str = ''
    keywords: List[str] = []
    mesh: List[str] = []
    pub_type: List[str] = []
    journal_country: str = ''
    source: str = 'PubMed'


class ClinicalTrial(BaseModel):
    nct_id: str = ''
    title: str = ''
    phase: str = ''
    status: str = ''
    conditions: List[str] = []
    interventions: List[str] = []
    sponsor: str = ''
    start_date: str = ''
    completion_date: str = ''
    brief_summary: str = ''
    enrollment: int = 0
    countries: List[str] = []
    source: str = 'ClinicalTrials.gov'


class DrugInfo(BaseModel):
    name: str = ''
    company: str = ''
    phase: str = ''
    mechanism_of_action: str = ''
    disease: str = ''


class PatentInfo(BaseModel):
    number: str = ''
    title: str = ''
    assignee: str = ''
    year: str = ''
    date: str = ''
    abstract: str = ''
    claims: int = 0
    citations: List[str] = []
    cpc: List[str] = []
    detail_legal: str = ''
    inventors: str = ''
    cited_by: int = 0
    countries: List[str] = []
    mcp_status: str = ''
    mcp_classification: str = ''
    mcp_abstract: str = ''
    link: str = ''
    snippet: str = ''
    priority_date: str = ''
    application_no: str = ''
    cn_download: str = ''
    cn_fulltext: str = ''


class MoleculeInfo(BaseModel):
    name: str = ''
    chembl: str = ''
    pchembl: str = ''
    best_type: str = ''
    best_val: str = ''
    details: List[Dict[str, str]] = []


class HPAData(BaseModel):
    gene: str = ''
    gene_syn: str = ''
    protein_class: str = ''
    molecular_func: str = ''
    evidence: str = ''
    rna_cancer: str = ''
    rna_cancer_score: str = ''
    protein_tissue: str = ''
    protein_tissue_score: str = ''
    subcell: str = ''
    sub_main: str = ''
    secretome: str = ''
    blood_c: str = ''
    blood_cl: str = ''
    antibody: str = ''
    antibody_reliab: str = ''
    antibody_n: int = 0
    link: str = ''


class PDBEntry(BaseModel):
    pdb_id: str = ''
    title: str = ''
    method: str = ''
    resolution: str = ''
    year: str = ''
    pubmed: str = ''
    doi: str = ''
    link: str = ''
    source: str = 'PDB'


class AlphaFoldEntry(BaseModel):
    uniprot_acc: str = ''
    pdb_url: str = ''
    cif_url: str = ''
    confidence: float = 0.0
    sequence: str = ''
    gene: str = ''
    source: str = 'AlphaFold'


class StringInteraction(BaseModel):
    protein_id: str = ''
    preferred_name: str = ''
    score: float = 0.0
    source: str = 'STRING'


class PubChemCompound(BaseModel):
    cid: int = 0
    name: str = ''
    formula: str = ''
    molecular_weight: float = 0.0
    logp: float = 0.0
    toxicity: str = ''
    link: str = ''
    source: str = 'PubChem'


class KEGGPathway(BaseModel):
    kegg_id: str = ''
    name: str = ''
    source: str = 'KEGG'


class DGIdbInteraction(BaseModel):
    drug_name: str = ''
    interaction_type: str = ''
    sources: List[str] = []
    source: str = 'DGIdb'


class ClinPGxAssociation(BaseModel):
    drug_name: str = ''
    gene_symbol: str = ''
    association_type: str = ''
    significance: str = ''
    level: str = ''
    disease: str = ''
    evidence: str = ''
    source: str = 'ClinPGx'


class ClinVarRecord(BaseModel):
    rcv_id: str = ''
    gene_symbol: str = ''
    clinical_significance: str = ''
    condition: str = ''
    review_status: str = ''
    source: str = 'ClinVar'


class GTExExpression(BaseModel):
    tissue: str = ''
    median_tpm: float = 0.0
    source: str = 'GTEx'


class TargetDetail(BaseModel):
    target_name: str = ''
    gene_symbol: str = ''
    uniprot_id: str = ''
    protein_name: str = ''
    protein_class: str = ''
    description: str = ''
    synonyms: List[str] = []
    subcellular: str = ''
    protein_families: str = ''
    sequence: str = ''
    sequence_length: int = 0
    tractability: List[str] = []
    related_diseases: List[str] = []
    hpa: Optional[HPAData] = None
    pdb: List[PDBEntry] = []
    alphafold: List[AlphaFoldEntry] = []
    string_interactions: List[StringInteraction] = []
    pubchem: List[PubChemCompound] = []
    kegg_pathways: List[KEGGPathway] = []
    kegg_disease: List[str] = []
    kegg_drugs: List[str] = []
    dgidb: List[DGIdbInteraction] = []
    clinpgx: List[ClinPGxAssociation] = []
    clinvar: List[ClinVarRecord] = []
    gtex: List[GTExExpression] = []
    mutation: str = ''
    hpa_link: str = ''


class ReportContent(BaseModel):
    target_name: str
    gene: str = ''
    mutation: str = ''
    target_detail: Optional[TargetDetail] = None
    papers: List[Paper] = []
    trials: List[ClinicalTrial] = []
    drugs: List[DrugInfo] = []
    patents: List[PatentInfo] = []
    molecules: List[MoleculeInfo] = []
    patent_insight: Optional[Dict] = None
    patent_note: str = ''
    cached_flags: List[str] = []
    increment_note: str = ''
    errors: List[str] = []
    target_overview: str = ''
    research_progress: str = ''
    clinical_landscape: str = ''
    key_findings: str = ''
    future_outlook: str = ''
    web_summary: str = ''
    patent_landscape: str = ''
    citations: Dict[str, List[str]] = {}
    target_overview_raw: str = ''
    research_progress_raw: str = ''
    clinical_landscape_raw: str = ''
    key_findings_raw: str = ''
    future_outlook_raw: str = ''

# ─── PubMed Literature Search ──────────────────────────────────────────────


async def search_papers(target: str, max_results: int = 40, years_back: int = 20) -> List[Paper]:
    now_year = datetime.now().year
    terms = _expand_query_terms(target)
    query_parts = [f'({t}[Title/Abstract])' for t in terms]
    query = ' OR '.join(query_parts)

    params = {
        'db': 'pubmed',
        'term': f'({query}) AND (review[pt] OR clinical trial[pt] OR systematic review[pt])',
        'retmax': max_results, 'retmode': 'json', 'sort': 'relevance',
        'email': PUBMED_EMAIL,
    }
    if years_back and years_back > 0:
        params['mindate'] = str(now_year - years_back)
        params['maxdate'] = str(now_year)
        params['datetype'] = 'pdat'
    if NCBI_API_KEY:
        params['api_key'] = NCBI_API_KEY

    async with httpx.AsyncClient(timeout=60) as client:
        try:
            resp = await _request_with_retry(client, 'GET', PUBMED_SEARCH_URL, params=params)
            resp.raise_for_status()
            id_list = resp.json().get('esearchresult', {}).get('idlist', [])
        except Exception:
            id_list = []

        if not id_list:
            params['term'] = f'({query})'
            try:
                resp = await _request_with_retry(client, 'GET', PUBMED_SEARCH_URL, params=params)
                resp.raise_for_status()
                id_list = resp.json().get('esearchresult', {}).get('idlist', [])
            except Exception:
                return []

        if not id_list:
            return []

        return await _fetch_papers(id_list)


async def _fetch_papers(pmids: List[str]) -> List[Paper]:
    params = {
        'db': 'pubmed', 'id': ','.join(pmids),
        'retmode': 'xml', 'rettype': 'abstract',
        'email': PUBMED_EMAIL,
    }
    if NCBI_API_KEY:
        params['api_key'] = NCBI_API_KEY
    async with httpx.AsyncClient(timeout=60) as client:
        try:
            resp = await _request_with_retry(client, 'GET', PUBMED_FETCH_URL, params=params, max_retries=4)
            resp.raise_for_status()
            return _parse_pubmed_xml(resp.text)
        except Exception:
            return []


def _parse_pubmed_xml(xml_str: str) -> List[Paper]:
    papers = []
    try:
        root = ET.fromstring(xml_str)
        for article_elem in root.findall('.//PubmedArticle'):
            try:
                paper = _parse_single_article(article_elem)
                if paper:
                    papers.append(paper)
            except Exception:
                continue
    except Exception:
        pass
    return papers


def _parse_single_article(article_elem) -> Optional[Paper]:
    medline = article_elem.find('.//MedlineCitation')
    article = medline.find('.//Article') if medline is not None else None
    if medline is None or article is None:
        return None

    title_elem = article.find('./ArticleTitle')
    title = ''.join(title_elem.itertext()) if title_elem is not None else ''

    journal_elem = article.find('./Journal/Title')
    journal = ''.join(journal_elem.itertext()) if journal_elem is not None else ''

    year_elem = article.find('./Journal/JournalIssue/PubDate/Year')
    year = ''.join(year_elem.itertext()) if year_elem is not None else ''

    pmid_elem = medline.find('./PMID')
    pmid = ''.join(pmid_elem.itertext()) if pmid_elem is not None else ''

    abstract = ''
    abstract_elem = article.find('./Abstract')
    if abstract_elem is not None:
        parts = []
        for elem in abstract_elem.iter():
            if elem.tag == 'AbstractText':
                label = elem.get('Label', '')
                text = ''.join(elem.itertext())
                parts.append(f'{label}: {text}' if label else text)
        abstract = ' '.join(parts)

    authors = []
    author_list = article.find('./AuthorList')
    if author_list is not None:
        for author in author_list:
            last = author.find('./LastName')
            fore = author.find('./ForeName')
            if last is not None:
                name = last.text or ''
                if fore is not None and fore.text:
                    name += ' ' + fore.text
                authors.append(name)

    keywords = []
    kw_list = medline.find('.//KeywordList')
    if kw_list is not None:
        for kw in kw_list:
            if kw.text:
                keywords.append(kw.text)

    mesh = []
    for mh in medline.findall('.//MeshHeading'):
        d = mh.find('./DescriptorName')
        d_name = d.text if d is not None else ''
        qs = [q.text for q in mh.findall('./QualifierName') if q.text]
        if d_name:
            sep = ' / '
            mesh.append(f'{d_name}{sep}{"/".join(qs)}' if qs else d_name)

    pub_type = []
    for pt in medline.findall('.//PublicationType'):
        if pt.text:
            pub_type.append(pt.text)

    jc = ''
    jc_elem = medline.find('.//MedlineJournalInfo/Country')
    if jc_elem is not None:
        jc = jc_elem.text or ''

    doi = ''
    for eid in article_elem.findall('.//ArticleIdList/ArticleId'):
        if eid.get('IdType') == 'doi':
            doi = eid.text or ''
            break

    return Paper(
        pmid=pmid, title=title, authors=authors[:5],
        journal=journal, year=year, doi=doi,
        abstract=abstract[:2000], keywords=keywords[:10],
        mesh=mesh[:8], pub_type=pub_type[:6],
        journal_country=jc,
    )


# ─── OpenAlex Literature Search ──────────────────────────────────────────


async def search_openalex(target: str, max_results: int = 10) -> List[Paper]:
    if not OPENALEX_KEY:
        return []
    terms = _expand_query_terms(target)
    query = ' OR '.join(terms)
    params = {
        'search': query,
        'per_page': min(max_results, 50),
        'sort': 'cited_by_count:desc',
        'filter': 'type:article|review',
    }
    headers = {}
    if OPENALEX_KEY:
        headers['api-key'] = OPENALEX_KEY
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            resp = await _request_with_retry(client, 'GET', f'{OPENALEX_URL}/works',
                                             params=params, headers=headers, max_retries=2)
            resp.raise_for_status()
            data = resp.json()
            results = data.get('results', [])
            papers = []
            for r in results[:max_results]:
                yr = r.get('publication_year', '')
                authors = [a.get('author', {}).get('display_name', '')
                          for a in (r.get('authorships') or [])[:5]]
                doi = (r.get('doi') or '').replace('https://doi.org/', '')
                abstract = r.get('abstract_inverted_index')
                abstract_text = ''
                if abstract:
                    word_positions = {}
                    for word, positions in abstract.items():
                        for pos in positions:
                            word_positions[pos] = word
                    abstract_text = ' '.join(word_positions[i] for i in sorted(word_positions))
                papers.append(Paper(
                    pmid='', title=r.get('title', '') or '',
                    authors=authors,
                    journal=(r.get('primary_location') or {}).get('source', {}).get('display_name', '') or '',
                    year=str(yr) if yr else '',
                    doi=doi,
                    abstract=abstract_text[:2000],
                    keywords=[], mesh=[], pub_type=[],
                    source='OpenAlex',
                ))
            return papers
        except Exception:
            return []


# ─── Semantic Scholar Literature Search ──────────────────────────────────


async def search_semantic_scholar(target: str, max_results: int = 10) -> List[Paper]:
    terms = _expand_query_terms(target)
    query = ' '.join(terms[:3])
    params = {
        'query': query,
        'limit': min(max_results, 50),
        'fields': 'title,authors,year,journal,externalIds,abstract,citationCount',
        'sort': 'citationCount:desc',
    }
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            resp = await _request_with_retry(client, 'GET',
                                             f'{SEMANTIC_SCHOLAR_URL}/paper/search',
                                             params=params, max_retries=2)
            resp.raise_for_status()
            data = resp.json()
            results = data.get('data', [])
            papers = []
            for r in results[:max_results]:
                title = r.get('title', '') or ''
                if not title:
                    continue
                authors = [a.get('name', '') for a in (r.get('authors') or [])[:5]]
                journal_info = r.get('journal') or {}
                journal = journal_info.get('name', '') or ''
                ext_ids = r.get('externalIds') or {}
                doi = ext_ids.get('DOI', '') or ''
                pmid = ext_ids.get('PubMed', '') or ''
                papers.append(Paper(
                    pmid=pmid, title=title, authors=authors,
                    journal=journal,
                    year=str(r.get('year', '')) if r.get('year') else '',
                    doi=doi,
                    abstract=(r.get('abstract') or '')[:2000],
                    keywords=[], mesh=[], pub_type=[],
                    source='SemanticScholar',
                ))
            return papers
        except Exception:
            return []


# ─── ClinicalTrials.gov API ────────────────────────────────────────────────


async def search_trials(target: str, max_results: int = 20) -> List[ClinicalTrial]:
    terms = _expand_query_terms(target)
    query = ' OR '.join(terms)
    params = {
        'query.term': query,
        'pageSize': min(max_results, 100), 'format': 'json',
        'fields': 'NCTId|BriefTitle|OverallStatus|Phase|Condition|'
                  'InterventionType|InterventionName|LeadSponsorName|StartDate|CompletionDate|'
                  'BriefSummary|EnrollmentCount|LocationCountry',
        'filter.overallStatus': 'ACTIVE_NOT_RECRUITING|COMPLETED|RECRUITING|NOT_YET_RECRUITING|ENROLLING_BY_INVITATION|AVAILABLE',
    }
    async with httpx.AsyncClient(timeout=60) as client:
        try:
            resp = await _request_with_retry(client, 'GET', f'{CT_API_BASE}/studies', params=params)
            if resp.status_code != 200:
                return []
            return [_parse_study(s) for s in resp.json().get('studies', [])]
        except Exception:
            return []


def _parse_study(study: dict) -> ClinicalTrial:
    p = study.get('protocolSection', {})
    idm = p.get('identificationModule', {})
    dm = p.get('designModule', {})
    smod = p.get('statusModule', {})
    cm = p.get('conditionsModule', {})
    aim = p.get('armsInterventionsModule', {})
    scm = p.get('sponsorCollaboratorsModule', {})
    desc = p.get('descriptionModule', {})
    clm = p.get('contactsLocationsModule', {})
    locations = clm.get('locations', []) if clm else []
    countries = list(set(l.get('country', '') for l in locations if l.get('country')))

    phase_raw = (dm.get('phases') or [''])[0]
    phase_map = {
        'EARLY_PHASE1': 'Phase 1(早期)', 'PHASE1': 'Phase 1',
        'PHASE2': 'Phase 2', 'PHASE3': 'Phase 3', 'PHASE4': 'Phase 4',
    }
    interventions = [
        f"[{a.get('interventionType', '')}] {a['interventionName']}"
        for a in aim.get('interventions', []) if a.get('interventionName')
    ]

    return ClinicalTrial(
        nct_id=idm.get('nctId', ''),
        title=idm.get('briefTitle', ''),
        status=smod.get('overallStatus', ''),
        start_date=_fmt_date_v2(smod.get('startDate', {}) or {}),
        completion_date=_fmt_date_v2(smod.get('completionDate', {}) or {}),
        phase=phase_map.get(phase_raw, phase_raw),
        conditions=cm.get('conditions', []),
        interventions=interventions,
        sponsor=(scm.get('leadSponsor') or {}).get('name', ''),
        brief_summary=(desc.get('briefSummary') or '')[:1000],
        enrollment=(dm.get('enrollmentInfo') or {}).get('count', 0),
        countries=countries[:10],
    )


# ─── ISRCTN ──────────────────────────────────────────────────────────────


def _parse_isrctn_page(html: str, max_results: int) -> List[ClinicalTrial]:
    """Parse ISRCTN search page. Result rows look like
    <a href="https://www.isrctn.com/ISRCTN12345678">Study title</a>."""
    trials = []
    seen = set()
    for m in re.finditer(r'<a[^>]+href="([^"]*ISRCTN\d{6,9}[^"]*)"[^>]*>(.*?)</a>',
                         html, re.I | re.S):
        id_match = re.search(r'(ISRCTN\d{6,9})', m.group(1))
        if not id_match:
            continue
        rid = id_match.group(1)
        if rid in seen:
            continue
        seen.add(rid)
        title = re.sub(r'<[^>]+>', '', m.group(2))
        title = re.sub(r'\s+', ' ', title).strip()
        trials.append(ClinicalTrial(nct_id=rid, title=title or '', source='ISRCTN'))
        if len(trials) >= max_results:
            break
    if not trials:
        for m in re.finditer(r'<h3[^>]*>([^<]*ISRCTN\d{6,9}[^<]*)</h3>', html, re.I):
            text = re.sub(r'\s+', ' ', m.group(1)).strip()
            id_match = re.search(r'(ISRCTN\d{6,9})', text)
            if id_match and id_match.group(1) not in seen:
                seen.add(id_match.group(1))
                trials.append(ClinicalTrial(
                    nct_id=id_match.group(1),
                    title=re.sub(r'\s*ISRCTN\d{6,9}\s*', '', text),
                    source='ISRCTN'))
                if len(trials) >= max_results:
                    break
    return trials


async def search_isrctn(target: str, max_results: int = 5) -> List[ClinicalTrial]:
    """
    ISRCTN retired the JSON API. The search page runs a Play framework
    cookie challenge: GET /search -> 303 -> /holding page whose JS sets a
    CHK cookie (urlencoded redirect) -> actual results page.
    We replay that flow manually and parse the HTML.
    """
    terms = _expand_query_terms(target)
    query = ' OR '.join(terms[:3])
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
        'Accept': 'text/html,application/xhtml+xml',
        'Accept-Language': 'en',
    }
    async with httpx.AsyncClient(timeout=30, follow_redirects=False, headers=headers) as client:
        try:
            resp = await client.get(ISRCTN_URL, params={'q': query, 'pageSize': min(max_results, 20)})
            if resp.status_code not in (200, 303):
                return []
            base = str(resp.url)
            for _ in range(5):
                if resp.status_code == 303 and resp.headers.get('location'):
                    loc = resp.headers['location']
                    base = str(httpx.URL(base).join(loc))
                    resp = await client.get(base)
                else:
                    break
            if 'holding' in str(resp.url).lower():
                chk = None
                try:
                    chk = client.cookies.get('CHK')
                except Exception:
                    chk = None
                if chk:
                    try:
                        nxt = unquote(chk)
                    except Exception:
                        nxt = chk
                    if nxt.startswith('/'):
                        resp = await client.get('https://www.isrctn.com' + nxt)
                    elif nxt.startswith('http'):
                        resp = await client.get(nxt)
            if resp.status_code != 200:
                return []
            return _parse_isrctn_page(resp.text, max_results)
        except Exception:
            return []


# ─── ChiCTR (中国临床试验注册中心, 新平台) ─────────────────────────────────


async def search_chictr(target: str, max_results: int = 5) -> List[ClinicalTrial]:
    """
    ChiCTR migrated to a SPA at trialsearch.chictr.org.cn (2025+).
    No stable public API is documented, so we probe several plausible
    REST endpoints defensively and degrade to [] on any failure.
    """
    terms = _expand_query_terms(target)
    query = ' '.join(terms[:2])
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
        'Accept': 'application/json,text/html',
        'Accept-Language': 'zh-CN,zh;q=0.9',
    }
    async with httpx.AsyncClient(timeout=20, follow_redirects=True, headers=headers) as client:
        data = None
        for base in ('https://trialsearch.chictr.org.cn', 'https://www.chictr.org.cn'):
            if data:
                break
            for path, params in (
                ('/api/search', {'keyword': query, 'page': 1, 'size': max_results}),
                ('/api/trials/search', {'keyword': query, 'pageNum': 1, 'pageSize': max_results}),
                ('/api/clinicalTrial/query', {'title': query, 'pageNum': 1, 'pageSize': max_results}),
            ):
                try:
                    resp = await client.get(base + path, params=params)
                    if resp.status_code == 200 and 'json' in (resp.headers.get('content-type') or ''):
                        try:
                            j = resp.json()
                        except Exception:
                            continue
                        if j is None:
                            continue
                        data = j.get('data') if isinstance(j, dict) else j
                        if isinstance(data, dict) and not data:
                            data = None
                        if data is not None:
                            break
                except Exception:
                    continue
        if not data:
            return []
        if isinstance(data, dict):
            rows = (data.get('list') or data.get('records') or data.get('rows')
                    or data.get('items') or data.get('content') or [])
        else:
            rows = data if isinstance(data, list) else []
        if not rows:
            return []
        trials = []
        for r in rows[:max_results]:
            if not isinstance(r, dict):
                continue
            tid = str(r.get('regNo') or r.get('registrationNumber') or r.get('chiCTR')
                      or r.get('registrationNo') or r.get('id') or '')
            title = str(r.get('title') or r.get('publicTitle') or r.get('trialTitle') or '')
            cond = r.get('condition') or r.get('disease') or ''
            trials.append(ClinicalTrial(
                nct_id=tid, title=title,
                status=str(r.get('status') or ''),
                conditions=[cond] if cond else [],
                sponsor=str(r.get('sponsor') or r.get('applicant') or ''),
                source='ChiCTR',
            ))
        return [t for t in trials if t.nct_id or t.title]


# ─── Open Targets GraphQL ──────────────────────────────────────────────────


async def _get_ensembl_id(target_name: str) -> Optional[str]:
    q = 'query($q:String!){search(queryString:$q,entityNames:["target"]){hits{id}}}'
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            resp = await _request_with_retry(client, 'POST', OPEN_TARGETS_URL,
                                             json={'query': q, 'variables': {'q': target_name}})
            hits = resp.json().get('data', {}).get('search', {}).get('hits', [])
            return hits[0]['id'] if hits else None
        except Exception:
            return None


async def get_target_detail(target_name: str) -> TargetDetail:
    detail = TargetDetail(target_name=target_name)
    eid = await _get_ensembl_id(target_name)
    if not eid:
        return detail

    q = """query($id:String!){
      target(ensemblId:$id){
        approvedSymbol functionDescriptions
        nameSynonyms{label source}
        targetClass{label}
        tractability{label value}
        associatedDiseases{rows{score disease{name}}}
      }
    }"""
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            resp = await _request_with_retry(client, 'POST', OPEN_TARGETS_URL,
                                             json={'query': q, 'variables': {'id': eid}})
            data = resp.json().get('data', {}).get('target', {})
            if not data:
                return detail
            detail.gene_symbol = data.get('approvedSymbol', '')
            descs = data.get('functionDescriptions') or []
            detail.description = descs[0] if descs else ''
            detail.synonyms = [s['label'] for s in (data.get('nameSynonyms') or [])]
            cls_list = data.get('targetClass') or []
            detail.protein_class = cls_list[0].get('label', '') if cls_list else ''
            detail.tractability = [t.get('label', '') for t in (data.get('tractability') or []) if t.get('value')]
            assoc = data.get('associatedDiseases') or {}
            rows = assoc.get('rows') or []
            rows_sorted = sorted(rows, key=lambda r: r.get('score', 0) or 0, reverse=True)
            detail.related_diseases = [
                r['disease']['name'] for r in rows_sorted[:10]
                if r.get('disease', {}).get('name')
            ]
        except Exception:
            pass
    return detail


async def get_drugs(target_name: str) -> List[DrugInfo]:
    eid = await _get_ensembl_id(target_name)
    if not eid:
        return []
    q = """query($id:String!){
      target(ensemblId:$id){
        drugAndClinicalCandidates{
          rows{
            maxClinicalStage
            drug{
              name
              mechanismsOfAction{rows{mechanismOfAction}}
            }
            diseases{diseaseFromSource disease{name}}
          }
        }
      }
    }"""
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            resp = await _request_with_retry(client, 'POST', OPEN_TARGETS_URL,
                                             json={'query': q, 'variables': {'id': eid}})
            rows = (resp.json().get('data', {}).get('target', {})
                    .get('drugAndClinicalCandidates', {})).get('rows', [])
            drugs = []
            for r in rows:
                drug_data = r.get('drug', {})
                moas = drug_data.get('mechanismsOfAction', {}).get('rows') or []
                diseases = r.get('diseases') or []
                d0 = diseases[0] if diseases else {}
                drugs.append(DrugInfo(
                    name=drug_data.get('name', ''),
                    mechanism_of_action=moas[0].get('mechanismOfAction', '') if moas else '',
                    phase=r.get('maxClinicalStage', ''),
                    disease=d0.get('disease', {}).get('name', '') or d0.get('diseaseFromSource', ''),
                ))
            return drugs
        except Exception:
            return []


async def _enrich_drugs_with_company(drugs: List[DrugInfo]) -> None:
    names = [d.name for d in drugs if d.name]
    if not names:
        return
    query = ' OR '.join(f'"{n}"' for n in names)
    params = {'query.term': query, 'pageSize': min(len(names) * 5, 100),
              'format': 'json',
              'fields': 'NCTId|InterventionType|InterventionName|LeadSponsorName'}
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            resp = await _request_with_retry(client, 'GET', f'{CT_API_BASE}/studies', params=params)
            if resp.status_code != 200:
                return
            mapping = {}
            for s in resp.json().get('studies', []):
                p = s.get('protocolSection', {})
                sponsor = (p.get('sponsorCollaboratorsModule', {})
                           .get('leadSponsor') or {}).get('name', '')
                if not sponsor:
                    continue
                for inv in (p.get('armsInterventionsModule', {})
                            .get('interventions') or []):
                    inv_name = inv.get('interventionName', '') or inv.get('name', '')
                    if inv_name:
                        mapping.setdefault(inv_name.lower(), sponsor)
            for drug in drugs:
                key = drug.name.lower()
                if key in mapping:
                    drug.company = mapping[key]
                else:
                    for inv_name, sponsor in mapping.items():
                        if key in inv_name or inv_name in key:
                            drug.company = sponsor
                            break
        except Exception:
            pass


# ─── UniProt ────────────────────────────────────────────────────────────────


async def _uniprot_search(query_str: str, size: int = 5) -> List[dict]:
    params = {
        'query': query_str, 'format': 'json', 'size': str(size),
        'fields': 'accession,protein_name,cc_subcellular_location,cc_function,protein_families,sequence',
    }
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            resp = await _request_with_retry(client, 'GET', UNIPROT_URL, params=params, max_retries=2)
            resp.raise_for_status()
            return resp.json().get('results', [])
        except Exception:
            return []


def _uniprot_pick(results: List[dict], gene: str) -> Optional[dict]:
    human = [r for r in results if r.get('organism', {}).get('scientificName') == 'Homo sapiens']
    pool = human or results
    reviewed = [r for r in pool if 'reviewed' in str(r.get('entryType', '')).lower()]
    return reviewed[0] if reviewed else (pool[0] if pool else None)


async def get_uniprot(gene: str, aliases: List[str] = None) -> Optional[dict]:
    if not gene:
        return None
    results = await _uniprot_search(f'gene:{gene} AND organism_id:9606', 5)
    if not results and aliases:
        for alias in aliases[:3]:
            if not alias:
                continue
            results = await _uniprot_search(f'gene_name:{alias} AND organism_id:9606', 5)
            if results:
                break
    if not results:
        return None
    r0 = _uniprot_pick(results, gene)
    if not r0:
        return None
    comments = r0.get('comments') or []

    def comment_val(ctype: str, key: str) -> str:
        items = []
        for c in comments:
            if c.get('commentType') == ctype:
                for item in c.get(key) or []:
                    if isinstance(item, dict):
                        if 'location' in item:
                            items.append(item['location'].get('value', ''))
                        else:
                            items.append(str(item.get('value', item.get('text', ''))))
                    elif isinstance(item, str):
                        items.append(item)
                break
        return '; '.join(filter(None, items))

    fam_list = r0.get('proteinFamilies') or []
    families = '; '.join(f.get('value', '') for f in fam_list if f.get('value'))
    seq_data = r0.get('sequence') or {}
    return {
        'accession': r0.get('primaryAccession', ''),
        'proteinName': (r0.get('proteinDescription', {}).get('recommendedName', {}).get('fullName', {}).get('value', '')
                        or (r0.get('proteinDescription', {}).get('submissionNames') or [{}])[0].get('fullName', {}).get('value', '')
                        or ''),
        'subcellular': comment_val('SUBCELLULAR LOCATION', 'subcellularLocations'),
        'func': comment_val('FUNCTION', 'text'),
        'proteinFamilies': families,
        'sequence': seq_data.get('value', ''),
        'sequenceLength': seq_data.get('length', 0),
    }


# ─── PDB ────────────────────────────────────────────────────────────────────


async def get_pdb(gene: str, uniprot_acc: str = '') -> List[PDBEntry]:
    if not gene and not uniprot_acc:
        return []
    try:
        nodes = []
        if uniprot_acc:
            nodes.append({'type': 'terminal', 'service': 'text', 'parameters': {
                'attribute': 'rcsb_polymer_entity_container_identifiers.reference_sequence_identifiers.database_accession',
                'operator': 'exact_match', 'value': uniprot_acc}})
            nodes.append({'type': 'terminal', 'service': 'text', 'parameters': {
                'attribute': 'rcsb_polymer_entity_container_identifiers.reference_sequence_identifiers.database_name',
                'operator': 'exact_match', 'value': 'UniProt'}})
        elif gene:
            nodes.append({'type': 'terminal', 'service': 'text', 'parameters': {
                'attribute': 'rcsb_entity_source_organism.rcsb_gene_name.value',
                'operator': 'exact_match', 'value': gene, 'case_sensitive': False}})
        if not nodes:
            return []
        body = {
            'query': nodes[0] if len(nodes) == 1 else {'type': 'group', 'logical_operator': 'and', 'nodes': nodes},
            'return_type': 'entry',
            'request_options': {
                'results_content_type': ['experimental'],
                'sort': [{'sort_by': 'score', 'direction': 'desc'}],
                'paginate': {'start': 0, 'rows': 10},
            },
        }
        async with httpx.AsyncClient(timeout=30) as client:
            search_data = await _request_json(client, 'POST', RCSB_SEARCH_URL, json=body, max_retries=2)
            pdb_ids = [r['identifier'] for r in search_data.get('result_set', []) if r.get('identifier')]
            if not pdb_ids:
                return []
            gql = """query($ids:[String!]!){
              entries(entry_ids:$ids){
                rcsb_id struct{title} exptl{method} refine{ls_d_res_high}
                rcsb_primary_citation{year pdbx_database_id_PubMed pdbx_database_id_DOI}
              }
            }"""
            detail_data = await _request_json(client, 'POST', RCSB_GRAPHQL_URL, json={
                'query': gql, 'variables': {'ids': pdb_ids}}, max_retries=2)
            entries = detail_data.get('data', {}).get('entries', [])
            return [PDBEntry(
                pdb_id=e.get('rcsb_id', ''),
                title=(e.get('struct') or {}).get('title', '') or '',
                method=(e.get('exptl') or [{}])[0].get('method', '') if e.get('exptl') else '',
                resolution=str((e.get('refine') or [{}])[0].get('ls_d_res_high', '') or ''),
                year=str((e.get('rcsb_primary_citation') or {}).get('year', '') or ''),
                pubmed=str((e.get('rcsb_primary_citation') or {}).get('pdbx_database_id_PubMed', '') or ''),
                doi=str((e.get('rcsb_primary_citation') or {}).get('pdbx_database_id_DOI', '') or ''),
                link=f"https://www.rcsb.org/structure/{e.get('rcsb_id', '')}" if e.get('rcsb_id') else '',
            ) for e in entries]
    except Exception:
        return []


# ─── AlphaFold ──────────────────────────────────────────────────────────────


async def get_alphafold(gene: str, uniprot_acc: str = '') -> List[AlphaFoldEntry]:
    if not uniprot_acc and not gene:
        return []
    try:
        async with httpx.AsyncClient(timeout=30) as client:
            url = f'{ALPHAFOLD_URL}/uniprot/{uniprot_acc}' if uniprot_acc else f'{ALPHAFOLD_URL}/search?gene={gene}'
            resp = await _request_with_retry(client, 'GET', url, max_retries=2)
            if resp.status_code != 200:
                return []
            data = resp.json()
            entries = data if isinstance(data, list) else [data]
            return [AlphaFoldEntry(
                uniprot_acc=e.get('uniprotAccession', '') or e.get('uniprot_id', '') or '',
                pdb_url=e.get('pdbUrl', '') or '',
                cif_url=e.get('cifUrl', '') or '',
                confidence=float(e.get('confidence', 0) or 0),
                sequence=e.get('sequence', '') or '',
                gene=e.get('gene', '') or '',
            ) for e in entries[:5] if e.get('uniprotAccession') or e.get('pdbUrl')]
    except Exception:
        return []


# ─── Human Protein Atlas ────────────────────────────────────────────────────


async def get_hpa(gene: str) -> Optional[HPAData]:
    if not gene:
        return None
    cols = 'g,gs,pc,up_mf,pe,rnacas,rnacss,prts,prtss,ab,relih,scl,scml,scal,secl,blconcms,ecblood'
    params = {'search': gene, 'format': 'json', 'columns': cols, 'compress': 'no'}
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            resp = await _request_with_retry(client, 'GET', HPA_URL, params=params, max_retries=2)
            resp.raise_for_status()
            data = resp.json()
            arr = data if isinstance(data, list) else []
            if not arr:
                return None
            gene_lower = gene.lower()
            match = next((x for x in arr if x and str(x.get('Gene', '') or '').lower() == gene_lower), None)
            if not match:
                match = next((x for x in arr if x and (x.get('Gene') or x.get('Protein class'))), None)
            if not match or (not match.get('Gene') and not match.get('Protein class')):
                return None

            def L(k):
                v = match.get(k)
                if isinstance(v, list):
                    return '; '.join(filter(None, v))
                return str(v or '')

            antibody_str = L('Antibody')
            antibody_count = len([s for s in antibody_str.split(';') if s.strip()])
            return HPAData(
                gene=L('Gene'), gene_syn=L('Gene synonym'),
                protein_class=L('Protein class'), molecular_func=L('Molecular function'),
                evidence=L('Evidence'),
                rna_cancer=L('RNA cancer specificity'), rna_cancer_score=L('RNA cancer specificity score'),
                protein_tissue=L('Protein tissue specificity'), protein_tissue_score=L('Protein tissue specificity score'),
                subcell=L('Subcellular location'), sub_main=L('Subcellular main location'),
                secretome=L('Secretome location'),
                blood_c=L('Blood concentration - Conc. blood MS [pg/L]'),
                blood_cl=L('Blood expression cluster'),
                antibody=antibody_str, antibody_reliab=L('Reliability (IH)'),
                antibody_n=antibody_count,
                link=f'https://www.proteinatlas.org/search/{quote(gene)}',
            )
        except Exception:
            return None


# ─── GTEx ────────────────────────────────────────────────────────────────────


async def get_gtex(gene: str) -> List[GTExExpression]:
    if not gene:
        return []
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            resp = await _request_with_retry(client, 'GET',
                                             f'{GTEX_URL}/expression/geneExpression',
                                             params={'geneId': gene, 'format': 'json'},
                                             max_retries=2)
            if resp.status_code != 200:
                return []
            data = resp.json()
            results = []
            tissue_data = (data.get('data', {}).get('tissueExpression', {})
                          or data.get('tissueExpression', []) or [])
            if isinstance(tissue_data, list):
                for item in tissue_data[:20]:
                    results.append(GTExExpression(
                        tissue=item.get('tissueSiteDetail', '') or item.get('tissue', '') or '',
                        median_tpm=float(item.get('medianExpression', 0) or item.get('tpm', 0) or 0),
                    ))
            elif isinstance(tissue_data, dict):
                for tissue, tpm in list(tissue_data.items())[:20]:
                    results.append(GTExExpression(
                        tissue=tissue,
                        median_tpm=float(tpm or 0),
                    ))
            return results
        except Exception:
            return []


# ─── ChEMBL ──────────────────────────────────────────────────────────────────


async def get_chembl(name: str) -> List[MoleculeInfo]:
    try:
        params = {'q': name, 'only_documented': 'true', 'limit': '5'}
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await _request_with_retry(client, 'GET', CHEMBL_SEARCH_URL, params=params,
                                             headers={'Accept': 'application/json'}, max_retries=2)
            resp.raise_for_status()
            data = resp.json()
            targets = data.get('targets', [])
            t = next((x for x in targets if x.get('organism') == 'Homo sapiens'), None) or (targets[0] if targets else None)
            if not t or not t.get('target_chembl_id'):
                return []
            params2 = {'target_chembl_id': t['target_chembl_id'], 'limit': '100', 'offset': '0'}
            resp2 = await _request_with_retry(client, 'GET', CHEMBL_ACTIVITY_URL, params=params2,
                                              headers={'Accept': 'application/json'}, max_retries=2)
            resp2.raise_for_status()
            activities = resp2.json().get('activities', [])
            mol_map = {}
            type_order = ['IC50', 'EC50', 'Ki', 'Kd', 'CC50', 'GI50', 'ED50']
            for a in activities:
                mid = a.get('molecule_chembl_id')
                if not mid:
                    continue
                pc = a.get('pchembl_value')
                pc_val = float(pc) if pc else None
                stype = str(a.get('standard_type', '') or '').upper()
                sval = a.get('standard_value')
                sval_num = float(sval) if sval else None
                rec = {'stype': stype, 'sval': sval_num, 'units': a.get('standard_units', '') or '',
                       'pchembl': pc_val}
                if mid not in mol_map:
                    mol_map[mid] = {'potency': pc_val, 'chembl': mid,
                                    'name': a.get('molecule_pref_name', '') or '',
                                    'details': [rec] if stype and sval_num else []}
                else:
                    prev = mol_map[mid]
                    if pc_val is not None and (prev['potency'] is None or pc_val > prev['potency']):
                        prev['potency'] = pc_val
                        prev['name'] = a.get('molecule_pref_name', '') or prev['name']
                    if stype and sval_num:
                        existing = next((d for d in prev['details'] if d['stype'] == stype), None)
                        if not existing:
                            prev['details'].append(rec)
                        elif pc_val and pc_val > (existing.get('pchembl') or 0):
                            prev['details'][prev['details'].index(existing)] = rec
            out = []
            for mol in sorted(mol_map.values(), key=lambda x: -(x['potency'] or -1)):
                mol['details'].sort(key=lambda d: type_order.index(d['stype']) if d['stype'] in type_order else 99)
                best = mol['details'][0] if mol['details'] else {}
                best_val_str = f"{best['sval']} {best['units']}" if best.get('sval') is not None and best.get('units') else (str(best['sval']) if best.get('sval') is not None else '')
                out.append(MoleculeInfo(
                    name=mol['name'] or mol['chembl'] or '',
                    chembl=mol['chembl'],
                    pchembl=f"{mol['potency']:.2f}" if mol['potency'] is not None else '',
                    best_type=best.get('stype', '') or '',
                    best_val=best_val_str,
                    details=[{'type': d['stype'],
                              'val': str(d['sval'] or '') + (' ' + d['units'] if d.get('units') else ''),
                              'pchembl': f"{d['pchembl']:.2f}" if d.get('pchembl') is not None else ''}
                             for d in mol['details'][:8]],
                ))
            return out[:8]
    except Exception:
        return []


# ─── STRING DB (Protein-Protein Interactions) ─────────────────────────────


async def get_string_interactions(gene: str) -> List[StringInteraction]:
    if not gene:
        return []
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            params = {
                'identifiers': gene,
                'species': 9606,
                'limit': 20,
                'caller_identity': PUBMED_EMAIL,
            }
            resp = await _request_with_retry(client, 'GET',
                                             f'{STRING_API_URL}/json/interaction_partners',
                                             params=params, max_retries=2)
            if resp.status_code != 200:
                return []
            data = resp.json()
            if not isinstance(data, list):
                return []
            return [StringInteraction(
                protein_id=item.get('stringId', '') or item.get('proteinId', '') or '',
                preferred_name=item.get('preferredName', '') or '',
                score=float(item.get('score', 0) or 0),
            ) for item in data[:20] if item.get('preferredName')]
        except Exception:
            return []


# ─── PubChem ────────────────────────────────────────────────────────────────


async def get_pubchem(gene: str) -> List[PubChemCompound]:
    if not gene:
        return []
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            url = f'{PUBCHEM_URL}/assay/type/JSON'
            params = {'list_return': 'True', 'gene': gene, 'maxRecords': 10}
            resp = await _request_with_retry(client, 'GET', url, params=params, max_retries=2)
            if resp.status_code != 200:
                return []
            data = resp.json()
            results = []
            assays = data.get('PC_AssaySubmit', []) or data.get('assay', []) or data.get('results', []) or []
            for item in assays[:10]:
                aid = item.get('aid', 0) or item.get('id', 0) or 0
                name = item.get('name', '') or item.get('assay_name', '') or ''
                if aid:
                    results.append(PubChemCompound(
                        cid=aid,
                        name=name,
                        link=f'https://pubchem.ncbi.nlm.nih.gov/assay/{aid}',
                    ))
            return results
        except Exception:
            return []


async def get_pubchem_compound_by_name(name: str) -> Optional[PubChemCompound]:
    if not name:
        return None
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            url = f'{PUBCHEM_URL}/compound/name/{quote(name)}/JSON'
            resp = await _request_with_retry(client, 'GET', url, max_retries=2)
            if resp.status_code != 200:
                return None
            data = resp.json()
            props = data.get('PC_Compounds', [{}])[0].get('props', []) if data.get('PC_Compounds') else []
            cid = (data.get('PC_Compounds') or [{}])[0].get('id', {}).get('id', {}).get('cid', 0)
            if not cid:
                return None
            name_val = ''
            formula = ''
            mw = 0.0
            logp_val = 0.0
            for p in props:
                urn = p.get('urn', {}).get('label', '')
                val = p.get('value', {}).get('sval', '') or p.get('value', {}).get('fval', 0)
                if urn == 'IUPAC Name':
                    name_val = str(val)
                elif urn == 'Molecular Formula':
                    formula = str(val)
                elif urn == 'Molecular Weight':
                    mw = float(val or 0)
                elif urn == 'Log P':
                    logp_val = float(val or 0)
            return PubChemCompound(
                cid=cid,
                name=name_val or name,
                formula=formula,
                molecular_weight=mw,
                logp=logp_val,
                link=f'https://pubchem.ncbi.nlm.nih.gov/compound/{cid}',
            )
        except Exception:
            return None


# ─── KEGG ────────────────────────────────────────────────────────────────────


async def get_kegg_pathways(gene: str) -> List[KEGGPathway]:
    if not gene:
        return []
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            kegg_id = gene
            if not gene.startswith('hsa:'):
                find_resp = await _request_with_retry(client, 'GET',
                    f'{KEGG_URL}/find/genes/{gene}', max_retries=2)
                if find_resp.status_code == 200:
                    for line in find_resp.text.strip().split('\n'):
                        if '\t' in line:
                            kegg_id = line.split('\t')[0].strip()
                            break
            url = f'{KEGG_URL}/link/pathway/{kegg_id}'
            resp = await _request_with_retry(client, 'GET', url, max_retries=2)
            if resp.status_code != 200:
                return []
            lines = resp.text.strip().split('\n')
            pathways = []
            for line in lines[:15]:
                parts = line.strip().split('\t')
                if len(parts) >= 2:
                    pid = parts[0].strip()
                    pname = parts[1].strip()
                    pathways.append(KEGGPathway(kegg_id=pid, name=pname))
            return pathways
        except Exception:
            return []


async def _kegg_gene_id(client: httpx.AsyncClient, gene: str) -> str:
    """Resolve a gene symbol to its KEGG entry id (e.g. hsa:1956)."""
    if not gene:
        return ''
    if gene.startswith('hsa:'):
        return gene
    try:
        find_resp = await _request_with_retry(
            client, 'GET', f'{KEGG_URL}/find/genes/{gene}', max_retries=2)
        if find_resp.status_code == 200:
            for line in find_resp.text.strip().split('\n'):
                if '\t' in line:
                    return line.split('\t')[0].strip()
    except Exception:
        pass
    return ''


async def get_kegg_disease(gene: str) -> List[str]:
    if not gene:
        return []
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            kegg_id = await _kegg_gene_id(client, gene)
            if not kegg_id:
                return []
            url = f'{KEGG_URL}/link/disease/{kegg_id}'
            resp = await _request_with_retry(client, 'GET', url, max_retries=2)
            if resp.status_code != 200:
                return []
            lines = resp.text.strip().split('\n')
            diseases = []
            for line in lines[:10]:
                parts = line.strip().split('\t')
                if len(parts) >= 2:
                    diseases.append(parts[1].strip())
            return diseases
        except Exception:
            return []


async def get_kegg_drugs(gene: str) -> List[str]:
    if not gene:
        return []
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            kegg_id = await _kegg_gene_id(client, gene)
            if not kegg_id:
                return []
            url = f'{KEGG_URL}/link/drug/{kegg_id}'
            resp = await _request_with_retry(client, 'GET', url, max_retries=2)
            if resp.status_code != 200:
                return []
            lines = resp.text.strip().split('\n')
            drugs = []
            for line in lines[:10]:
                parts = line.strip().split('\t')
                if len(parts) >= 2:
                    drugs.append(parts[1].strip())
            return drugs
        except Exception:
            return []


# ─── DGIdb ──────────────────────────────────────────────────────────────────


async def get_dgidb(gene: str) -> List[DGIdbInteraction]:
    if not gene:
        return []
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            resp = await _request_with_retry(client, 'GET',
                                             f'{DGIDB_URL}/interactions',
                                             params={'genes': gene},
                                             max_retries=2)
            if resp.status_code != 200:
                return []
            data = resp.json()
            interactions = []
            for item in (data.get('matchedTerms') or []):
                for interaction in (item.get('interactions') or [])[:5]:
                    interactions.append(DGIdbInteraction(
                        drug_name=interaction.get('drugName', '') or '',
                        interaction_type=interaction.get('interactionType', '') or '',
                        sources=interaction.get('sources', []) or [],
                    ))
            return interactions[:15]
        except Exception:
            return []


# ─── ClinPGx (原 PharmGKB，2025-07 更名；免 Key) ──────────────────────────


async def get_clinpgx(gene: str) -> List[ClinPGxAssociation]:
    """
    ClinPGx REST API v1 (api.clinpgx.org, 免 Key, 限速 ~2 req/s):
    GET /data/summaryAnnotation?location.genes.symbol=<gene>&view=max
    Old api.pharmgkb.org was shut down 2026-07-20.
    """
    if not gene:
        return []
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            resp = await _request_with_retry(
                client, 'GET', f'{CLINPGX_URL}/data/summaryAnnotation',
                params={'location.genes.symbol': gene, 'view': 'max'}, max_retries=2)
            if resp.status_code != 200:
                return []
            data = (resp.json() or {}).get('data', []) or []
            if not isinstance(data, list):
                return []
            results = []
            for item in data[:10]:
                if not isinstance(item, dict):
                    continue
                chems = item.get('relatedChemicals') or []
                drug = ''
                if isinstance(chems, list) and chems:
                    c0 = chems[0]
                    drug = str(c0.get('name', '') or '') if isinstance(c0, dict) else str(c0 or '')
                level = item.get('levelOfEvidence') or {}
                level_txt = ''
                if isinstance(level, dict):
                    level_txt = str(level.get('description') or '')
                name_txt = str(item.get('name') or '')
                m = re.search(r'(level\s+\d+)', name_txt, re.I)
                if m:
                    level_txt = m.group(1).title()
                elif level_txt:
                    level_txt = level_txt[:120]
                types = item.get('types') or []
                if not isinstance(types, list):
                    types = [types]
                type_txt = ', '.join(str(t) for t in types if t)
                diseases = item.get('relatedDiseases') or []
                disease = ''
                if isinstance(diseases, list) and diseases:
                    d0 = diseases[0]
                    disease = str(d0.get('name', '') or '') if isinstance(d0, dict) else str(d0 or '')
                results.append(ClinPGxAssociation(
                    drug_name=drug,
                    gene_symbol=gene,
                    association_type=type_txt,
                    significance=name_txt,
                    level=level_txt,
                    disease=disease,
                    evidence=name_txt,
                ))
            return results
        except Exception:
            return []


# ─── ClinVar ────────────────────────────────────────────────────────────────


async def get_clinvar(gene: str) -> List[ClinVarRecord]:
    if not gene:
        return []
    params = {
        'db': 'clinvar',
        'term': f'{gene}[gene]',
        'retmax': 10,
        'retmode': 'json',
    }
    if NCBI_API_KEY:
        params['api_key'] = NCBI_API_KEY
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            resp = await _request_with_retry(client, 'GET', CLINVAR_URL, params=params, max_retries=2)
            resp.raise_for_status()
            id_list = resp.json().get('esearchresult', {}).get('idlist', [])
            if not id_list:
                return []
            fetch_params = {
                'db': 'clinvar',
                'id': ','.join(id_list[:10]),
                'retmode': 'xml',
            }
            if NCBI_API_KEY:
                fetch_params['api_key'] = NCBI_API_KEY
            resp2 = await _request_with_retry(client, 'GET', CLINVAR_FETCH_URL, params=fetch_params, max_retries=2)
            resp2.raise_for_status()
            xml_str = resp2.text
            records = []
            try:
                root = ET.fromstring(xml_str)
                for var in root.findall('.//ClinVarSet'):
                    try:
                        rcv = var.find('.//ReferenceClinVarAssertion')
                        if rcv is None:
                            continue
                        rcv_id = rcv.get('ID', '') or ''
                        gene_symbol = (rcv.find('.//Gene/Symbol') or
                                       rcv.find('.//GeneSymbol'))
                        gene_name = gene_symbol.text if gene_symbol is not None else ''
                        cls_name = rcv.find('.//Description')
                        clin_sig = rcv.find('.//ClinicalSignificance/Description')
                        review = rcv.find('.//ClinicalSignificance/ReviewStatus')
                        sig = clin_sig.text if clin_sig is not None else ''
                        review_status = review.text if review is not None else ''
                        condition = cls_name.text if cls_name is not None else ''
                        if rcv_id:
                            records.append(ClinVarRecord(
                                rcv_id=rcv_id,
                                gene_symbol=gene_name,
                                clinical_significance=sig,
                                condition=condition,
                                review_status=review_status,
                            ))
                    except Exception:
                        continue
            except Exception:
                pass
            return records[:10]
        except Exception:
            return []

# ─── Patent Search ──────────────────────────────────────────────────────────

PATENT_MCP_FALLBACK = 'https://gateway.pipeworx.io/patents/mcp'

def _dedupe_patents(lst: List[dict]) -> List[dict]:
    seen = set()
    out = []
    for p in (lst or []):
        k = str(p.get('number', '') or p.get('title', '') or '').strip()
        if k and k not in seen:
            seen.add(k)
            out.append(p)
    return out

def _google_num(num: str) -> str:
    s = re.sub(r'^patent/', '', str(num or ''))
    s = re.sub(r'/.*$', '', s)
    s = re.sub(r'^US0+', '', s)
    return s.strip()

def _strip_tags(s: str) -> str:
    return re.sub(r'\s+', ' ', re.sub(r'<[^>]*>', ' ', str(s or ''))).strip()

async def uspto_search(target: str, key: str) -> List[dict]:
    if not key:
        return []
    terms = _expand_query_terms(target)[:3]
    safe = lambda t: t.replace('"', ' ').replace('\\', ' ')
    q_parts = [f'inventionTitle:({safe(t)})' for t in terms] + [f'({safe(t)})' for t in terms]
    q = ' OR '.join(q_parts) if q_parts else f'({safe(target)})'
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            resp = await _request_with_retry(client, 'GET',
                'https://api.uspto.gov/api/v1/patent/applications/search',
                params={'q': q, 'limit': '50', 'offset': '0'},
                headers={'X-API-Key': key}, max_retries=2)
            if resp.status_code != 200:
                return []
            data = resp.json()
            arr = (data if isinstance(data, list)
                   else data.get('patentBag') or data.get('data', {}).get('patentBag') or data.get('data') or [])
            if not isinstance(arr, list):
                return []
            out = []
            for x in arr:
                num = re.sub(r'^US0+', '', str(x.get('patentNumber', '') or x.get('publicationNumber', '') or ''))
                aff = ''
                if isinstance(x.get('assignee'), str):
                    aff = x['assignee']
                elif isinstance(x.get('assignee'), list):
                    aff = ', '.join(a.get('name', '') or a.get('organization', '') or a.get('assignee_organization', '') or a.get('org_name', '') or '' for a in x['assignee'])
                elif x.get('assignee_organization'):
                    aff = x['assignee_organization']
                title = x.get('inventionTitle', '') or x.get('title', '') or ''
                date_str = x.get('grantDate', '') or x.get('publicationDate', '') or x.get('filingDate', '') or ''
                if title:
                    out.append({
                        'number': num, 'title': title, 'assignee': aff,
                        'year': date_str[:4], 'date': date_str,
                        'link': f'https://patents.google.com/patent/US{re.sub(r"^US", "", num)}' if num else '',
                    })
            return out
        except Exception:
            return []

async def lens_search(name: str, key: str) -> List[dict]:
    if not key:
        return []
    body = {
        'query_number': 1, 'query_text_1': 'title_claims_and_abstract',
        'query_value_1': name, 'size': 30,
        'include': ['title', 'publication_number', 'publication_date', 'assignee'],
    }
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            resp = await _request_with_retry(client, 'POST',
                'https://api.lens.org/scholarly/search',
                json=body,
                headers={'Authorization': f'Bearer {key}', 'Content-Type': 'application/json'},
                max_retries=2)
            if resp.status_code != 200:
                return []
            data = resp.json()
            results = data.get('results') or []
            if not isinstance(results, list):
                results = data.get('data', results) or []
            out = []
            for x in results:
                pub_num = x.get('publication_number', '') or ''
                title = x.get('title', '') or ''
                assignee_list = x.get('assignee') or []
                aff = ', '.join(a.get('assignee_organization', '') or '' for a in assignee_list if isinstance(a, dict)) if isinstance(assignee_list, list) else str(assignee_list or '')
                date_str = (x.get('publication_date') or '')[:10]
                if pub_num or title:
                    out.append({
                        'number': pub_num, 'title': title, 'assignee': aff,
                        'year': date_str[:4] if len(date_str) >= 4 else '',
                        'date': date_str,
                        'link': f'https://www.lens.org/lens/patent/search/result?patent_id={pub_num}' if pub_num else '',
                    })
            return out
        except Exception:
            return []

async def _google_fetch_text(url: str, timeout: float = 8.0) -> str:
    """Fetch a Google Patents URL trying direct + fallback proxies (Python has no CORS,
    so direct is tried first; proxies only as fallback)."""
    attempts = [('直连', url, None)]
    attempts.append(('r.jina.ai', 'https://r.jina.ai/' + url, 5.0))
    attempts.append(('r.jina.ai', 'https://r.jina.ai/' + url, 10.0))
    attempts.append(('allorigins', 'https://api.allorigins.win/raw?url=' + quote(url, safe=''), None))
    attempts.append(('whateverorigin', 'https://whateverorigin.org/get?url=' + quote(url, safe=''), None))
    errors = []
    prev_jina = False
    for label, u, t in attempts:
        try:
            async with httpx.AsyncClient(timeout=t or timeout, follow_redirects=True) as client:
                resp = await _request_with_retry(client, 'GET', u,
                                                 headers={'Referer': 'https://patents.google.com/',
                                                          'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'},
                                                 max_retries=2)
                if resp.status_code != 200:
                    errors.append(f'{label}: HTTP {resp.status_code}')
                    continue
                body = resp.text
                # r.jina.ai returns Markdown-wrapped content
                if label == 'r.jina.ai' and 'Markdown Content:' in body:
                    body = body.split('Markdown Content:', 1)[1].strip()
                # whateverorigin wraps JSON as {"contents":"..."}
                if label == 'whateverorigin':
                    try:
                        w = json.loads(body)
                        if isinstance(w, dict) and isinstance(w.get('contents'), str):
                            body = w['contents']
                    except Exception:
                        pass
                if not body:
                    errors.append(f'{label}: 空响应')
                    continue
                return body
        except Exception as e:
            errors.append(f'{label}: {str(e)[:80]}')
        if label == 'r.jina.ai' and prev_jina:
            await asyncio.sleep(0.4)
        prev_jina = (label == 'r.jina.ai')
    raise RuntimeError('Google Patents 请求失败（已尝试: ' + '；'.join(errors) + '）')

async def google_xhr(q: str, num: int = 30) -> dict:
    exp = urlencode({'q': q, 'num': num})
    text = await _google_fetch_text('https://patents.google.com/xhr/query?url=' + exp)
    parsed = json.loads(text)
    if not isinstance(parsed, dict):
        raise RuntimeError('Google Patents XHR 响应非 JSON')
    return parsed

async def google_patent_detail(num_raw: str) -> Optional[dict]:
    num = _google_num(num_raw)
    if not num:
        return None
    try:
        html = await _google_fetch_text(f'https://patents.google.com/patent/{num}/en', timeout=7.0)
    except Exception:
        return None
    if '<meta' not in html and '<html' not in html:
        return None
    from lxml import html as lh
    try:
        doc = lh.fromstring(html)
    except Exception:
        return None

    def meta_val(name: str) -> str:
        els = doc.xpath(f'//meta[@name="{name}"]/@content')
        return (els[0] if els else '').strip()

    inventors = [v.strip() for v in doc.xpath('//meta[@name="DC.contributor"]/@content') if v.strip()][:15]
    citations = [v.strip() for v in doc.xpath('//meta[@name="DC.relation"]/@content') if v.strip()][:40]
    cpc_set = set()
    for c in doc.xpath('//*[@itemprop="Code"]/text()'):
        c = (c or '').strip()
        if len(c) >= 4 and re.search(r'\d', c):
            cpc_set.add(c)
    claim_els = doc.xpath('//*[contains(@class, "claim-text")]') or doc.xpath('//*[contains(@class, "claim")]')
    legal_m = re.search(r'Legal status[\s\S]*?\)\s*(?:<[^>]*>\s*)*([A-Za-z][A-Za-z ,.\-]{0,50})', html)
    return {
        'abstract': meta_val('DC.description'),
        'title': meta_val('DC.title'),
        'inventors': ', '.join(inventors),
        'citations': citations,
        'cpc': list(cpc_set)[:12],
        'legal': legal_m.group(1).strip() if legal_m else '',
        'claims': len(claim_els),
    }

async def google_cites(num_raw: str) -> int:
    num = _google_num(num_raw)
    if not num:
        return 0
    try:
        d = await google_xhr('cites:' + num, 5)
        n = ((d.get('results') or {}).get('total_num_results') or 0)
        return int(n)
    except Exception:
        return 0

async def google_assignee_portfolio(assignee: str) -> List[dict]:
    a = str(assignee or '').replace('"', ' ').replace('\\', ' ').strip()
    if not a or not re.search(r'[\w\u4e00-\u9fa5]', a):
        return []
    try:
        d = await google_xhr(f'assignee:"{a}"', 20)
        return google_patents_parse(d)[:12]
    except Exception:
        return []

async def mcp_patent_details(patents: List[dict], tools, endpoint: str, key: str, session, max_n: int) -> None:
    detail_tool = None
    for t in (tools or []):
        if re.search(r'^get[_\-\s]?patent|patent[_\-\s]?detail', t.get('name', ''), re.I):
            detail_tool = t
            break
    if not detail_tool:
        return
    async with httpx.AsyncClient(timeout=60) as client:
        for p in patents[:max_n]:
            num = _google_num(p.get('number', ''))
            if not num:
                continue
            try:
                args = _build_mcp_args(detail_tool.get('schema'), detail_tool.get('required', []), num, num)
                text = await _mcp_tool_call(client, endpoint, key, session, detail_tool['name'], args)
                parsed = _parse_mcp_text(text)
                d = parsed.get('data') if isinstance(parsed, dict) else parsed
                if isinstance(d, dict) and isinstance(d.get('result'), dict):
                    d = d['result']
                if not isinstance(d, dict):
                    continue
                inv = ','.join(i if isinstance(i, str) else (i.get('name') or i.get('inventor_name') or '') for i in (d.get('inventors') or [])) if isinstance(d.get('inventors'), list) else (d.get('first_inventor') or '')
                p['mcp_status'] = d.get('status') or d.get('status_date') or d.get('legal_status') or ''
                if inv:
                    p['inventors'] = inv
                p['mcp_classification'] = d.get('classification') or ''
                p['mcp_abstract'] = d.get('abstract') or d.get('description') or ''
                if p.get('mcp_abstract') or p.get('mcp_status') or p.get('mcp_classification'):
                    p['mcp_detail'] = True
            except Exception:
                continue

async def deep_enrich_patents(patents: List[dict], *, src: str, key: str = '', target: str = '',
                               tools=None, endpoint: str = '', session=None,
                               max_detail: int = 5, progress_cb=None) -> Tuple[List[dict], dict]:
    out = []
    for p in patents:
        out.append(dict(p))
    insight = {'cites': {}, 'assignees': {}}
    max_n = min(max_detail, len(out))
    if src == 'mcp':
        if progress_cb:
            progress_cb('🔍 MCP 专利详情（状态/发明人/分类）…')
        try:
            await mcp_patent_details(out, tools, endpoint, key, session, max_n)
        except Exception:
            pass
    else:
        for i, p in enumerate(out[:max_n]):
            if progress_cb:
                progress_cb(f'🔍 专利详情增强 {i+1}/{max_n} …')
            num = _google_num(p.get('number', ''))
            if not num:
                continue
            try:
                d = await google_patent_detail(num)
                if d:
                    if d.get('abstract'):
                        p['abstract'] = d['abstract']
                    if d.get('claims'):
                        p['claims'] = d['claims']
                    if d.get('citations'):
                        p['citations'] = d['citations']
                    if d.get('cpc'):
                        p['cpc'] = d['cpc']
                    if d.get('legal'):
                        p['detail_legal'] = d['legal']
                    if d.get('inventors') and not p.get('inventors'):
                        p['inventors'] = d['inventors']
            except Exception:
                pass
            try:
                c = await google_cites(num)
                if c:
                    p['cited_by'] = c
            except Exception:
                pass
            if i < max_n - 1:
                await asyncio.sleep(0.3)
        if src == 'google':
            counts = {}
            for p in out:
                if p.get('assignee'):
                    counts[p['assignee']] = counts.get(p['assignee'], 0) + 1
            top_a = [a for a, _ in sorted(counts.items(), key=lambda x: -x[1])[:2]]
            for a in top_a:
                if progress_cb:
                    progress_cb(f'🧭 申请人布局分析: {a}')
                try:
                    insight['assignees'][a] = await google_assignee_portfolio(a)
                except Exception:
                    pass
                await asyncio.sleep(0.3)
    cn_patents = [p for p in out if p.get('application_no')]
    for i, p in enumerate(cn_patents[:4]):
        if progress_cb:
            progress_cb(f'🇨🇳 中国专利深度增强 {i + 1}/{min(4, len(cn_patents))} …')
        try:
            d9 = await patent9_details(p['application_no'])
            if d9:
                if d9.get('abstract'):
                    p['abstract'] = d9['abstract']
                if d9.get('cpc'):
                    p['cpc'] = d9['cpc']
                if d9.get('inventors'):
                    p['inventors'] = d9['inventors']
                if d9.get('main_ipc'):
                    p['main_ipc'] = d9['main_ipc']
                pr = d9.get('priority', '')
                if pr:
                    pm = re.search(r'(\d{4}\.\d{2}\.\d{2})', pr)
                    if pm:
                        p['priority_date'] = pm.group(1)
                    pc = re.findall(r'\b([A-Z]{2})\b', pr.split('\n')[0] if '\n' in pr else pr)
                    if pc and not p.get('countries'):
                        p['countries'] = pc
                if d9.get('agent'):
                    p['agent'] = d9['agent']
                if d9.get('wo'):
                    p['wo'] = d9['wo']
        except Exception:
            pass
        if not p.get('cn_fulltext') and p.get('cn_download'):
            p['cn_fulltext'] = await drugfuture_cn_fulltext(p['application_no'], progress_cb)
        if i < len(cn_patents[:4]) - 1:
            await asyncio.sleep(0.8)
    return out, insight

async def google_patents_search(target: str) -> List[dict]:
    terms = _expand_query_terms(target)[:3]
    q_terms = [f'"{t.replace(chr(34), " ").strip()}"' for t in terms]
    q = ' OR '.join(q_terms) + ' in:title' if q_terms else f'"{target}"'
    data = await google_xhr(q, 30)
    return google_patents_parse(data)

def google_patents_parse(data: dict) -> List[dict]:
    items = []
    clusters = data.get('results', {}).get('cluster', []) if isinstance(data.get('results'), dict) else []
    for cl in clusters:
        for r in (cl.get('result') or []):
            if r and r.get('patent'):
                items.append(r)
    out = []
    for x in items:
        p = x.get('patent') or {}
        num = str(p.get('publication_number', '') or x.get('id', '') or '')
        dt = p.get('grant_date', '') or p.get('publication_date', '') or p.get('filing_date', '') or p.get('priority_date', '') or ''
        fm = p.get('family_metadata') or {}
        fm_agg = fm.get('aggregated', {}).get('country_status', []) or []
        title = _strip_tags(p.get('title', '') or '')
        assignee = _strip_tags(p.get('assignee', '') or '')
        link_num = re.sub(r'/.*$', '', re.sub(r'^patent/', '', num)) if num else ''
        if title and num:
            out.append({
                'number': num, 'title': title, 'assignee': assignee,
                'year': dt[:4], 'date': dt, 'priority_date': p.get('priority_date', '') or '',
                'inventors': _strip_tags(p.get('inventor', '') or ''),
                'snippet': _strip_tags(p.get('snippet', '') or ''),
                'legal': _family_legal(fm),
                'countries': [c.get('country_code', '') for c in fm_agg if c.get('country_code')][:12],
                'link': f'https://patents.google.com/patent/{link_num}/en' if link_num else '',
            })
    return out

def _family_legal(fm: dict) -> str:
    try:
        agg = fm.get('aggregated') or {}
        status_list = agg.get('country_status') or []
        best = next((c for c in status_list if str(c.get('country_code', '')).upper() == 'US'), None) or (status_list[0] if status_list else None)
        if not best:
            return ''
        st = (best.get('best_patent_stage') or {}).get('state', '') or ''
        cc = best.get('country_code', '') or ''
        return f'{st} · {cc}' if st and cc else (st or cc)
    except Exception:
        return ''

async def espacenet_search(target: str, key: str) -> List[dict]:
    parts = (key or '').split()
    if len(parts) < 2:
        return []
    async with httpx.AsyncClient(timeout=30) as client:
        try:
            tr = await _request_with_retry(client, 'POST',
                'https://auth.epo.org/oauth2/token',
                data={'grant_type': 'client_credentials', 'client_id': parts[0], 'client_secret': parts[1]},
                max_retries=2)
            if tr.status_code != 200:
                return []
            tj = tr.json()
            token = tj.get('access_token', '')
            if not token:
                return []
            terms = _expand_query_terms(target)[:3]
            q_terms = [f'"{t.replace(chr(34), " ").strip()}"' for t in terms]
            q = ' OR '.join(q_terms) if q_terms else target
            url = f'https://ops.epo.org/3.2/rest-services/published-data/search?{urlencode({"q": q, "Range": "1-25"})}'
            resp = await _request_with_retry(client, 'GET', url,
                headers={'Authorization': f'Bearer {token}', 'Accept': 'application/xml'},
                max_retries=2)
            if resp.status_code != 200:
                return []
            return _parse_espacenet_xml(resp.text)
        except Exception:
            return []

def _parse_espacenet_xml(xml: str) -> List[dict]:
    out = []
    try:
        root = ET.fromstring(xml)
        for doc in root.findall('.//{http://www.epo.org/exchange}exchange-document') or root.findall('.//exchange-document'):
            try:
                bd = doc.find('.//{http://www.epo.org/exchange}bibliographic-data') or doc.find('.//bibliographic-data')
                pub = bd.find('.//{http://www.epo.org/exchange}publication-reference') or bd.find('.//publication-reference') if bd else None
                did = pub.find('.//{http://www.epo.org/exchange}document-id') or pub.find('.//document-id') if pub else None
                cc = (did.findtext('{http://www.epo.org/exchange}country', '') or did.findtext('country', '')) if did else ''
                dn = (did.findtext('{http://www.epo.org/exchange}doc-number', '') or did.findtext('doc-number', '')) if did else ''
                kind = (did.findtext('{http://www.epo.org/exchange}kind', '') or did.findtext('kind', '')) if did else ''
                date = (did.findtext('{http://www.epo.org/exchange}date', '') or did.findtext('date', '')) if did else ''
                def _join_text(parent, *tags):
                    parts = []
                    for tag in tags:
                        for el in (parent.findall(f'.//{{http://www.epo.org/exchange}}{tag}') or parent.findall(f'.//{tag}')):
                            t = ''.join(el.itertext()).strip()
                            if t:
                                parts.append(t)
                    return ', '.join(parts)
                title = _join_text(bd, 'invention-title') if bd else ''
                applicant = _join_text(bd, 'applicant-name') if bd else ''
                number = (cc + dn + kind).strip()
                year = date[:4]
                if number or title:
                    out.append({
                        'number': number, 'title': title, 'assignee': applicant,
                        'year': year, 'date': date,
                        'link': f'https://worldwide.espacenet.com/patent/search?q=pn={quote(number)}' if number else '',
                    })
            except Exception:
                continue
    except Exception:
        pass
    return out

try:
    import ddddocr as _ddddocr
    _OCR = _ddddocr.DdddOcr(show_ad=False)
except Exception:
    _OCR = None


# ─── Patent9 (中国专利, 免Key) ────────────────────────────────────────────────


def _p9_text(part: str) -> str:
    t = re.sub(r'<br\s*/?>', '\n', part)
    t = re.sub(r'</(?:div|li|p|tr|td)>', '\n', t, flags=re.I)
    t = re.sub(r'<[^>]+>', '', t)
    t = re.sub(r'&nbsp;', ' ', t)
    t = re.sub(r'[ \t\u3000]+', ' ', t)
    t = re.sub(r'\n\s*\n+', '\n', t)
    return t


def _p9_grab(text: str, label: str) -> str:
    m = re.search(re.escape(label) + r'[：:]\s*([^\n]*?)(?=\s*(?:申请号|公开号|主分类号|申请人|申请日|公开日|发明人|摘要)[：:]|\n|$)', text)
    return m.group(1).strip() if m else ''


def _parse_patent9_results(html: str, max_results: int) -> List[dict]:
    parts = re.split(r'<div class="title">', html)
    out = []
    for part in parts[1:]:
        if len(out) >= max_results:
            break
        title_line = re.sub(r'<[^>]+>', '', part[:part.find('</div>')])
        title_line = re.sub(r'\s+', ' ', title_line).strip()
        m = re.match(r'\d+[：:]\s*\[([^\]]*)\]\s*(.*)$', title_line)
        kind = m.group(1) if m else ''
        title = re.sub(r'\s*#imgabs\d+#\s*', '', (m.group(2) if m else title_line)).strip()
        no_m = re.search(r"PatentNo:'([^']+)'", part)
        appno = no_m.group(1) if no_m else ''
        text = _p9_text(part)
        pub = _p9_grab(text, '公开号')
        pubdate = _p9_grab(text, '公开日')
        abstract = _p9_grab(text, '摘要')
        abstract = re.sub(r'\s*#imgabs\d+#\s*', '', abstract)
        cn_dl = ''
        cdm = re.search(r'href="(http[^"]*verify\.aspx\?cnpatentno=[^"]+)"', part)
        if cdm:
            cn_dl = cdm.group(1)
        elif appno:
            cn_dl = f'{DRUGFUTURE_URL}/verify.aspx?cnpatentno={appno}'
        if not appno and not pub:
            continue
        out.append({
            'number': pub or appno,
            'application_no': appno,
            'title': title,
            'kind': kind,
            'assignee': _p9_grab(text, '申请人'),
            'year': pubdate[:4],
            'date': pubdate,
            'inventors': _p9_grab(text, '发明人'),
            'abstract': abstract,
            'snippet': abstract,
            'priority_date': '',
            'cpc': [],
            'link': '',
            'cn_download': cn_dl,
            'countries': ['CN'],
            'source': 'Patent9',
        })
    return out


async def patent9_search(target: str, max_results: int = 10) -> List[dict]:
    terms = _expand_query_terms(target)[:3]
    query = ' '.join(terms) if terms else target
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
        'Accept-Language': 'zh-CN,zh;q=0.9',
    }
    async with httpx.AsyncClient(timeout=30, follow_redirects=True, headers=headers) as client:
        try:
            resp = await _request_with_retry(client, 'GET', PATENT9_URL + '/search.aspx',
                                             params={'k': query, 't': 'all',
                                                     'c': 'PatentName,Abstract', 'p': '1'},
                                             max_retries=2)
            if resp.status_code != 200:
                return []
            return _parse_patent9_results(resp.text, max_results)
        except Exception:
            return []


def _parse_patent9_details(html: str, application_no: str) -> dict:
    fields = {}
    for m in re.finditer(r'<tr>\s*<th>([^<]*)</th>\s*<td>(.*?)</td>\s*</tr>', html, re.I | re.S):
        label = m.group(1).strip()
        value = re.sub(r'<[^>]+>', '', m.group(2))
        value = re.sub(r'\s+', ' ', value).strip()
        if value:
            fields[label] = value
    if not fields:
        return {}
    abstract_m = re.search(r'<div class="abstract-box">(.*?)</div>', html, re.I | re.S)
    if abstract_m:
        fields['摘要'] = re.sub(r'\s+', ' ', re.sub(r'<[^>]+>', '', abstract_m.group(1))).strip()
    return {
        'application_no': fields.get('专利申请号', application_no),
        'app_date': fields.get('申请日', ''),
        'pub_no': fields.get('公开(公告)号', ''),
        'pub_date': fields.get('公开(公告)日', ''),
        'main_ipc': fields.get('主分类号', ''),
        'cpc': fields.get('分类号', '').split(),
        'priority': fields.get('优先权', ''),
        'assignee': fields.get('申请(专利权)人', ''),
        'address': fields.get('地址', ''),
        'inventors': fields.get('发明(设计)人', ''),
        'pct': fields.get('国际申请', ''),
        'wo': fields.get('国际公布', ''),
        'agent': fields.get('专利代理机构', ''),
        'patent_type': fields.get('专利类型', ''),
        'abstract': fields.get('摘要', ''),
    }


async def patent9_details(application_no: str) -> dict:
    if not application_no:
        return {}
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
        'Accept-Language': 'zh-CN,zh;q=0.9',
    }
    async with httpx.AsyncClient(timeout=30, follow_redirects=True, headers=headers) as client:
        try:
            resp = await _request_with_retry(client, 'POST', PATENT9_URL + '/PatentDetails.aspx',
                                             data={'PatentNo': application_no}, max_retries=2)
            if resp.status_code != 200:
                return {}
            return _parse_patent9_details(resp.text, application_no)
        except Exception:
            return {}


# ─── drugfuture 中国专利全文 (验证码 OCR, 可选依赖) ──────────────────────────


async def drugfuture_cn_fulltext(application_no: str, progress_cb=None) -> str:
    """
    Resolve the drugfuture PDF full-text URL for a CN application number.
    Requires the optional ddddocr package for the graphic CAPTCHA.
    Returns '' (not None) on any failure so callers can skip quietly.
    """
    if not _OCR or not application_no:
        return ''
    if progress_cb:
        progress_cb(f'🛡️ drugfuture 验证码识别 {application_no} …')
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
        'Referer': f'{DRUGFUTURE_URL}/verify.aspx',
    }
    async with httpx.AsyncClient(timeout=30, follow_redirects=True, headers=headers) as client:
        try:
            r = await client.post(f'{DRUGFUTURE_URL}/verify.aspx',
                                  data={'cnpatentno': application_no, 'Common': '1'})
            if r.status_code != 200:
                return ''
            fm = re.search(r'<form[^>]*action="([^"]+)"[^>]*>', r.text, re.I)
            action = fm.group(1) if fm else 'search.aspx'
            if action.startswith('/'):
                action = f'{DRUGFUTURE_URL}{action}'
            elif not action.startswith('http'):
                action = f'{DRUGFUTURE_URL}/{action}'
            img_m = re.search(r'src="([^"]*VerifyCode\.aspx[^"]*)"', r.text, re.I)
            if not img_m:
                return ''
            img_url = img_m.group(1)
            if img_url.startswith('/'):
                img_url = 'http://www2.drugfuture.com' + img_url
            elif not img_url.startswith('http'):
                img_url = f'{DRUGFUTURE_URL}/{img_url}'
            ir = await client.get(img_url)
            if ir.status_code != 200 or not ir.content:
                return ''
            try:
                code = str(_OCR.classification(ir.content)).strip()
            except Exception:
                return ''
            if not code:
                return ''
            rr = await client.post(action, data={'cnpatentno': application_no, 'common': '1', 'ValidCode': code})
            if rr.status_code != 200:
                return ''
            for m in re.finditer(r'href="([^"]+\.pdf[^"]*)"', rr.text, re.I):
                url = m.group(1)
                if url.startswith('/'):
                    url = 'http://www2.drugfuture.com' + url
                return url
            return ''
        except Exception:
            return ''


async def mcp_patent_search(target: str, endpoint: str, key: str, tool_pref: str = '') -> dict:
    if not endpoint:
        raise ValueError('MCP endpoint required')
    async with httpx.AsyncClient(timeout=60) as client:
        init = await _mcp_request(client, endpoint, key, {'jsonrpc': '2.0', 'id': 1, 'method': 'initialize',
            'params': {'protocolVersion': '2025-06-18', 'capabilities': {}, 'clientInfo': {'name': 'TargetInfo', 'version': '1.0'}}})
        session = (init or {}).get('_sessionId')
        tools_data = await _mcp_request(client, endpoint, key, {'jsonrpc': '2.0', 'id': 2, 'method': 'tools/list', 'params': {}}, session)
        tools = []
        for t in ((tools_data or {}).get('result', {}).get('tools') or []):
            tools.append({
                'name': t.get('name', ''),
                'schema': (t.get('inputSchema') or {}).get('properties') if t.get('inputSchema') else None,
                'required': (t.get('inputSchema') or {}).get('required', []) if t.get('inputSchema') else [],
            })
        if not tools:
            raise ValueError('MCP server provided no tools')
        tool_name = tool_pref if tool_pref and any(t['name'] == tool_pref for t in tools) else tools[0]['name']
        terms = _expand_query_terms(target)[:2]
        q = ' OR '.join(terms) if terms else target
        t_schema = next((t for t in tools if t['name'] == tool_name), {})
        args = _build_mcp_args(t_schema.get('schema'), t_schema.get('required', []), q, target)
        try:
            text = await _mcp_tool_call(client, endpoint, key, session, tool_name, args)
        except Exception as e:
            if 'required' in str(e).lower() or 'invalid' in str(e).lower() or 'schema' in str(e).lower():
                alt = _build_mcp_args(t_schema.get('schema'), t_schema.get('required', []), target, q)
                text = await _mcp_tool_call(client, endpoint, key, session, tool_name, alt)
            else:
                raise
        parsed = _parse_mcp_text(text)
        data = parsed.get('data') if isinstance(parsed, dict) else None
        if data is None:
            data = parsed
        arr = data if isinstance(data, list) else (data.get('patents') or data.get('results') or data.get('data') or data.get('items') or [])
        if not isinstance(arr, list):
            arr = arr.get('patents') or arr.get('results') or arr.get('items') or arr.get('hits') or []
        patents = [_normalize_patent(p) for p in arr] if isinstance(arr, list) else []
        patents = [p for p in patents if p and (p.get('number') or p.get('title'))]
        return {'patents': patents, 'tools': tools, 'session': session, 'endpoint': endpoint, 'key': key}

async def _mcp_request(client, endpoint, key, body, session=None):
    headers = {'Content-Type': 'application/json', 'Accept': 'application/json, text/event-stream', 'MCP-Protocol-Version': '2025-06-18'}
    if key:
        headers['Authorization'] = f'Bearer {key}'
    if session:
        headers['Mcp-Session-Id'] = session
    resp = await _request_with_retry(client, 'POST', endpoint, json=body, headers=headers, max_retries=2)
    raw = resp.text
    if resp.status_code != 200:
        raise RuntimeError(f'MCP HTTP {resp.status_code}: {raw[:200]}')
    sid = resp.headers.get('mcp-session-id') if hasattr(resp.headers, 'get') else None
    p = _parse_mcp_text(raw)
    if isinstance(p, dict) and p.get('kind') == 'json' and p.get('data'):
        result = p['data']
        if isinstance(result, dict) and result.get('error'):
            raise RuntimeError(f'MCP error: {result["error"].get("message", str(result["error"][:200]))}')
        if sid:
            result['_sessionId'] = sid
        return result
    if isinstance(p, dict) and p.get('kind') == 'text':
        return {'_text': p['data'], '_sessionId': sid}
    raise RuntimeError(f'MCP response parse failed: {raw[:200]}')

async def _mcp_tool_call(client, endpoint, key, session, name, args):
    body = {'jsonrpc': '2.0', 'id': 3, 'method': 'tools/call', 'params': {'name': name, 'arguments': args}}
    d = await _mcp_request(client, endpoint, key, body, session)
    result = d.get('result') or {}
    if result.get('isError'):
        raise RuntimeError(f'MCP tool {name} error: {str(result)[:200]}')
    texts = []
    for c in (result.get('content') or []):
        if c.get('type') == 'text' and c.get('text'):
            texts.append(c['text'])
        elif c.get('type') == 'resource':
            texts.append(json.dumps(c['resource']))
    return '\n'.join(texts) or d.get('_text', '')

def _parse_mcp_text(text: str) -> dict:
    t = (text or '').strip()
    if not t:
        return {'kind': 'none', 'data': None}
    if t.startswith('{'):
        try:
            return {'kind': 'json', 'data': json.loads(t)}
        except json.JSONDecodeError:
            pass
    lines = [l for l in t.split('\n') if re.match(r'^\s*data\s*:', l)]
    if lines:
        joined = '\n'.join(re.sub(r'^\s*data\s*:\s*', '', l) for l in lines)
        try:
            return {'kind': 'json', 'data': json.loads(joined)}
        except json.JSONDecodeError:
            return {'kind': 'text', 'data': joined}
    return {'kind': 'text', 'data': t}

def _build_mcp_args(schema, required, *values) -> dict:
    if not schema:
        return {'q': values[0], 'limit': 30}
    keys = list(schema.keys())
    req_key = next((k for k in required if re.match(r'^(q|query|search|text|question|prompt|input|keyword|keywords|terms?)$', k, re.I)), None)
    q_key = req_key or next((k for k in keys if re.match(r'^(query|search|q|keyword|keywords)$', k, re.I)), None) or next((k for k in keys if re.match(r'^(text|question|prompt|input|terms?)$', k, re.I)), None) or keys[0]
    out = {q_key: ''.join(str(v or '') for v in values)}
    if 'limit' in keys:
        out['limit'] = 30
    return out

def _normalize_patent(obj) -> Optional[dict]:
    if not obj or not isinstance(obj, dict):
        return None
    num = str(obj.get('number', '') or obj.get('publication_number', '') or obj.get('application_number', '') or obj.get('patent_number', '') or obj.get('patent_id', '') or obj.get('patentNumber', '') or obj.get('id', '') or '')
    title = obj.get('title', '') or obj.get('invention_title', '') or obj.get('inventionTitle', '') or obj.get('name', '') or ''
    dt = obj.get('grant_date', '') or obj.get('publication_date', '') or obj.get('filing_date', '') or obj.get('date', '') or obj.get('grantDate', '') or obj.get('publicationDate', '') or obj.get('filingDate', '') or ''
    aff = ''
    if isinstance(obj.get('assignee'), str):
        aff = obj['assignee']
    elif isinstance(obj.get('assignee'), list):
        aff = ', '.join(str(a.get('assignee_organization', '') or a.get('name', '') or a.get('organization', '') or '') for a in obj['assignee'] if isinstance(a, dict))
    elif isinstance(obj.get('applicants'), list):
        aff = ', '.join(str(a.get('name', '') or a.get('assignee_organization', '') or '') for a in obj['applicants'] if isinstance(a, dict))
    elif obj.get('assignee_organization'):
        aff = obj['assignee_organization']
    elif obj.get('first_applicant'):
        aff = obj['first_applicant']
    elif obj.get('inventor'):
        aff = ', '.join(str(i.get('inventor_name', '') or i.get('name', '') or '') for i in (obj['inventor'] if isinstance(obj['inventor'], list) else [obj['inventor']]))
    elif obj.get('first_inventor'):
        aff = obj['first_inventor']
    link_num = re.sub(r'/.*$', '', re.sub(r'^patent/', '', re.sub(r'^US0+', '', num))) if num else ''
    inv = ''
    if isinstance(obj.get('inventors'), list):
        inv = ', '.join(str(i.get('name', '') or i.get('inventor_name', '') or '') for i in obj['inventors'] if isinstance(i, dict)) or ', '.join(str(i) for i in obj['inventors'] if isinstance(i, str))
    elif obj.get('first_inventor'):
        inv = obj['first_inventor']
    return {
        'number': num, 'title': title, 'assignee': aff,
        'year': dt[:4], 'date': dt, 'inventors': inv, 'priority_date': obj.get('priority_date', '') or obj.get('priorityDate', '') or '',
        'snippet': obj.get('snippet', '') or obj.get('abstract', '') or '',
        'legal': obj.get('legal', '') or obj.get('legal_status', '') or '',
        'countries': obj.get('countries', []) if isinstance(obj.get('countries'), list) else [],
        'link': obj.get('link', '') or (f'https://patents.google.com/patent/{link_num}/en' if link_num else ''),
        'claims': int(obj.get('claims', 0) or 0),
        'citations': obj.get('citations', []) if isinstance(obj.get('citations'), list) else [],
        'cpc': obj.get('cpc', []) if isinstance(obj.get('cpc'), list) else [],
        'application_no': obj.get('application_no', '') or '',
        'cn_download': obj.get('cn_download', '') or '',
        'cn_fulltext': obj.get('cn_fulltext', '') or '',
        'kind': obj.get('kind', '') or '',
    }

async def search_patents(target: str, src: str = 'google', key: str = '',
                         deep: bool = True, mcp_url: str = '', mcp_tool: str = '',
                         progress_cb=None) -> dict:
    note, eff_src, insight, mcp_ctx = '', src, None, None
    if src == 'google':
        try:
            patents = await google_patents_search(target)
        except Exception:
            try:
                mcp_result = await mcp_patent_search(target, PATENT_MCP_FALLBACK, '', '')
                patents = mcp_result.get('patents', [])
                eff_src = 'mcp'
                mcp_ctx = mcp_result
                note = 'Google 代理链不可用，已自动回退到 MCP 免Key 来源（pipeworx）。'
            except Exception:
                note = '专利主源检索失败，仅返回 Patent9 中国专利结果。'
                try:
                    p9 = await patent9_search(target, 10)
                    normalized = []
                    for p in p9:
                        np = _normalize_patent(p)
                        if np and (np.get('number') or np.get('application_no')):
                            normalized.append(np)
                    return {'patents': normalized, 'insight': None, 'note': note}
                except Exception:
                    return {'patents': [], 'insight': None, 'note': '专利检索全部失败'}
        if not patents and not note:
            try:
                mcp_result = await mcp_patent_search(target, PATENT_MCP_FALLBACK, '', '')
                if mcp_result.get('patents'):
                    patents = mcp_result['patents']
                    eff_src = 'mcp'
                    mcp_ctx = mcp_result
                    note = 'Google 未检索到结果，已自动回退到 MCP 免Key 来源（pipeworx）。'
            except Exception:
                pass
    elif src == 'uspto':
        patents = await uspto_search(target, key)
    elif src == 'lens':
        patents = await lens_search(target, key)
    elif src == 'espacenet':
        patents = await espacenet_search(target, key)
    elif src == 'mcp':
        try:
            endpoint = (mcp_url or '').strip().rstrip('/')
            if not endpoint:
                return {'patents': [], 'insight': None, 'note': '请先在设置中填写 MCP endpoint'}
            mcp_result = await mcp_patent_search(target, endpoint, key, mcp_tool or '')
            patents = mcp_result.get('patents', [])
            mcp_ctx = mcp_result
        except Exception as e:
            try:
                p9 = await patent9_search(target, 10)
                normalized = []
                for p in p9:
                    np = _normalize_patent(p)
                    if np and (np.get('number') or np.get('application_no')):
                        normalized.append(np)
                if normalized:
                    return {'patents': normalized, 'insight': None,
                            'note': f'MCP 检索失败（{str(e)[:80]}），仅返回 Patent9 中国专利结果。'}
            except Exception:
                pass
            return {'patents': [], 'insight': None, 'note': f'MCP 检索失败: {str(e)[:120]}'}
    else:
        patents = await google_patents_search(target)
    patents = _dedupe_patents(patents or [])
    normalized = []
    for p in patents:
        np = _normalize_patent(p)
        if np:
            normalized.append(np)
    try:
        p9 = await patent9_search(target, 10)
        seen_n = {p.get('number') for p in normalized}
        for p in p9:
            np = _normalize_patent(p)
            if np and (np.get('number') or np.get('application_no')) and np.get('number') not in seen_n:
                normalized.append(np)
                seen_n.add(np['number'])
                if not note:
                    note = f'Patent9 自动补充中国专利 {len(p9)} 条。'
    except Exception:
        pass
    if deep and normalized:
        if progress_cb:
            progress_cb('🔍 专利深度增强…')
        try:
            enriched, insight = await deep_enrich_patents(
                normalized, src=eff_src, key=key, target=target,
                tools=(mcp_ctx or {}).get('tools'),
                endpoint=(mcp_ctx or {}).get('endpoint'),
                session=(mcp_ctx or {}).get('session'),
            )
            normalized = enriched
        except Exception:
            pass
    return {'patents': normalized, 'insight': insight, 'note': note}


# ─── LLM / AI ──────────────────────────────────────────────────────────────

LLM_PROVIDERS = {
    'deepseek': {
        'name': 'DeepSeek', 'base': 'https://api.deepseek.com/v1',
        'models': ['deepseek-chat', 'deepseek-v4-flash'],
    },
    'mimo': {
        'name': '小米 MiMo', 'base': 'https://api.xiaomimimo.com/v1',
        'models': ['mimo-v2.5', 'mimo-v2.5-pro', 'mimo-v2-flash'],
    },
    'zhipu': {
        'name': '智谱 GLM', 'base': 'https://open.bigmodel.cn/api/paas/v4',
        'models': ['GLM-4.7-Flash', 'GLM-4.7-Plus'],
        'web_search': True,
    },
    'custom': {
        'name': '自定义', 'base': '', 'editable': True,
        'models': ['gpt-4o-mini', 'deepseek-chat', 'moonshot-v1-8k', 'qwen-max'],
    },
}

LLM_PROMPTS = {
    'target_overview': '你是一个生物医药领域的专家。请根据以下靶点信息，生成一份专业的靶点概述（中文，200-500字），涵盖基因符号、蛋白类别、功能描述和相关疾病。\n\n{context}',
    'research_progress': '你是一个生物医药领域的专家。以下是关于某靶点的文献数据（年份分布和标题列表），请分析研究趋势，指出近三年的研究热点和方向（中文，150-400字）。\n\n{context}',
    'clinical_landscape': '你是一个临床研究专家。以下是关于某靶点的临床试验数据（阶段分布、适应症列表），请分析临床开发格局，指出当前阶段分布特征和主要探索方向（中文，150-400字）。\n\n{context}',
    'key_findings': '你是一个生物医药领域的专家。请根据以下文献标题列表，提取最有代表性的3-5个研究发现，逐条列出并简要说明其意义（中文，200-600字）。\n\n{context}',
    'future_outlook': '你是一个药物研发专家。以下是针对某靶点的药物管线数据（已获批和在研药物），请分析其未来发展趋势和值得关注的方向（中文，150-400字）。\n\n{context}',
    'web_summary': lambda name, gene, cls: (
        f'你是一名生物医药行业情报分析师。请利用联网搜索能力，针对以下靶点，搜集并整理最新全网关键情报（中文，400-800字），涵盖以下方面：\n'
        f'1. 关键里程碑事件：重要的研究发现、学术突破、指南更新等\n'
        f'2. 重大收购与合作：药企之间的并购、授权引进、战略合作等\n'
        f'3. 临床试验重大进展：近期成功的 III 期结果、FDA 突破性疗法认定、重大失败或终止的试验及其影响\n'
        f'4. 监管动态：FDA/EMA/NMPA 批准、加速审批、孤儿药认定等\n'
        f'5. 竞争格局变化：新进入者、专利到期、生物类似药进展等\n\n'
        f'请确保每条信息标注来源（如新闻标题或机构名称），并说明该信息对药物研发的潜在影响。\n\n'
        f'靶点名称: {name}\n基因符号: {gene or "—"}\n蛋白类别: {cls or "—"}'
    ),
    'web_summary_legacy': lambda ctx: (
        f'你是一名生物医药行业情报分析师，擅长基于结构化检索数据做专业研判。请完全基于下面已核实检索到的数据（文献、临床、药物、专利、靶点功能），不要臆造或引用检索数据之外的新信息，生成一份「综合情报研判」（中文，300-600字），涵盖：\n'
        f'1. 研究热度与判断\n'
        f'2. 临床开发格局与关键进展\n'
        f'3. 药物管线态势（含已获批/在研/活性分子线索）\n'
        f'4. 专利布局概况与FTO/竞争风险提示\n'
        f'5. 总体研发建议与后续关注点\n\n'
        f'检索数据：\n{ctx}'
    ),
    'patent_landscape': lambda name, txt: (
        f'你是一名生物医药知识产权分析师。请基于以下围绕靶点「{name}」的专利检索数据，撰写专业专利调研报告（中文，300-700字），涵盖：\n'
        f'1. 专利布局概况：主要申请人/权利人及活跃度\n'
        f'2. 申请年份趋势与所处阶段\n'
        f'3. 主要技术方向（化合物/制剂/用途/制备等）\n'
        f'4. 值得关注的近期关键专利\n'
        f'5. FTO/自由实施风险提示与竞争格局研判\n\n'
        f'检索数据：\n{txt}'
    ),
}

async def llm_chat(user: str, provider: str, api_key: str, model: str, system: str = '', base_url: str = '') -> str:
    if not api_key or not model:
        return ''
    p = LLM_PROVIDERS.get(provider, LLM_PROVIDERS['deepseek'])
    base = base_url or p['base']
    messages = [{'role': 'system', 'content': system}, {'role': 'user', 'content': user}] if system else [{'role': 'user', 'content': user}]
    body = {'model': model, 'messages': messages, 'temperature': 0.6, 'stream': False}
    if provider == 'zhipu':
        body['max_tokens'] = 4096
    async with httpx.AsyncClient(timeout=120) as client:
        try:
            resp = await _request_with_retry(client, 'POST', f'{base}/chat/completions',
                json=body, headers={'Authorization': f'Bearer {api_key}', 'Content-Type': 'application/json'},
                max_retries=2)
            if resp.status_code != 200:
                return ''
            text = (resp.json().get('choices') or [{}])[0].get('message', {}).get('content', '') or ''
            return text.strip()
        except Exception:
            return ''

async def llm_json(user: str, provider: str, api_key: str, model: str, system: str = '', base_url: str = '') -> dict:
    if not api_key or not model:
        return {}
    p = LLM_PROVIDERS.get(provider, LLM_PROVIDERS['deepseek'])
    base = base_url or p['base']
    messages = [{'role': 'system', 'content': system}, {'role': 'user', 'content': user}] if system else [{'role': 'user', 'content': user}]

    def _clean_and_parse(text: str) -> Optional[dict]:
        if not text:
            return None
        cleaned = re.sub(r'^```json\s*', '', text).replace('```', '').strip()
        try:
            o = json.loads(cleaned)
            if isinstance(o, dict):
                return o
        except Exception:
            pass
        m = re.search(r'\{[\s\S]*\}', text)
        if m:
            try:
                o = json.loads(m.group(0))
                if isinstance(o, dict):
                    return o
            except Exception:
                pass
        return None

    async def _post(body: dict):
        try:
            resp = await _request_with_retry(client, 'POST', f'{base}/chat/completions',
                json=body, headers={'Authorization': f'Bearer {api_key}', 'Content-Type': 'application/json'},
                max_retries=2)
            if resp.status_code != 200:
                return None
            return (resp.json().get('choices') or [{}])[0].get('message', {}).get('content', '') or ''
        except Exception:
            return None

    async with httpx.AsyncClient(timeout=120) as client:
        # Strategy 1: response_format json_object
        try:
            text = await _post({'model': model, 'messages': messages, 'temperature': 0.4,
                                'stream': False, 'response_format': {'type': 'json_object'}})
            o = _clean_and_parse(text)
            if o:
                return o
        except Exception:
            pass
        # Strategy 2: plain chat then extract JSON (some providers reject response_format)
        try:
            text = await _post({'model': model, 'messages': messages, 'temperature': 0.4,
                                'stream': False})
            o = _clean_and_parse(text)
            if o:
                return o
        except Exception:
            pass
        return {}

async def llm_chat_with_search(user: str, provider: str, api_key: str, model: str, system: str = '', base_url: str = '') -> str:
    if not api_key or not model:
        return ''
    p = LLM_PROVIDERS.get(provider, LLM_PROVIDERS['deepseek'])
    base = base_url or p['base']
    messages = [{'role': 'system', 'content': system}, {'role': 'user', 'content': user}] if system else [{'role': 'user', 'content': user}]
    async with httpx.AsyncClient(timeout=120) as client:
        # Strategy 1: web_search param
        body = {'model': model, 'messages': messages, 'temperature': 0.5, 'stream': False, 'web_search': True}
        try:
            resp = await _request_with_retry(client, 'POST', f'{base}/chat/completions',
                json=body, headers={'Authorization': f'Bearer {api_key}', 'Content-Type': 'application/json'},
                max_retries=1)
            if resp.status_code == 200:
                text = (resp.json().get('choices') or [{}])[0].get('message', {}).get('content', '') or ''
                if text.strip():
                    return text.strip()
        except Exception:
            pass
        # Strategy 2: tools param
        body = {'model': model, 'messages': messages, 'temperature': 0.5, 'stream': False,
                'tools': [{'type': 'web_search', 'web_search': {'enable': True}}]}
        try:
            resp = await _request_with_retry(client, 'POST', f'{base}/chat/completions',
                json=body, headers={'Authorization': f'Bearer {api_key}', 'Content-Type': 'application/json'},
                max_retries=1)
            if resp.status_code == 200:
                text = (resp.json().get('choices') or [{}])[0].get('message', {}).get('content', '') or ''
                if text.strip():
                    return text.strip()
        except Exception:
            pass
    # Fallback: regular chat
    return await llm_chat(user, provider, api_key, model, system, base_url)

def _extract_refs(text: str) -> List[str]:
    refs = set()
    for m in re.findall(r'PMID:\s*\d+', text or ''):
        refs.add(m)
    for m in re.findall(r'NCT\d+', text or ''):
        refs.add(m)
    for m in re.findall(r'US\d+[A-Z]\d*', text or ''):
        refs.add(m)
    for m in re.findall(r'EP\d+[A-Z]\d*', text or ''):
        refs.add(m)
    for m in re.findall(r'https?://[^\s\]\\)]+', text or ''):
        refs.add(m)
    return list(refs)[:20]


# ─── Report Building ──────────────────────────────────────────────────────────


def build_report(target_name: str, gene: str, mutation: str,
                 target_detail: Optional[TargetDetail],
                 papers: List[Paper], trials: List[ClinicalTrial],
                 drugs: List[DrugInfo], patents: List[dict],
                 molecules: List[MoleculeInfo]) -> ReportContent:
    report = ReportContent(
        target_name=target_name, gene=gene, mutation=mutation,
        target_detail=target_detail,
        papers=papers, trials=trials, drugs=drugs,
        molecules=molecules,
    )
    # Convert patent dicts to PatentInfo objects
    report.patents = []
    for p in (patents or []):
        report.patents.append(PatentInfo(
            number=str(p.get('number', '') or ''),
            title=str(p.get('title', '') or ''),
            assignee=str(p.get('assignee', '') or ''),
            year=str(p.get('year', '') or ''),
            date=str(p.get('date', '') or ''),
            link=str(p.get('link', '') or ''),
            snippet=str(p.get('snippet', '') or ''),
            abstract=str(p.get('abstract', '') or ''),
            claims=int(p.get('claims', 0) or 0),
            citations=list(p.get('citations', []) or []),
            cpc=list(p.get('cpc', []) or []),
            detail_legal=str(p.get('detail_legal', '') or ''),
            inventors=str(p.get('inventors', '') or ''),
            cited_by=int(p.get('cited_by', 0) or 0),
            countries=list(p.get('countries', []) or []),
            mcp_status=str(p.get('mcp_status', '') or ''),
            mcp_classification=str(p.get('mcp_classification', '') or ''),
            mcp_abstract=str(p.get('mcp_abstract', '') or ''),
            priority_date=str(p.get('priority_date', '') or ''),
            application_no=str(p.get('application_no', '') or ''),
            cn_download=str(p.get('cn_download', '') or ''),
            cn_fulltext=str(p.get('cn_fulltext', '') or ''),
        ))

    # Target overview (rule-based)
    d = target_detail
    if d and d.gene_symbol:
        parts = [f'靶点: {target_name}']
        if d.gene_symbol: parts.append(f'基因: {d.gene_symbol}')
        if mutation: parts.append(f'突变: {mutation}')
        if d.protein_class: parts.append(f'蛋白类别: {d.protein_class}')
        if d.protein_name: parts.append(f'蛋白全名: {d.protein_name}')
        if d.description: parts.append(f'功能: {d.description}')
        if d.subcellular: parts.append(f'亚细胞定位: {d.subcellular}')
        if d.protein_families: parts.append(f'蛋白家族: {d.protein_families}')
        if d.tractability: parts.append(f'成药性评估(Tractability): {", ".join(d.tractability)}')
        if d.pdb: parts.append(f'PDB结构: {"; ".join(p.pdb_id + " (" + p.method + (", " + p.resolution + "Å" if p.resolution else "") + ")" for p in d.pdb[:5])}')
        if d.alphafold: parts.append(f'AlphaFold预测: {len(d.alphafold)}个结构')
        if d.string_interactions: parts.append(f'蛋白互作(STRING): {len(d.string_interactions)}个互作伙伴')
        if d.kegg_pathways: parts.append(f'KEGG通路: {"; ".join(p.name for p in d.kegg_pathways[:5])}')
        if d.hpa: parts.append(f'HPA表达谱: 组织特异 {d.hpa.protein_tissue or "—"}')
        if d.gtex: parts.append(f'GTEx组织表达: {len(d.gtex)}组织')
        if d.clinvar: parts.append(f'ClinVar变异: {len(d.clinvar)}条记录')
        if d.clinpgx: parts.append(f'ClinPGx药物基因组关联: {len(d.clinpgx)}条')
        if d.dgidb: parts.append(f'DGIdb药物互作: {len(d.dgidb)}条')
        if d.related_diseases: parts.append(f'相关疾病: {", ".join(d.related_diseases)}')
        report.target_overview_raw = '\n'.join(parts)
    else:
        report.target_overview_raw = f'靶点 {target_name} 的基本信息暂未获取到。'

    # Research progress
    yc = _count_by_year(papers)
    yrs = list(yc.keys())
    now = datetime.now().year
    if yrs:
        recent_years = [str(now - 2), str(now - 1), str(now)]
        trend = '上升' if len(yrs) >= 2 and yc[yrs[-1]] > yc[yrs[0]] else '波动'
        rcnt = sum(yc.get(y, 0) for y in recent_years)
        ptc = {}
        for p in papers:
            for t in p.pub_type:
                ptc[t] = ptc.get(t, 0) + 1
        pt_str = '、'.join(f'{k} {v}篇' for k, v in sorted(ptc.items(), key=lambda x: -x[1])[:4])
        report.research_progress_raw = (
            f'共检索到 {len(papers)} 篇文献，覆盖 {len(yrs)} 个年份'
            f'（{yrs[0]}–{yrs[-1]}），整体呈 {trend} 趋势。'
            f'近三年（{recent_years[0]}–{recent_years[-1]}）共发表 {rcnt} 篇。'
            + (f'\n文献类型：{pt_str}。' if pt_str else '')
        )
    else:
        report.research_progress_raw = f'共检索到 {len(papers)} 篇文献。'

    # Clinical landscape
    ph = _count_by_phase(trials)
    pk = list(ph.keys())
    cls_text = f'共 {len(trials)} 项试验，阶段：{", ".join(k + " " + str(ph[k]) + "项" for k in pk)}。' if pk else f'共 {len(trials)} 项试验。'
    sources = list(set(t.source for t in trials if t.source != 'ClinicalTrials.gov'))
    report.clinical_landscape_raw = cls_text + f'状态：{_collect_conditions(trials)}' + (f'\n含 {len(sources)} 项来自 {", ".join(sources)} 的试验。' if sources else '')

    # Key findings
    top = sorted(papers, key=lambda p: p.year or '0', reverse=True)[:5]
    report.key_findings_raw = '\n'.join(f'{i+1}. {p.title} — {p.journal} ({p.year})' for i, p in enumerate(top)) if top else '暂无文献数据。'

    # Future outlook
    if drugs:
        approved = [d for d in drugs if 'appr' in d.phase.lower() or '上市' in d.phase]
        report.future_outlook_raw = f'已有 {len(approved)} 个药物获批，{len(drugs) - len(approved)} 个在研。'
        if approved:
            report.future_outlook_raw += f'已获批药物：{", ".join(d.name for d in approved)}。'
        if molecules:
            report.future_outlook_raw += f'另检索到 {len(molecules)} 个相关活性分子/候选（ChEMBL）线索。'
    else:
        report.future_outlook_raw = '暂未检索到针对该靶点的药物信息。' + (f'\n已检索到 {len(molecules)} 个相关活性分子线索（ChEMBL）。' if molecules else '')

    # Copy raw to main fields
    report.target_overview = report.target_overview_raw
    report.research_progress = report.research_progress_raw
    report.clinical_landscape = report.clinical_landscape_raw
    report.key_findings = report.key_findings_raw
    report.future_outlook = report.future_outlook_raw
    return report


def _build_llm_context(report: ReportContent, field: str) -> str:
    parts = [f'靶点: {report.target_name}']
    d = report.target_detail
    if d:
        if d.gene_symbol: parts.append(f'基因: {d.gene_symbol}')
        if d.protein_class: parts.append(f'蛋白类别: {d.protein_class}')
        if d.protein_name: parts.append(f'蛋白全名: {d.protein_name}')
        if d.description: parts.append(f'功能: {d.description}')
        if d.subcellular: parts.append(f'亚细胞定位: {d.subcellular}')
        if d.tractability: parts.append(f'成药性评估: {", ".join(d.tractability)}')
        if d.pdb: parts.append(f'PDB结构: {"; ".join(p.pdb_id for p in d.pdb[:5])}')
        if d.sequence_length: parts.append(f'氨基酸序列长度: {d.sequence_length} aa')
        if d.synonyms: parts.append(f'别名: {"; ".join(d.synonyms[:10])}')
        if d.hpa: parts.append(f'HPA: {d.hpa.link}')
        if d.related_diseases: parts.append(f'相关疾病: {", ".join(d.related_diseases)}')
    if field in ('research_progress', 'key_findings'):
        parts.append(f'\n文献总数: {len(report.papers)}')
        yc = _count_by_year(report.papers)
        if yc:
            parts.append(f'年份分布: {", ".join(f"{y}({c}篇)" for y, c in yc.items())}')
        parts.append(_clip('\n'.join(f'{i+1}. {p.title} ({p.journal}, {p.year})' for i, p in enumerate(report.papers[:10])), 2200))
    if field in ('clinical_landscape', 'future_outlook'):
        parts.append(f'\n试验总数: {len(report.trials)}')
        ph = _count_by_phase(report.trials)
        if ph:
            parts.append(f'阶段: {json.dumps(ph)}')
        c = _collect_conditions(report.trials)
        if c != '—':
            parts.append(f'适应症: {c}')
    if field == 'future_outlook':
        parts.append(f'\n药物: {len(report.drugs)}')
        for d in report.drugs:
            parts.append(f' - {d.name} [{d.phase}] {d.company} — {d.mechanism_of_action}')
        if report.molecules:
            parts.append(f'活性分子线索(ChEMBL): {", ".join(m.name + (" (" + m.chembl + ")" if m.chembl else "") for m in report.molecules[:8])}')
    return '\n'.join(parts)


def _build_sys_context(report: ReportContent) -> str:
    parts = [f'靶点: {report.target_name}']
    d = report.target_detail
    if d:
        if d.gene_symbol: parts.append(f'基因: {d.gene_symbol}')
        if d.protein_class: parts.append(f'蛋白类别: {d.protein_class}')
        if d.protein_name: parts.append(f'蛋白全名: {d.protein_name}')
        if d.subcellular: parts.append(f'亚细胞定位: {d.subcellular}')
        if d.tractability: parts.append(f'成药性: {", ".join(d.tractability)}')
        if d.pdb: parts.append(f'PDB: {", ".join(p.pdb_id for p in d.pdb[:5])}')
        if d.sequence_length: parts.append(f'序列长度: {d.sequence_length} aa')
        if d.hpa and d.hpa.link: parts.append(f'HPA: {d.hpa.link}')
        if d.related_diseases: parts.append(f'相关疾病: {", ".join(d.related_diseases)}')
        if d.pubchem: parts.append(f'PubChem活性分子: {len(d.pubchem)}条')
        if d.kegg_disease: parts.append(f'KEGG疾病: {", ".join(d.kegg_disease[:5])}')
        if d.kegg_drugs: parts.append(f'KEGG药物: {", ".join(d.kegg_drugs[:5])}')
    parts.append(f'\n[文献-{len(report.papers)}篇]' + _clip('\n'.join(f'{i+1}. {p.title} ({p.journal}, {p.year})' for i, p in enumerate(report.papers[:8])), 1800))
    parts.append(f'\n[临床-{len(report.trials)}项]' + _clip('\n'.join(f'{t.nct_id} {t.phase} {t.status} {t.title}' for t in report.trials[:8]), 1500))
    parts.append(f'\n[药物-{len(report.drugs) + len(report.molecules)}个]' + '\n'.join(f'{d.name} [{d.phase}] {d.company} — {d.mechanism_of_action}' for d in report.drugs[:10]) + (f'\n活性分子: {", ".join(m.name for m in report.molecules[:8])}' if report.molecules else ''))
    parts.append(f'\n[专利-{len(report.patents)}项]' + _clip('\n'.join(f' - {p.year} | {p.assignee or "—"} | {p.number} | {p.title}' for p in report.patents[:15]), 2000))
    return _clip('\n'.join(parts), 5000)


async def enhance_report(report: ReportContent, provider: str, api_key: str, model: str,
                          progress_callback=None, use_cites: bool = True,
                          base_url: str = '') -> ReportContent:
    report.citations = {}
    fields = [
        ('target_overview', 'target_overview'),
        ('research_progress', 'research_progress'),
        ('clinical_landscape', 'clinical_landscape'),
        ('key_findings', 'key_findings'),
        ('future_outlook', 'future_outlook'),
    ]
    for idx, (field, prompt_key) in enumerate(fields):
        if progress_callback:
            progress_callback(72 + idx * 5, f'AI 分析 {field} ...')
        context = _build_llm_context(report, field)
        prompt = LLM_PROMPTS[prompt_key].format(context=context)
        try:
            if use_cites:
                sys_prompt = '你是一名生物医药专家。基于给定数据生成专业分析。请严格返回 JSON：{"text":"分析正文","claims":[{"text":"论断","refs":["PMID:123","NCT04567890"]}]}。refs 只能引用数据中真实出现的编号，没有就留空数组。'
                o = await llm_json(prompt, provider, api_key, model, system=sys_prompt, base_url=base_url)
                if o and o.get('text'):
                    setattr(report, field, o['text'])
                    claims = [c for c in (o.get('claims') or []) if c and c.get('text')]
                    refs = set()
                    for c in claims:
                        for r in (c.get('refs') or []):
                            refs.add(r)
                    combined = o['text'] + ' ' + ' '.join(c.get('text', '') for c in claims)
                    for r in _extract_refs(combined):
                        refs.add(r)
                    if refs:
                        report.citations[field] = list(refs)
                    continue
        except Exception:
            pass
        text = await llm_chat(prompt, provider, api_key, model, base_url=base_url)
        if text:
            setattr(report, field, text)
            refs = _extract_refs(text)
            if refs:
                report.citations[field] = refs

    # Web summary / comprehensive analysis
    p = LLM_PROVIDERS.get(provider, {})
    if p.get('web_search') and provider == 'zhipu':
        if progress_callback:
            progress_callback(95, 'AI 联网情报分析 ...')
        d = report.target_detail
        prompt = LLM_PROMPTS['web_summary'](report.target_name,
                                            d.gene_symbol if d else '',
                                            d.protein_class if d else '')
        text = await llm_chat_with_search(prompt, provider, api_key, model, base_url=base_url)
        if text:
            report.web_summary = text
            refs = _extract_refs(text)
            if refs:
                report.citations['web_summary'] = refs
    else:
        if progress_callback:
            progress_callback(96, 'AI 综合研判 ...')
        ctx = _build_sys_context(report)
        prompt = LLM_PROMPTS['web_summary_legacy'](ctx)
        text = await llm_chat(prompt, provider, api_key, model, base_url=base_url)
        if text:
            report.web_summary = text
            refs = _extract_refs(text)
            if refs:
                report.citations['web_summary'] = refs

    # Patent landscape
    if report.patents:
        if progress_callback:
            progress_callback(97, 'AI 专利格局解读 ...')
        summary_parts = []
        for p in report.patents[:15]:
            s = f' - {p.year} | {p.assignee or "—"} | {p.number} | {p.title}'
            bits = []
            if p.detail_legal: bits.append(f'状态:{p.detail_legal}')
            if p.cited_by: bits.append(f'被引{p.cited_by}')
            if p.claims: bits.append(f'约{p.claims}项权利要求')
            if bits: s += ' [' + ', '.join(bits) + ']'
            if p.abstract: s += '\n    摘要: ' + p.abstract[:160]
            summary_parts.append(s)
        summary = _clip('\n'.join(summary_parts), 4000)
        prompt = LLM_PROMPTS['patent_landscape'](report.target_name, summary)
        text = await llm_chat(prompt, provider, api_key, model, base_url=base_url)
        if text:
            report.patent_landscape = text
            refs = _extract_refs(text)
            if refs:
                report.citations['patent_landscape'] = refs

    return report

# ─── Export: Markdown ──────────────────────────────────────────────────────


def to_markdown(report: ReportContent) -> str:
    name = str(report.target_name or '').upper()
    NL = chr(10)  # newline
    parts = []
    parts.append('# \u9776\u70b9 ' + name + ' \u7814\u7a76\u8fdb\u5c55\u4e0e\u4e34\u5e8a\u5206\u6790\u62a5\u544a' + NL)
    srcs = ['PubMed', 'OpenAlex', 'Semantic Scholar', 'ClinicalTrials.gov', 'ISRCTN', 'ChiCTR',
            'Open Targets', 'UniProt', 'PDB', 'AlphaFold', 'HPA', 'GTEx', 'STRING', 'ChEMBL',
            'PubChem', 'KEGG', 'ClinPGx', 'ClinVar', 'Google Patents', 'Patent9', 'drugfuture',
            'USPTO', 'Lens.org', 'Espacenet']
    parts.append('> \u751f\u6210\uff1a' + datetime.now().strftime("%Y-%m-%d %H:%M") + '  \u6570\u636e\u6e90\uff1a' + ' | '.join(srcs) + NL)
    if report.cached_flags:
        parts.append('> 缓存: ' + '、'.join(report.cached_flags) + NL)
    if report.increment_note:
        parts.append('> ' + report.increment_note + NL)
    d = report.target_detail
    if d:
        det = []
        if d.gene_symbol: det.append('基因符号: ' + d.gene_symbol)
        if d.uniprot_id: det.append('UniProt: ' + d.uniprot_id)
        if d.protein_name: det.append('蛋白名: ' + d.protein_name)
        if d.protein_class: det.append('蛋白类别: ' + d.protein_class)
        if d.tractability: det.append('成药性: ' + ', '.join(d.tractability[:5]))
        if d.related_diseases: det.append('相关疾病: ' + ', '.join(d.related_diseases[:8]))
        if d.pubchem: det.append('PubChem活性分子: ' + str(len(d.pubchem)) + '条')
        if d.kegg_disease: det.append('KEGG疾病: ' + ', '.join(d.kegg_disease[:5]))
        if d.kegg_drugs: det.append('KEGG药物: ' + ', '.join(d.kegg_drugs[:5]))
        if d.subcellular: det.append('亚细胞定位: ' + d.subcellular)
        if d.protein_families: det.append('蛋白家族: ' + d.protein_families)
        if d.clinpgx:
            det.append('ClinPGx药物基因组关联: ' + str(len(d.clinpgx)) + '条')
            for a in d.clinpgx[:12]:
                bits = []
                if a.drug_name: bits.append('药物: ' + a.drug_name)
                if a.level: bits.append('证据级别: ' + a.level)
                if a.disease: bits.append('疾病: ' + a.disease)
                if a.association_type: bits.append('类型: ' + a.association_type)
                if a.evidence: bits.append('证据: ' + a.evidence)
                det.append('  - ' + ' | '.join(bits))
        if report.mutation: det.append('变异: ' + report.mutation)
        if det:
            parts.append(NL + '## 一、靶点详情' + NL + NL.join('- ' + x for x in det) + NL)
    if report.target_overview:
        parts.append(NL + '## 二、靶点概述' + NL + report.target_overview + NL)
    if report.web_summary:
        parts.append(NL + '## AI 综合情报研判' + NL + report.web_summary + NL)
    if report.patent_landscape:
        parts.append(NL + '## 专利调研' + NL + report.patent_landscape + NL)
    parts.append(NL + '## 三、研究进展' + NL + report.research_progress + NL)
    parts.append(NL + '## 四、核心文献' + NL + report.key_findings + NL)
    parts.append(NL + '## 五、临床概况' + NL + report.clinical_landscape + NL)
    parts.append(NL + '## 六、药物展望' + NL + report.future_outlook + NL)
    if report.papers:
        parts.append(NL + '## 文献（' + str(len(report.papers)) + '篇）' + NL)
        for i, p in enumerate(report.papers, 1):
            refs = []
            if p.pmid: refs.append('PMID:' + p.pmid)
            if p.doi: refs.append('DOI: ' + p.doi)
            extra = ' | ' + ' | '.join(refs) if refs else ''
            parts.append(str(i) + '. ' + p.title + ' - ' + p.journal + ' (' + p.year + ')' + extra + NL)
    if report.trials:
        parts.append(NL + '## 临床试验（' + str(len(report.trials)) + '项）' + NL)
        for t in report.trials:
            src = ' [' + t.source + ']' if t.source and t.source != 'ClinicalTrials.gov' else ''
            extra = ''
            if t.conditions: extra += ' 适应症: ' + ', '.join(t.conditions[:4])
            if t.enrollment: extra += ' 入组: ' + str(t.enrollment)
            parts.append('- ' + t.nct_id + ' | ' + t.title + ' | ' + t.phase + ' | ' + t.status + src + extra + NL)
    if report.drugs:
        parts.append(NL + '## 药物（' + str(len(report.drugs)) + '个）' + NL)
        for d in report.drugs:
            extra = ' 适应症: ' + d.disease if d.disease else ''
            parts.append('- ' + d.name + ' [' + d.phase + '] ' + d.company + ' - ' + d.mechanism_of_action + extra + NL)
    if report.molecules:
        parts.append(NL + '## 靶向活性分子（' + str(len(report.molecules)) + '个）' + NL)
        for m in report.molecules:
            vals = []
            if m.pchembl: vals.append('pChEMBL: ' + m.pchembl)
            if m.best_val: vals.append(m.best_type + ': ' + m.best_val)
            extra = ' | ' + ' | '.join(vals) if vals else ''
            parts.append('- ' + m.name + ' (' + m.chembl + ')' + extra + NL)
    if report.patents:
        parts.append(NL + '## 专利（' + str(len(report.patents)) + '项）' + NL)
        if report.patent_note:
            parts.append('> ' + report.patent_note + NL)
        for p in report.patents:
            line = '- ' + p.number + ' | ' + p.title + ' | ' + p.assignee + ' | ' + p.year
            detl = []
            if p.date: detl.append('公告: ' + p.date)
            if p.priority_date: detl.append('优先权: ' + p.priority_date)
            if p.cited_by: detl.append('被引: ' + str(p.cited_by))
            if p.claims: detl.append('权利要求数: ' + str(p.claims))
            if p.detail_legal: detl.append('法律状态: ' + p.detail_legal)
            if p.application_no: detl.append('申请号: ' + p.application_no)
            if p.countries: detl.append('国家: ' + ', '.join(p.countries[:6]))
            if p.inventors: detl.append('发明人: ' + p.inventors)
            if p.cpc: detl.append('CPC: ' + ', '.join(p.cpc[:4]))
            if p.mcp_status: detl.append('MCP状态: ' + p.mcp_status)
            if p.cn_fulltext: detl.append('全文: ' + p.cn_fulltext)
            elif p.cn_download: detl.append('全文(下载页): ' + p.cn_download)
            if detl:
                line += NL + '  ' + ' | '.join(detl)
            parts.append(line + NL)
            if p.abstract and p.abstract != p.snippet:
                parts.append('  > ' + p.abstract[:400] + NL)
    if report.patent_insight:
        parts.append(NL + '## 专利申请人布局' + NL + str(report.patent_insight) + NL)
    td = report.target_detail
    if td and td.sequence:
        parts.append(NL + '## 蛋白序列 (FASTA)' + NL)
        parts.append('>' + (td.gene_symbol or name) + '|' + td.uniprot_id + ' length=' + str(td.sequence_length or len(td.sequence)) + NL)
        seq = td.sequence
        for j in range(0, len(seq), 80):
            parts.append(seq[j:j + 80] + NL)
    if td and td.hpa:
        h = td.hpa
        parts.append(NL + '## 蛋白表达 (HPA)' + NL)
        for lbl, val in [('蛋白组织表达', h.protein_tissue), ('RNA 肿瘤表达', h.rna_cancer),
                         ('蛋白肿瘤表达', h.rna_cancer_score), ('亚细胞定位', h.subcell),
                         ('分子功能', h.molecular_func)]:
            if val:
                parts.append('- ' + lbl + ': ' + val + NL)
    if td and td.pdb:
        parts.append(NL + '## PDB 结构' + NL)
        for s in td.pdb:
            parts.append('- ' + (s.pdb_id or '') + ' | ' + (s.resolution or '') + ' | ' + (s.method or '') + NL)
    if td and td.alphafold:
        parts.append(NL + '## AlphaFold 结构' + NL)
        for s in td.alphafold:
            parts.append('- ' + (s.uniprot_acc or '') + ' | confidence ' + str(s.confidence) + ' | ' + (s.pdb_url or '') + NL)
    if report.citations:
        parts.append(NL + '## 参考文献引用' + NL)
        for field, refs in report.citations.items():
            field_zh = {'target_overview': '靶点概述', 'research_progress': '研究进展',
                        'clinical_landscape': '临床概况', 'key_findings': '核心文献',
                        'future_outlook': '药物展望'}.get(field, field)
            parts.append('### ' + field_zh + NL)
            for r in refs:
                parts.append('- ' + str(r) + NL)
    if report.errors:
        parts.append(NL + '## 数据源失败提示' + NL)
        for e in report.errors:
            parts.append('- ' + e + NL)
    return ''.join(parts)


# ─── Export: JSON ─────────────────────────────────────────────────────────


def to_json(report: ReportContent) -> str:
    return json.dumps(report.model_dump(), ensure_ascii=False, indent=2)


# ─── Export: HTML ─────────────────────────────────────────────────────────


def to_html(report: ReportContent) -> str:
    esc = escape
    md = to_markdown(report)
    body = ''
    for line in md.split(chr(10)):
        if line.startswith('# ') and not line.startswith('## '):
            body += '<h1>' + esc(line[2:]) + '</h1>' + chr(10)
        elif line.startswith('## '):
            body += '<h2>' + esc(line[3:]) + '</h2>' + chr(10)
        elif line.startswith('> '):
            body += '<blockquote>' + esc(line[2:]) + '</blockquote>' + chr(10)
        elif line.startswith('- '):
            body += '<li>' + esc(line[2:]) + '</li>' + chr(10)
        elif line.strip():
            body += '<p>' + esc(line) + '</p>' + chr(10)
    t = esc(report.target_name)
    html = '<!DOCTYPE html>' + chr(10)
    html += '<html lang="zh-CN">' + chr(10)
    html += '<head><meta charset="utf-8"><title>' + t + ' 靶点调研报告</title>' + chr(10)
    html += '<style>' + chr(10)
    html += 'body{font-family:-apple-system,"Microsoft YaHei",sans-serif;max-width:900px;margin:24px auto;padding:0 20px;color:#2C3E50;line-height:1.6}' + chr(10)
    html += 'h1{color:#1B4F72;border-bottom:3px solid #2E86C1;padding-bottom:8px}' + chr(10)
    html += 'h2{color:#2E86C1;border-left:4px solid #2E86C1;padding-left:10px;margin-top:24px}' + chr(10)
    html += 'blockquote{background:#fffbe6;border-left:3px solid #f1c40f;padding:8px 12px;margin:10px 0}' + chr(10)
    html += '</style>' + chr(10)
    html += '</head>' + chr(10)
    html += '<body><button onclick="window.print()">打印/保存PDF</button>' + chr(10)
    html += body + chr(10)
    html += '</body></html>'
    return html
# ─── SQLite History ────────────────────────────────────────────────────────

NL = chr(10)

DB_PATH = str(Path.home() / '.targetinfo.db')


def _get_db() -> sqlite3.Connection:
    db = sqlite3.connect(DB_PATH)
    db.row_factory = sqlite3.Row
    db.execute('CREATE TABLE IF NOT EXISTS reports ('
               'id TEXT PRIMARY KEY, target TEXT, created_at REAL, '
               'report_json TEXT, errs TEXT, counts TEXT, ai_enhanced INTEGER DEFAULT 0)')
    db.execute('CREATE TABLE IF NOT EXISTS cache ('
               'key TEXT PRIMARY KEY, value TEXT, cached_at REAL)')
    db.commit()
    return db


def report_save(report: ReportContent, errs=None) -> str:
    rid = report.target_name + '_' + str(int(time.time() * 1000))
    counts = json.dumps({
        'papers': len(report.papers), 'trials': len(report.trials),
        'drugs': len(report.drugs), 'patents': len(report.patents),
    })
    ai = 1 if report.web_summary or report.patent_landscape else 0
    db = _get_db()
    db.execute('INSERT OR REPLACE INTO reports VALUES (?,?,?,?,?,?,?)',
               (rid, report.target_name, time.time(),
                report.model_dump_json(), json.dumps(errs or []), counts, ai))
    db.commit()
    db.close()
    return rid


def report_list() -> List[dict]:
    db = _get_db()
    rows = db.execute(
        'SELECT id, target, created_at, counts, ai_enhanced '
        'FROM reports ORDER BY created_at DESC').fetchall()
    db.close()
    return [dict(r) for r in rows]


def report_load(rid: str) -> Optional[ReportContent]:
    db = _get_db()
    row = db.execute(
        'SELECT report_json FROM reports WHERE id=?', (rid,)).fetchone()
    db.close()
    if row:
        try:
            return ReportContent(**json.loads(row['report_json']))
        except Exception:
            return None
    return None


def report_delete(rid: str) -> bool:
    db = _get_db()
    db.execute('DELETE FROM reports WHERE id=?', (rid,))
    db.commit()
    db.close()
    return True


def report_clear() -> bool:
    db = _get_db()
    db.execute('DELETE FROM reports')
    db.commit()
    db.close()
    return True


def cache_get(key: str) -> Optional[str]:
    db = _get_db()
    row = db.execute(
        'SELECT value, cached_at FROM cache WHERE key=?', (key,)).fetchone()
    db.close()
    if row and time.time() - row['cached_at'] < 86400:
        return row['value']
    return None


def cache_set(key: str, value: str) -> None:
    db = _get_db()
    db.execute('INSERT OR REPLACE INTO cache VALUES (?,?,?)',
               (key, value, time.time()))
    db.commit()
    db.close()


# ─── Cached search wrapper (24h TTL, mirrors target.html cachedSearch) ─────


_MODEL_MAP = {
    'papers': Paper, 'trials': ClinicalTrial,
    'drugs': DrugInfo, 'molecules': MoleculeInfo,
}


def _year_backs(v) -> int:
    """Map UI year_range string to years_back; 0/'0'/'' => 0 (all years)."""
    try:
        n = int(v or 0)
    except Exception:
        n = 0
    return max(n, 0)


def _trial_key(t) -> str:
    """Dedup key: prefer registry id, fall back to normalized title."""
    if t.nct_id:
        return 'id:' + t.nct_id.lower()
    if t.title:
        return 'title:' + re.sub(r'[^a-z0-9]+', '', t.title.lower())
    return id(t)


def _paper_key(p) -> str:
    if p.pmid:
        return 'pmid:' + p.pmid.lower()
    if p.doi:
        return 'doi:' + p.doi.lower()
    if p.title:
        return 'title:' + re.sub(r'[^a-z0-9]+', '', p.title.lower())
    return id(p)


async def _cached_search(key: str, kind: str, fn) -> Tuple[Any, bool]:
    """Return (data, cached). data is a list of models/dicts, or a dict for patents."""
    try:
        raw = await asyncio.to_thread(cache_get, key)
        if raw is not None:
            arr = json.loads(raw)
            if kind == 'patents' and isinstance(arr, dict):
                return arr, True
            if isinstance(arr, list):
                if kind == 'patents':
                    return [dict(x) for x in arr], True
                model = _MODEL_MAP.get(kind)
                if model and arr and isinstance(arr[0], dict):
                    return [model(**x) for x in arr], True
                if not arr:
                    return [], True
    except Exception:
        pass
    data = await fn()
    try:
        await asyncio.to_thread(cache_set, key, json.dumps(
            data if kind == 'patents' else [x.model_dump() for x in data],
            ensure_ascii=False, default=list))
    except Exception:
        pass
    return data, False


# ─── BM25 + Follow-up Q&A ─────────────────────────────────────────────


def _tokenize(s: str) -> List[str]:
    return [t for t in re.split(
        r'[^a-z0-9\u4e00-\u9fa5]+', (s or '').lower()) if len(t) > 1]


def _bm25_score(query: str, doc_text: str, k: float = 1.5,
                b: float = 0.75) -> float:
    qt = _tokenize(query)
    dt = _tokenize(doc_text)
    if not dt:
        return 0.0
    tf = {}
    for w in dt:
        tf[w] = tf.get(w, 0) + 1
    dl = len(dt)
    score = 0.0
    for w in qt:
        f = tf.get(w, 0)
        if f:
            score += (f * (k + 1)) / (f + k * (1 - b + b * dl))
    return score


def _search_docs(report, query, top_k=8):
    docs = []
    for p in report.papers:
        if p.title:
            docs.append({'id': 'PMID:' + p.pmid,
                         'text': p.title + ' ' + p.abstract,
                         'kind': '\u6587\u732e'})
    for t in report.trials:
        if t.title:
            docs.append({'id': t.nct_id,
                         'text': t.title + ' ' + t.brief_summary,
                         'kind': '\u4e34\u5e8a'})
    for p in report.patents:
        if p.title:
            docs.append({'id': p.number,
                         'text': p.title + ' ' + p.abstract,
                         'kind': '\u4e13\u5229'})
    for d in report.drugs:
        if d.name:
            docs.append({'id': d.name,
                         'text': d.name + ' ' + d.mechanism_of_action,
                         'kind': '\u836f\u7269'})
    scored = [(d, _bm25_score(query, d['text'])) for d in docs if _bm25_score(query, d['text']) > 0]
    scored.sort(key=lambda x: -x[1])
    return [x[0] for x in scored[:top_k]]


def _link_for_ref(ref_id: str) -> str:
    if ref_id.startswith('PMID:'):
        return 'https://pubmed.ncbi.nlm.nih.gov/' + ref_id[5:] + '/'
    if ref_id.startswith('NCT') or ref_id.startswith('ISRCTN'):
        return 'https://clinicaltrials.gov/study/' + ref_id + '/'
    return '#'

# ─── Export: PPT ──────────────────────────────────────────────────────────


DARK_BLUE = PPT_RGBColor(0x1B, 0x4F, 0x72)
MID_BLUE = PPT_RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BLUE = PPT_RGBColor(0xD6, 0xEA, 0xF8)
DARK_GRAY = PPT_RGBColor(0x2C, 0x3E, 0x50)
MED_GRAY = PPT_RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_GRAY = PPT_RGBColor(0xEC, 0xF0, 0xF1)
WHITE_COLOR = PPT_RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT_NAME = 'Microsoft YaHei'

_AVAIL_FONTS = {f.name for f in fontManager.ttflist}
for _f in ['Microsoft YaHei', 'SimHei', 'WenQuanYi Micro Hei',
           'Noto Sans CJK SC', 'WenQuanYi Zen Hei', 'DejaVu Sans']:
    if _f in _AVAIL_FONTS:
        plt.rcParams['font.family'] = _f
        break
else:
    plt.rcParams['font.family'] = 'sans-serif'
plt.rcParams['axes.unicode_minus'] = False

NL2 = chr(10)


def _ppt_set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _ppt_header(slide, title):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.15), Inches(11), Inches(0.8))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR


def _ppt_footer(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.0), SLIDE_W, Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()


def _ppt_add_table(slide, left, top, width, height, headers, rows):
    ts = slide.shapes.add_table(
        len(rows) + 1, len(headers), left, top, width, height)
    table = ts.table
    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.text = h
        for para in c.text_frame.paragraphs:
            para.font.size = Pt(11)
            para.font.bold = True
            para.font.color.rgb = WHITE_COLOR
            para.font.name = FONT_NAME
        c.fill.solid()
        c.fill.fore_color.rgb = MID_BLUE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = table.cell(i + 1, j)
            c.text = str(val)
            for para in c.text_frame.paragraphs:
                para.font.size = Pt(9)
                para.font.color.rgb = DARK_GRAY
                para.font.name = FONT_NAME
            c.fill.solid()
            c.fill.fore_color.rgb = LIGHT_GRAY if i % 2 == 1 else WHITE_COLOR


def generate_ppt(report: ReportContent) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_set_bg(s, DARK_BLUE)
    tb = s.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = '靶点调研报告' + NL2 + str(report.target_name).upper()
    p.font.name = FONT_NAME; p.font.size = Pt(40); p.font.bold = True
    p.font.color.rgb = WHITE_COLOR; p.alignment = PP_ALIGN.CENTER
    tb2 = s.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11), Inches(0.5))
    tb2.text_frame.paragraphs[0].text = '生成日期: ' + datetime.now().strftime('%Y-%m-%d')
    tb2.text_frame.paragraphs[0].font.name = FONT_NAME
    tb2.text_frame.paragraphs[0].font.size = Pt(14)
    tb2.text_frame.paragraphs[0].font.color.rgb = LIGHT_BLUE
    tb2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_header(s, '靶点概述'); _ppt_footer(s)
    if report.target_overview:
        tb = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
        tb.text_frame.word_wrap = True
        tb.text_frame.paragraphs[0].text = report.target_overview[:800]
        tb.text_frame.paragraphs[0].font.name = FONT_NAME
        tb.text_frame.paragraphs[0].font.size = Pt(13)
        tb.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

    if report.target_detail and report.target_detail.sequence:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _ppt_header(s, '蛋白序列 (FASTA)'); _ppt_footer(s)
        tb = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
        tb.text_frame.word_wrap = True
        seq = report.target_detail.sequence
        tb.text_frame.paragraphs[0].text = '>' + (report.target_detail.gene_symbol or str(report.target_name).upper()) + '|' + report.target_detail.uniprot_id + ' length=' + str(report.target_detail.sequence_length or len(seq))
        tb.text_frame.paragraphs[0].font.name = FONT_NAME
        tb.text_frame.paragraphs[0].font.size = Pt(11)
        for j in range(0, len(seq), 80):
            para = tb.text_frame.add_paragraph()
            para.text = seq[j:j + 80]
            para.font.name = FONT_NAME; para.font.size = Pt(11)

    td = report.target_detail
    if td and td.clinpgx:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _ppt_header(s, 'ClinPGx 药物基因组关联'); _ppt_footer(s)
        hd = ['药物', '证据级别', '疾病', '类型']
        r = [[a.drug_name, a.level, a.disease, a.association_type] for a in td.clinpgx[:10]]
        _ppt_add_table(s, Inches(0.3), Inches(1.4), Inches(12.5), Inches(5.0), hd, r)
        if td.clinpgx[0].evidence:
            tb = s.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9))
            tb.text_frame.word_wrap = True
            tb.text_frame.paragraphs[0].text = '示例证据: ' + td.clinpgx[0].evidence[:180]
            tb.text_frame.paragraphs[0].font.name = FONT_NAME
            tb.text_frame.paragraphs[0].font.size = Pt(10)
            tb.text_frame.paragraphs[0].font.color.rgb = MED_GRAY

    if report.research_progress or report.papers:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _ppt_header(s, '研究进展')
        if report.research_progress:
            tb = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.5))
            tb.text_frame.paragraphs[0].text = report.research_progress[:300]
            tb.text_frame.paragraphs[0].font.name = FONT_NAME
            tb.text_frame.paragraphs[0].font.size = Pt(13)
        if report.papers:
            hd = ['#', '标题', '期刊', '年份']
            r = [[str(i+1), (p.title or '')[:60], (p.journal or '')[:25], p.year or '']
                 for i, p in enumerate(report.papers[:10])]
            _ppt_add_table(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(3.5), hd, r)
        _ppt_footer(s)

    if report.key_findings:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _ppt_header(s, '核心文献发现'); _ppt_footer(s)
        tb = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
        tb.text_frame.word_wrap = True
        tb.text_frame.paragraphs[0].text = report.key_findings[:900]
        tb.text_frame.paragraphs[0].font.name = FONT_NAME
        tb.text_frame.paragraphs[0].font.size = Pt(13)
        tb.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY

    if report.trials:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _ppt_header(s, '临床试验')
        y_off = Inches(1.4)
        if report.clinical_landscape:
            tb = s.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(11.9), Inches(1.3))
            tb.text_frame.word_wrap = True
            tb.text_frame.paragraphs[0].text = report.clinical_landscape[:260]
            tb.text_frame.paragraphs[0].font.name = FONT_NAME
            tb.text_frame.paragraphs[0].font.size = Pt(11)
            y_off = Inches(2.7)
        hd = ['NCT', '标题', '阶段', '状态']
        r = [[t.nct_id or '', (t.title or '')[:50], t.phase or '', t.status or '']
             for t in report.trials[:12]]
        _ppt_add_table(s, Inches(0.5), y_off, Inches(12.3), Inches(5.0), hd, r)
        _ppt_footer(s)

    if report.drugs:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _ppt_header(s, '药物管线')
        hd = ['药物', '公司', '阶段', '机制']
        r = [[d.name or '', (d.company or '')[:20], d.phase or '',
              (d.mechanism_of_action or '')[:35]] for d in report.drugs[:12]]
        _ppt_add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.0), hd, r)
        _ppt_footer(s)

    if report.molecules:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _ppt_header(s, '靶向活性分子')
        hd = ['分子', 'ChEMBL ID', 'pChEMBL', '活性']
        r = [[m.name or '', m.chembl or '', m.pchembl or '',
              (m.best_type or '') + ' ' + (m.best_val or '')] for m in report.molecules[:12]]
        _ppt_add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.0), hd, r)
        _ppt_footer(s)

    if report.web_summary:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _ppt_header(s, 'AI 综合情报研判')
        tb = s.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
        tb.text_frame.paragraphs[0].text = report.web_summary[:1200]
        tb.text_frame.paragraphs[0].font.name = FONT_NAME
        tb.text_frame.paragraphs[0].font.size = Pt(13)
        tb.text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
        _ppt_footer(s)

    if report.patents:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _ppt_header(s, '专利调研')
        if report.patent_landscape:
            tb = s.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.9), Inches(2.0))
            tb.text_frame.paragraphs[0].text = report.patent_landscape[:400]
            tb.text_frame.paragraphs[0].font.name = FONT_NAME
            tb.text_frame.paragraphs[0].font.size = Pt(11)
            y_off = Inches(3.5)
        else:
            y_off = Inches(1.4)
        hd = ['专利号', '标题', '申请人', '年份']
        r = [[(p.number or '') + (' (引' + str(p.cited_by) + ')' if p.cited_by else ''),
              (p.title or '')[:45], (p.assignee or '')[:20], p.year or '']
             for p in report.patents[:10]]
        _ppt_add_table(s, Inches(0.5), y_off, Inches(12.3), Inches(3.5), hd, r)
        _ppt_footer(s)

    if report.patents:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _ppt_header(s, '专利详情')
        hd = ['专利号', '申请号', '法律状态', '被引', '权利要求', 'CPC']
        r = []
        for p in report.patents[:8]:
            r.append([p.number or '', (p.application_no or '')[:20], (p.detail_legal or '')[:15],
                      str(p.cited_by or ''), str(p.claims or ''),
                      ', '.join(p.cpc[:3])[:24]])
        if r:
            _ppt_add_table(s, Inches(0.3), Inches(1.4), Inches(12.5), Inches(5.0), hd, r)
        if report.patent_insight:
            tb = s.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9))
            tb.text_frame.word_wrap = True
            tb.text_frame.paragraphs[0].text = '申请人布局: ' + str(report.patent_insight)[:200]
            tb.text_frame.paragraphs[0].font.name = FONT_NAME
            tb.text_frame.paragraphs[0].font.size = Pt(10)
            tb.text_frame.paragraphs[0].font.color.rgb = MED_GRAY
        _ppt_footer(s)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_header(s, '数据来源')
    sources = [
        '文献: PubMed | OpenAlex | Semantic Scholar',
        '临床试验: ClinicalTrials.gov | ISRCTN | ChiCTR',
        '靶点: Open Targets | UniProt | ClinVar | ClinPGx',
        '蛋白: PDB | AlphaFold | STRING',
        '表达: Human Protein Atlas | GTEx',
        '药物: Open Targets | ChEMBL | PubChem | KEGG',
        '专利: Google Patents | Patent9 | drugfuture | USPTO | Lens.org | Espacenet | MCP',
        'AI: DeepSeek | 小米MiMo | 智谱GLM | 自定义',
    ]
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.0))
    for i, src in enumerate(sources):
        para = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
        para.text = '\u25b8 ' + src
        para.font.name = FONT_NAME; para.font.size = Pt(14)
        para.font.color.rgb = DARK_BLUE; para.space_after = Pt(8)
    _ppt_footer(s)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ─── Export: Word ──────────────────────────────────────────────────────────


def generate_docx(report: ReportContent) -> bytes:
    doc = Document()
    style = doc.styles['Normal']
    style.font.name = 'Microsoft YaHei'
    style.font.size = Pt(10.5)
    style.paragraph_format.space_after = Pt(4)

    def add_h(text, level=1):
        h = doc.add_heading(text, level=level)
        for run in h.runs: run.font.name = 'Microsoft YaHei'

    def add_p(text, bold=False, size=None):
        p = doc.add_paragraph()
        r = p.add_run(text); r.bold = bold
        r.font.name = 'Microsoft YaHei'
        if size: r.font.size = Pt(size)

    def add_tbl(headers, rows):
        t = doc.add_table(rows=1 + len(rows), cols=len(headers))
        t.style = 'Table Grid'; t.alignment = WD_TABLE_ALIGNMENT.CENTER
        for j, h in enumerate(headers): t.rows[0].cells[j].text = h
        for i, row in enumerate(rows):
            for j, val in enumerate(row): t.rows[i+1].cells[j].text = str(val)

    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run('靶点 ' + str(report.target_name).upper() + ' 研究进展与临床分析报告')
    r.bold = True; r.font.size = Pt(22); r.font.name = 'Microsoft YaHei'
    add_p('生成日期: ' + datetime.now().strftime('%Y-%m-%d %H:%M'), size=11)
    doc.add_page_break()

    add_h('一、靶点概述', 1); d = report.target_detail
    if d:
        add_p('靶点名称: ' + report.target_name)
        if d.gene_symbol: add_p('基因符号: ' + d.gene_symbol)
        if d.protein_class: add_p('蛋白类别: ' + d.protein_class)
        if d.protein_name: add_p('蛋白全名: ' + d.protein_name)
        if d.synonyms: add_p('别名: ' + '; '.join(d.synonyms[:8]))
        if d.description: add_p('功能描述: ' + d.description)
        if d.subcellular: add_p('亚细胞定位: ' + d.subcellular)
        if d.related_diseases: add_p('相关疾病: ' + ', '.join(d.related_diseases))
        if d.pubchem: add_p('PubChem活性分子: ' + str(len(d.pubchem)) + '条')
        if d.kegg_disease: add_p('KEGG疾病: ' + '; '.join(d.kegg_disease[:6]))
        if d.kegg_drugs: add_p('KEGG药物: ' + '; '.join(d.kegg_drugs[:6]))
        if d.clinpgx:
            add_tbl(['药物', '证据级别', '疾病', '类型', '证据'],
                    [[a.drug_name, a.level, a.disease, a.association_type, a.evidence]
                     for a in d.clinpgx[:12]])
    if report.target_overview: add_p(report.target_overview)
    doc.add_page_break()

    add_h('二、研究进展', 1)
    add_p(report.research_progress or ('共 ' + str(len(report.papers)) + ' 篇文献'))
    if report.papers:
        add_h('文献列表', 2)
        add_tbl(['#', '标题', '期刊', '年份', 'PMID'],
                [[str(i+1), p.title, p.journal, p.year, p.pmid] for i, p in enumerate(report.papers)])
    doc.add_page_break()

    add_h('三、临床试验', 1)
    add_p(report.clinical_landscape or ('共 ' + str(len(report.trials)) + ' 项试验'))
    if report.trials:
        add_tbl(['NCT ID', '标题', '阶段', '状态'],
                [[t.nct_id, t.title, t.phase, t.status] for t in report.trials])
    doc.add_page_break()

    add_h('四、药物研发管线', 1)
    add_p(report.future_outlook or ('共 ' + str(len(report.drugs)) + ' 个药物'))
    if report.drugs:
        add_tbl(['药物', '公司', '机制', '阶段'],
                [[d.name, d.company, d.mechanism_of_action, d.phase] for d in report.drugs])
    doc.add_page_break()

    if report.patents:
        add_h('五、专利调研', 1)
        if report.patent_landscape: add_p(report.patent_landscape)
        add_tbl(['专利号', '申请号', '标题', '申请人', '年份', '全文'],
                [[p.number, p.application_no, p.title[:60], p.assignee[:25], p.year,
                  p.cn_fulltext or p.cn_download or ''] for p in report.patents])
        doc.add_page_break()

    if report.web_summary:
        add_h('六、AI 综合情报研判', 1)
        for line in report.web_summary.split(NL2):
            if line.strip(): add_p(line)
        doc.add_page_break()

    add_h('七、附录', 1)
    if report.mutation:
        add_p('靶点变异: ' + report.mutation)
    if d and d.sequence:
        add_h('蛋白序列 (FASTA)', 2)
        add_p('>' + (d.gene_symbol or str(report.target_name).upper()) + '|' + d.uniprot_id +
              ' length=' + str(d.sequence_length or len(d.sequence)), size=9)
        seq = d.sequence
        for j in range(0, len(seq), 100):
            add_p(seq[j:j + 100], size=9)
    if d and d.hpa:
        add_h('蛋白表达 (HPA)', 2)
        h = d.hpa
        for lbl, val in [('蛋白组织表达', h.protein_tissue), ('RNA 肿瘤表达', h.rna_cancer),
                         ('蛋白肿瘤表达', h.rna_cancer_score), ('亚细胞定位', h.subcell),
                         ('分子功能', h.molecular_func)]:
            if val:
                add_p(lbl + ': ' + val, size=10)
    if d and d.string_interactions:
        add_h('蛋白互作 (STRING)', 2)
        add_tbl(['蛋白', '名称', '评分'],
                [[s.protein_id, s.preferred_name, str(s.score)] for s in d.string_interactions[:15]])
    if d and d.kegg_pathways:
        add_h('KEGG 通路', 2)
        add_tbl(['ID', '通路'], [[k.kegg_id, k.name] for k in d.kegg_pathways[:15]])
    if d and d.kegg_disease:
        add_h('KEGG 疾病', 2)
        for x in d.kegg_disease[:10]:
            add_p('\u25b8 ' + x, size=10)
    if d and d.kegg_drugs:
        add_h('KEGG 药物', 2)
        for x in d.kegg_drugs[:10]:
            add_p('\u25b8 ' + x, size=10)
    if d and d.pubchem:
        add_h('PubChem 活性分子', 2)
        add_tbl(['CID', '名称', '分子式', 'MW', 'LogP'],
                [[str(c.cid), c.name, c.formula, str(c.molecular_weight), str(c.logp)] for c in d.pubchem[:12]])
    if d and d.clinvar:
        add_h('ClinVar 临床变异', 2)
        add_tbl(['变异/RCV', '临床意义', '疾病'],
                [[c.rcv_id or c.gene_symbol, c.clinical_significance, c.condition] for c in d.clinvar[:12]])
    if d and d.pdb:
        add_h('PDB 结构', 2)
        add_tbl(['PDB', '方法', '分辨率', '年份'],
                [[s.pdb_id, s.method, s.resolution, s.year] for s in d.pdb[:12]])
    if d and d.alphafold:
        add_h('AlphaFold 结构', 2)
        add_tbl(['UniProt', '置信度', 'URL'],
                [[s.uniprot_acc, str(s.confidence), s.pdb_url] for s in d.alphafold[:5]])
    if report.patent_insight:
        add_h('专利申请人布局', 2)
        add_p(str(report.patent_insight), size=10)
    doc.add_page_break()

    add_h('八、参考文献引用', 1)
    if report.citations:
        for field, refs in report.citations.items():
            if refs:
                add_p(field + ': ' + ', '.join(str(x) for x in refs), size=10)
    else:
        add_p('无', size=10)
    doc.add_page_break()

    add_h('九、数据来源', 1)
    for src in ['PubMed', 'OpenAlex', 'Semantic Scholar', 'ClinicalTrials.gov',
                'ISRCTN', 'ChiCTR',
                'Open Targets', 'UniProt',
                'PDB', 'AlphaFold', 'Human Protein Atlas', 'GTEx',
                'ChEMBL', 'PubChem', 'STRING', 'KEGG',
                'ClinPGx', 'ClinVar', 'Google Patents', 'Patent9', 'drugfuture',
                'USPTO', 'Lens.org', 'Espacenet', 'MCP']:
        add_p('\u25b8 ' + src)
    buf = io.BytesIO(); doc.save(buf); buf.seek(0)
    return buf.getvalue()


# ─── PyQt5 UI ─────────────────────────────────────────────────────────────


NL = chr(10)


class PipelineWorker(QThread):
    progress = pyqtSignal(int, str)
    finished = pyqtSignal(object)
    error = pyqtSignal(str)

    def __init__(self, target_name, gene, mutation, provider, api_key, model, base_url,
                 use_llm=True, src_papers=True, src_trials=True, src_target=True,
                 src_drugs=True, src_patents=True, patent_src='google', patent_key='',
                 mcp_url='', mcp_tool='', pat_deep=True, year_range='20',
                 ncbi_key='', intl_trials=False, ai_cites=True, save_hist=True):
        super().__init__()
        self.target_name = target_name
        self.gene = gene; self.mutation = mutation
        self.provider = provider; self.api_key = api_key
        self.model = model; self.base_url = base_url
        self.use_llm = use_llm
        self.src_papers = src_papers; self.src_trials = src_trials
        self.src_target = src_target; self.src_drugs = src_drugs
        self.src_patents = src_patents
        self.patent_src = patent_src; self.patent_key = patent_key
        self.mcp_url = mcp_url; self.mcp_tool = mcp_tool
        self.pat_deep = pat_deep; self.year_range = year_range
        self.ncbi_key = ncbi_key; self.intl_trials = intl_trials
        self.ai_cites = ai_cites; self.save_hist = save_hist
        self._stop = False

    def request_stop(self):
        self._stop = True

    def _check_stop(self):
        if self._stop:
            raise RuntimeError('已取消')

    def is_cancelled(self) -> bool:
        return self._stop

    def run(self):
        try:
            asyncio.run(self._async_run())
        except RuntimeError as e:
            if str(e) == '已取消':
                self.finished.emit(None)
                return
            import traceback
            self.error.emit(str(e) + NL + NL + traceback.format_exc())
        except Exception as e:
            import traceback
            self.error.emit(str(e) + NL + NL + traceback.format_exc())

    async def _async_run(self):
        p = self.progress
        gene_sym = self.gene or self.target_name
        errs = []
        cached_flags = []

        # Set NCBI key globally
        if self.ncbi_key:
            global NCBI_API_KEY
            NCBI_API_KEY = self.ncbi_key

        papers = []; trials = []; detail = TargetDetail(); drugs = []
        molecules = []; patents = []; patent_insight = None; patent_note = ''
        self._check_stop()

        if self.src_papers:
            p.emit(5, '正在搜索文献 (PubMed)...')
            try:
                res, cached = await _limiter.run(_cached_search(
                    f'papers:{self.target_name}:{self.year_range}', 'papers',
                    lambda: search_papers(gene_sym, 40, _year_backs(self.year_range))))
                papers = res
                if cached: cached_flags.append('文献(缓存)')
                p.emit(12, '文献: PubMed ' + str(len(papers)) + '篇')
                self._check_stop()
                if self.mutation:
                    res2, _ = await _limiter.run(_cached_search(
                        f'papers:{self.target_name}:mut', 'papers',
                        lambda: search_papers(self.target_name, 20, _year_backs(self.year_range))))
                    seen = {_paper_key(p) for p in papers}
                    papers.extend(p for p in res2 if _paper_key(p) not in seen)
                oa = await _limiter.run(search_openalex(self.target_name, 10))
                ss = await _limiter.run(search_semantic_scholar(self.target_name, 10))
                papers.extend(oa); papers.extend(ss)
            except Exception as e:
                if not self.is_cancelled():
                    errs.append('文献: ' + str(e)[:150])

        self._check_stop()
        if self.src_trials:
            p.emit(18, '正在查询临床试验...')
            try:
                res, cached = await _limiter.run(_cached_search(
                    f'trials:{gene_sym}', 'trials',
                    lambda: search_trials(gene_sym, 20)))
                trials = res
                if cached: cached_flags.append('临床(缓存)')
                if self.mutation:
                    res2, _ = await _cached_search(f'trials:{self.target_name}:mut', 'trials',
                                                   lambda: search_trials(self.target_name, 10))
                    seen = {_trial_key(t) for t in trials}
                    trials.extend(t for t in res2 if _trial_key(t) not in seen)
                p.emit(25, '临床: ' + str(len(trials)) + '项')
                if self.intl_trials:
                    for fn in [search_isrctn, search_chictr]:
                        try:
                            extra = await _limiter.run(fn(self.target_name, 5))
                            seen = {_trial_key(t) for t in trials}
                            trials.extend(t for t in extra if _trial_key(t) not in seen)
                        except Exception:
                            pass
            except Exception as e:
                if not self.is_cancelled():
                    errs.append('临床: ' + str(e)[:150])

        self._check_stop()
        if self.src_target:
            p.emit(30, '正在获取靶点信息 (Open Targets, UniProt, HPA, PDB, AlphaFold, STRING, KEGG, PubChem, ClinVar, ClinPGx)...')
            try:
                detail = await _limiter.run(get_target_detail(self.target_name))
            except Exception as e:
                if not self.is_cancelled():
                    errs.append('靶点信息: ' + str(e)[:150])
            try:
                uni = await _limiter.run(get_uniprot(
                    detail.gene_symbol or gene_sym, detail.synonyms))
                if uni:
                    detail.uniprot_id = uni.get('accession', '')
                    detail.protein_name = uni.get('proteinName', '')
                    detail.subcellular = uni.get('subcellular', '')
                    detail.protein_families = uni.get('proteinFamilies', '')
                    detail.sequence = uni.get('sequence', '')
                    detail.sequence_length = uni.get('sequenceLength', 0)
                    acc = uni.get('accession', '')
                    detail.pdb = await _limiter.run(get_pdb(gene_sym, acc))
                    detail.alphafold = await _limiter.run(get_alphafold(gene_sym, acc))
                else:
                    detail.pdb = await _limiter.run(get_pdb(gene_sym, ''))
                    detail.alphafold = await _limiter.run(get_alphafold(gene_sym, ''))
                hpa = await _limiter.run(get_hpa(detail.gene_symbol or gene_sym))
                if hpa: detail.hpa = hpa
                detail.gtex = await _limiter.run(get_gtex(detail.gene_symbol or gene_sym))
                detail.string_interactions = await _limiter.run(
                    get_string_interactions(detail.gene_symbol or gene_sym))
                detail.kegg_pathways = await _limiter.run(
                    get_kegg_pathways(detail.gene_symbol or gene_sym))
                detail.pubchem = await _limiter.run(
                    get_pubchem(detail.gene_symbol or gene_sym))
                detail.kegg_disease = await _limiter.run(
                    get_kegg_disease(detail.gene_symbol or gene_sym))
                detail.kegg_drugs = await _limiter.run(
                    get_kegg_drugs(detail.gene_symbol or gene_sym))
                detail.clinvar = await _limiter.run(
                    get_clinvar(detail.gene_symbol or gene_sym))
                detail.clinpgx = await _limiter.run(
                    get_clinpgx(detail.gene_symbol or gene_sym))
                detail.dgidb = await _limiter.run(
                    get_dgidb(detail.gene_symbol or gene_sym))
                detail.mutation = self.mutation
            except Exception as e:
                if not self.is_cancelled():
                    errs.append('靶点详情: ' + str(e)[:150])

        self._check_stop()
        if self.src_drugs:
            p.emit(50, '正在获取药物信息 (Open Targets, ChEMBL)...')
            try:
                res, cached = await _limiter.run(_cached_search(
                    f'drugs:{gene_sym}', 'drugs', lambda: get_drugs(self.target_name)))
                drugs = res
                if cached: cached_flags.append('药物(缓存)')
                if not cached:
                    await _enrich_drugs_with_company(drugs)
                res_m, cached_m = await _limiter.run(_cached_search(
                    f'chembl:{gene_sym}', 'molecules', lambda: get_chembl(gene_sym)))
                molecules = res_m
                if cached_m: cached_flags.append('分子(缓存)')
            except Exception as e:
                if not self.is_cancelled():
                    errs.append('药物: ' + str(e)[:150])

        self._check_stop()
        if self.src_patents:
            p.emit(60, '正在检索专利...')
            try:
                async def _patents_fetch():
                    return await search_patents(
                        gene_sym, self.patent_src, self.patent_key,
                        deep=self.pat_deep, mcp_url=self.mcp_url, mcp_tool=self.mcp_tool,
                        progress_cb=lambda t: p.emit(66, t))
                res, cached = await _limiter.run(_cached_search(
                    f'patents:{gene_sym}:{self.patent_src}', 'patents', _patents_fetch))
                if res and isinstance(res, dict):
                    patents = res.get('patents') or []
                    patent_insight = res.get('insight')
                    patent_note = res.get('note') or ''
                else:
                    patents = res or []
                if cached and patents: cached_flags.append('专利(缓存)')
            except Exception as e:
                if not self.is_cancelled():
                    errs.append('专利: ' + str(e)[:150])

        self._check_stop()
        p.emit(68, '正在构建报告...')
        report = build_report(self.target_name, self.gene, self.mutation,
                              detail, papers, trials, drugs, patents, molecules)
        report.patent_insight = patent_insight
        report.patent_note = patent_note
        report.cached_flags = cached_flags

        # Increment note vs previous report (mirrors target.html)
        try:
            key = f'{(gene_sym or self.target_name).strip().lower()}'
            prev = None
            for r in report_list():
                if str(r.get('target', '')).strip().lower() == key and r.get('id'):
                    prev = r
                    break
            if prev:
                try:
                    prevc = json.loads(prev.get('counts') or '{}')
                except Exception:
                    prevc = {}
                d_p = len(papers) - int(prevc.get('papers', 0))
                d_t = len(trials) - int(prevc.get('trials', 0))
                d_pa = len(patents) - int(prevc.get('patents', 0))
                parts = []
                if d_p: parts.append(f'文献 {d_p:+.0f}')
                if d_t: parts.append(f'临床 {d_t:+.0f}')
                if d_pa: parts.append(f'专利 {d_pa:+.0f}')
                report.increment_note = '与上次生成相比: ' + '、'.join(parts) + '。' if parts else ''
        except Exception:
            pass

        if self.use_llm and self.api_key and self.model:
            p.emit(72, 'AI 增强报告中...')
            report = await enhance_report(
                report, self.provider, self.api_key, self.model,
                progress_callback=p.emit, use_cites=self.ai_cites, base_url=self.base_url)

        self._check_stop()
        p.emit(95, '正在生成文件...')
        p.emit(100, '完成!')
        report.errors = errs
        self.finished.emit(report)


class FollowUpWorker(QThread):
    result = pyqtSignal(str, list)

    def __init__(self, report, query, provider, api_key, model, base_url, parent=None):
        super().__init__(parent)
        self._report = report
        self._query = query
        self._provider = provider
        self._api_key = api_key
        self._model = model
        self._base_url = base_url

    def run(self):
        docs = _search_docs(self._report, self._query, top_k=8)
        if not docs:
            self.result.emit('', [])
            return
        try:
            ctx = '\n\n'.join('[%d] (%s %s) %s' % (i + 1, d['kind'], d['id'], (d['text'] or '')[:400])
                              for i, d in enumerate(docs))
            prompt = ('基于以下检索片段回答用户问题。请只用片段中的信息，并标注引用编号 [1][2]…。\n\n'
                      + ctx + '\n\n问题：' + self._query)
            text = asyncio.run(llm_chat(prompt, self._provider, self._api_key, self._model,
                                        base_url=self._base_url))
        except Exception:
            text = ''
        self.result.emit(text, docs)


class HistoryWidget(QWidget):
    selected = pyqtSignal(str)
    compare = pyqtSignal(str, str)

    def __init__(self):
        super().__init__()
        self.setup_ui()

    def setup_ui(self):
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        hdr = QHBoxLayout()
        title = QLabel('历史报告')
        title.setStyleSheet('font-weight:600;color:#2E86C1;font-size:14px')
        refresh_btn = QPushButton('刷新')
        refresh_btn.clicked.connect(self.refresh)
        clear_btn = QPushButton('清空')
        clear_btn.clicked.connect(self.clear_all)
        compare_btn = QPushButton('对比所选')
        compare_btn.setToolTip('先分别点击两行记录（第一行作基准），再点击"对比所选"')
        compare_btn.clicked.connect(self._on_compare)
        hdr.addWidget(title); hdr.addStretch()
        hdr.addWidget(compare_btn); hdr.addWidget(refresh_btn); hdr.addWidget(clear_btn)
        layout.addLayout(hdr)
        self.table = QTableWidget()
        self.table.setColumnCount(7)
        self.table.setHorizontalHeaderLabels(['靶点', '时间', '文献', '临床', '专利', 'AI', '操作'])
        self.table.horizontalHeader().setStretchLastSection(True)
        self.table.setSelectionBehavior(QAbstractItemView.SelectRows)
        self.table.setEditTriggers(QAbstractItemView.NoEditTriggers)
        self.table.cellDoubleClicked.connect(self._on_click)
        layout.addWidget(self.table)

    def refresh(self):
        items = report_list()
        revised = []
        for item in items:
            try:
                rd = report_load(item['id'])
                if rd:
                    item['_revised'] = rd
            except Exception:
                pass
            revised.append(item)
        self.table.setRowCount(len(revised))
        for i, item in enumerate(revised):
            counts = json.loads(item.get('counts', '{}'))
            created = datetime.fromtimestamp(item.get('created_at', 0))
            target_text = str(item.get('target', ''))
            ai_ext = 'AI' if item.get('ai_enhanced') else ''
            rd = item.get('_revised')
            pat = str(rd.patents.__len__()) if rd else str(counts.get('patents', 0))
            cell = QTableWidgetItem(target_text)
            cell.setData(Qt.UserRole, item['id'])
            self.table.setItem(i, 0, cell)
            self.table.setItem(i, 1, QTableWidgetItem(created.strftime('%Y-%m-%d %H:%M')))
            self.table.setItem(i, 2, QTableWidgetItem(str(counts.get('papers', 0))))
            self.table.setItem(i, 3, QTableWidgetItem(str(counts.get('trials', 0))))
            self.table.setItem(i, 4, QTableWidgetItem(pat))
            self.table.setItem(i, 5, QTableWidgetItem(ai_ext))
            w = QWidget(); wl = QHBoxLayout(w); wl.setContentsMargins(2, 2, 2, 2)
            exp = QPushButton('导出'); exp.setEnabled(rd is not None)
            exp.clicked.connect(lambda checked, rid=item['id']: self._export(rid))
            dl = QPushButton('删除'); dl.setStyleSheet('color:#C0392B')
            dl.clicked.connect(lambda checked, rid=item['id']: self._delete(rid))
            wl.addWidget(exp); wl.addWidget(dl); wl.addStretch()
            self.table.setCellWidget(i, 6, w)

    def _on_click(self, row, col):
        cell = self.table.item(row, 0)
        if cell and col != 6:
            rid = cell.data(Qt.UserRole) or cell.text()
            self.selected.emit(str(rid))

    def _on_compare(self):
        sel = self.table.selectedItems()
        rows = sorted({it.row() for it in sel})
        if len(rows) < 2:
            QMessageBox.warning(self, '提示', '请先点击选择两行记录（第一行作基准）')
            return
        cells = [self.table.item(r, 0) for r in rows]
        rids = [str(c.data(Qt.UserRole) or c.text()) for c in cells]
        self.compare.emit(rids[0], rids[1])

    def _export(self, rid):
        rd = report_load(rid)
        if not rd:
            QMessageBox.warning(self, '提示', '无法加载该历史记录')
            return
        path, _ = QFileDialog.getSaveFileName(
            self, '导出历史', slug(rd.target_name) + '_报告.md', 'Markdown (*.md)')
        if not path:
            return
        try:
            with open(path, 'w', encoding='utf-8') as f:
                f.write(to_markdown(rd))
            QMessageBox.information(self, '导出成功', 'Markdown 已导出: ' + path)
        except Exception as e:
            QMessageBox.critical(self, '导出失败', str(e))

    def _delete(self, rid):
        if QMessageBox.question(self, '确认', '确定删除该历史记录?') == QMessageBox.Yes:
            report_delete(rid); self.refresh()

    def clear_all(self):
        if QMessageBox.question(self, '确认', '确定清空全部历史?') == QMessageBox.Yes:
            report_clear(); self.refresh()


class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle('TargetInfo - 靶点调研报告生成器')
        self.setMinimumSize(950, 750)
        self.resize(1150, 880)
        self._report = None
        self.setup_ui()
        self._restore_settings()

    def _settings_sources(self):
        return [(self.chk_papers, 'srcPapers'), (self.chk_trials, 'srcTrials'),
                (self.chk_target, 'srcTarget'), (self.chk_drugs, 'srcDrugs'),
                (self.chk_patents, 'srcPatents')]

    def _restore_settings(self):
        s = QSettings('TargetInfo', 'TargetInfo')
        prov = s.value('providerSel', 'deepseek')
        if prov in ['deepseek', 'mimo', 'zhipu', 'custom']:
            self.provider_combo.setCurrentText(prov)
        self.model_combo.setCurrentText(s.value('model-' + self.provider_combo.currentText(), ''))
        self.api_key_input.setText(s.value('apiKey', ''))
        self.base_url_input.setText(s.value('customBase', ''))
        self.pat_src_combo.setCurrentText(s.value('patSrc', 'google'))
        self.pat_key_input.setText(s.value('patKey', ''))
        self.mcp_url_input.setText(s.value('patMcpUrl', ''))
        self.mcp_tool_input.setText(s.value('patMcpTool', ''))
        self.chk_pat_deep.setChecked(str(s.value('patDeep', 'true')).lower() in ('true', '1'))
        self.chk_intl.setChecked(str(s.value('intlTrials', 'false')).lower() in ('true', '1'))
        self.chk_ai_cites.setChecked(str(s.value('aiCites', 'true')).lower() in ('true', '1'))
        self.chk_save_hist.setChecked(str(s.value('saveHist', 'true')).lower() in ('true', '1'))
        self.chk_use_llm.setChecked(str(s.value('useLLM', 'true')).lower() in ('true', '1'))
        self.year_combo.setCurrentText(s.value('yearRange', '20'))
        self.ncbi_key_input.setText(s.value('ncbiKey', ''))
        for cb, key in self._settings_sources():
            v = s.value(key, None)
            if v is not None:
                cb.setChecked(str(v).lower() in ('true', '1'))
        self.provider_combo.currentIndexChanged.connect(self._settings_changed)
        self.model_combo.currentTextChanged.connect(self._settings_changed)
        self.api_key_input.textChanged.connect(self._settings_changed)
        self.base_url_input.textChanged.connect(self._settings_changed)
        self.pat_src_combo.currentIndexChanged.connect(self._settings_changed)
        self.pat_key_input.textChanged.connect(self._settings_changed)
        self.mcp_url_input.textChanged.connect(self._settings_changed)
        self.mcp_tool_input.textChanged.connect(self._settings_changed)
        self.chk_pat_deep.toggled.connect(self._settings_changed)
        self.chk_intl.toggled.connect(self._settings_changed)
        self.chk_ai_cites.toggled.connect(self._settings_changed)
        self.chk_save_hist.toggled.connect(self._settings_changed)
        self.chk_use_llm.toggled.connect(self._settings_changed)
        self.year_combo.currentIndexChanged.connect(self._settings_changed)
        self.ncbi_key_input.textChanged.connect(self._settings_changed)
        for cb, _ in self._settings_sources():
            cb.toggled.connect(self._settings_changed)

    def _settings_changed(self, *args):
        s = QSettings('TargetInfo', 'TargetInfo')
        s.setValue('providerSel', self.provider_combo.currentText())
        s.setValue('model-' + self.provider_combo.currentText(), self.model_combo.currentText())
        s.setValue('apiKey', self.api_key_input.text())
        s.setValue('customBase', self.base_url_input.text())
        s.setValue('patSrc', self.pat_src_combo.currentText())
        s.setValue('patKey', self.pat_key_input.text())
        s.setValue('patMcpUrl', self.mcp_url_input.text())
        s.setValue('patMcpTool', self.mcp_tool_input.text())
        s.setValue('patDeep', self.chk_pat_deep.isChecked())
        s.setValue('yearRange', self.year_combo.currentText())
        s.setValue('ncbiKey', self.ncbi_key_input.text())
        s.setValue('intlTrials', self.chk_intl.isChecked())
        s.setValue('aiCites', self.chk_ai_cites.isChecked())
        s.setValue('saveHist', self.chk_save_hist.isChecked())
        s.setValue('useLLM', self.chk_use_llm.isChecked())
        for cb, key in self._settings_sources():
            s.setValue(key, cb.isChecked())
        s.sync()

    def setup_ui(self):
        central = QWidget()
        self.setCentralWidget(central)
        outer = QVBoxLayout(central)
        outer.setContentsMargins(20, 15, 20, 15)

        title = QLabel('TargetInfo - 靶点调研报告生成器')
        title.setStyleSheet('font-size:20px;font-weight:700;color:#1B4F72;padding:6px 0;border-bottom:3px solid #2E86C1')
        outer.addWidget(title)
        sub = QLabel('PubMed | OpenAlex | Semantic Scholar | ClinicalTrials.gov | ISRCTN | ChiCTR | Open Targets | UniProt | PDB | AlphaFold | HPA | GTEx | STRING | ChEMBL | PubChem | KEGG | ClinPGx | ClinVar | Google Patents | Patent9 | drugfuture | USPTO | Lens.org | Espacenet')
        sub.setStyleSheet('color:#7F8C8D;font-size:11px;margin-bottom:6px')
        sub.setWordWrap(True)
        outer.addWidget(sub)

        # Input row
        inp = QHBoxLayout()
        lbl = QLabel('靶点名称:'); lbl.setStyleSheet('font-size:14px;font-weight:600')
        self.input_field = QLineEdit()
        self.input_field.setPlaceholderText('例如: PD-1, EGFR, KRAS G12C, HER2...')
        ss = 'QLineEdit{padding:8px 12px;font-size:14px;border:2px solid #ddd;border-radius:6px} QLineEdit:focus{border-color:#2E86C1}'
        self.input_field.setStyleSheet(ss)
        self.start_btn = QPushButton('开始生成')
        bss = 'QPushButton{padding:8px 24px;font-size:14px;font-weight:600;background:#2E86C1;color:white;border:none;border-radius:6px} QPushButton:hover{background:#2874A6} QPushButton:disabled{background:#BDC3C7}'
        self.start_btn.setStyleSheet(bss)
        self.start_btn.clicked.connect(self.start_pipeline)
        self.input_field.returnPressed.connect(self.start_pipeline)
        inp.addWidget(lbl); inp.addWidget(self.input_field, 1); inp.addWidget(self.start_btn)
        outer.addLayout(inp)

        # Settings panel (collapsible)
        self.settings_group = QGroupBox('设置 (数据源 | AI | 专利 | 文献)')
        self.settings_group.setCheckable(True)
        self.settings_group.setChecked(False)
        sgs = 'QGroupBox{font-weight:600;font-size:13px;padding:10px;border:1px solid #d0d0d0;border-radius:6px;margin:4px 0} QGroupBox::indicator{width:18px;height:18px}'
        self.settings_group.setStyleSheet(sgs)
        settings_widget = QWidget()
        settings_layout = QVBoxLayout(settings_widget)
        settings_layout.setSpacing(6)

        # Row 1: Data sources
        row1 = QHBoxLayout()
        row1.addWidget(QLabel('数据源:'))
        self.chk_papers = QCheckBox('文献'); self.chk_papers.setChecked(True)
        self.chk_trials = QCheckBox('临床'); self.chk_trials.setChecked(True)
        self.chk_target = QCheckBox('靶点'); self.chk_target.setChecked(True)
        self.chk_drugs = QCheckBox('药物'); self.chk_drugs.setChecked(True)
        self.chk_patents = QCheckBox('专利'); self.chk_patents.setChecked(True)
        for cb in [self.chk_papers, self.chk_trials, self.chk_target, self.chk_drugs, self.chk_patents]:
            row1.addWidget(cb)
        settings_layout.addLayout(row1)

        # Row 2: AI provider
        row2 = QHBoxLayout()
        row2.addWidget(QLabel('AI供应商:'))
        self.provider_combo = QComboBox()
        self.provider_combo.addItems(['deepseek', 'mimo', 'zhipu', 'custom'])
        self.provider_combo.currentIndexChanged.connect(self._on_provider_change)
        row2.addWidget(self.provider_combo)
        row2.addWidget(QLabel('模型:'))
        self.model_combo = QComboBox()
        self.model_combo.setEditable(True)
        self.model_combo.setInsertPolicy(QComboBox.NoInsert)
        row2.addWidget(self.model_combo, 1)
        row2.addWidget(QLabel('API Key:'))
        self.api_key_input = QLineEdit()
        self.api_key_input.setPlaceholderText('对应供应商 API Key')
        self.api_key_input.setEchoMode(QLineEdit.Password)
        row2.addWidget(self.api_key_input, 1)
        settings_layout.addLayout(row2)

        # Row 2b: Custom Base URL (hidden unless custom)
        self.custom_row = QWidget()
        cr_layout = QHBoxLayout(self.custom_row)
        cr_layout.setContentsMargins(0, 0, 0, 0)
        cr_layout.addWidget(QLabel('Base URL:'))
        self.base_url_input = QLineEdit()
        self.base_url_input.setPlaceholderText('https://.../v1')
        cr_layout.addWidget(self.base_url_input, 1)
        self.custom_row.setVisible(False)
        settings_layout.addWidget(self.custom_row)

        # Row 3: LLM enable + AI cites + save hist
        row3 = QHBoxLayout()
        self.chk_use_llm = QCheckBox('启用AI增强'); self.chk_use_llm.setChecked(True)
        self.chk_ai_cites = QCheckBox('AI结构化引用'); self.chk_ai_cites.setChecked(True)
        self.chk_save_hist = QCheckBox('自动保存历史'); self.chk_save_hist.setChecked(True)
        row3.addWidget(self.chk_use_llm); row3.addWidget(self.chk_ai_cites); row3.addWidget(self.chk_save_hist)
        row3.addStretch()
        settings_layout.addLayout(row3)

        # Row 4: Patent
        row4 = QHBoxLayout()
        row4.addWidget(QLabel('专利源:'))
        self.pat_src_combo = QComboBox()
        self.pat_src_combo.addItems(['google', 'uspto', 'lens', 'espacenet', 'mcp'])
        self.pat_src_combo.currentIndexChanged.connect(self._on_pat_src_change)
        row4.addWidget(self.pat_src_combo)
        row4.addWidget(QLabel('Key/Token:'))
        self.pat_key_input = QLineEdit()
        self.pat_key_input.setPlaceholderText('免Key可留空')
        row4.addWidget(self.pat_key_input, 1)
        self.chk_pat_deep = QCheckBox('专利深度增强'); self.chk_pat_deep.setChecked(True)
        row4.addWidget(self.chk_pat_deep)
        settings_layout.addLayout(row4)

        # Row 4b: MCP settings (hidden unless mcp)
        self.mcp_row = QWidget()
        mcp_layout = QHBoxLayout(self.mcp_row)
        mcp_layout.setContentsMargins(0, 0, 0, 0)
        mcp_layout.addWidget(QLabel('MCP URL:'))
        self.mcp_url_input = QLineEdit()
        self.mcp_url_input.setPlaceholderText('https://.../mcp')
        mcp_layout.addWidget(self.mcp_url_input, 1)
        mcp_layout.addWidget(QLabel('工具名:'))
        self.mcp_tool_input = QLineEdit()
        self.mcp_tool_input.setPlaceholderText('留空自动探测')
        mcp_layout.addWidget(self.mcp_tool_input, 1)
        self.mcp_row.setVisible(False)
        settings_layout.addWidget(self.mcp_row)

        # Row 5: Year range + NCBI key + Intl trials
        row5 = QHBoxLayout()
        row5.addWidget(QLabel('文献年份:'))
        self.year_combo = QComboBox()
        self.year_combo.addItems(['20', '10', '5', '0'])
        self.year_combo.setCurrentText('20')
        row5.addWidget(self.year_combo)
        row5.addWidget(QLabel('NCBI Key:'))
        self.ncbi_key_input = QLineEdit()
        self.ncbi_key_input.setPlaceholderText('可选，提升PubMed配额')
        self.ncbi_key_input.setEchoMode(QLineEdit.Password)
        row5.addWidget(self.ncbi_key_input, 1)
        self.chk_intl = QCheckBox('国际临床覆盖(ISRCTN/ChiCTR)')
        row5.addWidget(self.chk_intl)
        settings_layout.addLayout(row5)

        self.settings_group.setLayout(QVBoxLayout())
        self.settings_group.layout().addWidget(settings_widget)
        outer.addWidget(self.settings_group)

        # Progress
        self.progress_bar = QProgressBar()
        self.progress_bar.setFixedHeight(22)
        self.progress_label = QLabel('')
        self.progress_label.setAlignment(Qt.AlignCenter)
        self.progress_widget = QWidget()
        pw = QVBoxLayout(self.progress_widget)
        pw.setContentsMargins(0, 0, 0, 0)
        pw.addWidget(self.progress_bar); pw.addWidget(self.progress_label)
        self.stop_btn = QPushButton('停止生成')
        stopss = 'QPushButton{padding:6px 18px;font-size:12px;font-weight:600;background:#C0392B;color:white;border:none;border-radius:5px} QPushButton:hover{background:#A93226}'
        self.stop_btn.setStyleSheet(stopss)
        self.stop_btn.setVisible(False)
        self.stop_btn.clicked.connect(self._stop_pipeline)
        pw.addWidget(self.stop_btn, alignment=Qt.AlignRight)
        self.progress_widget.setVisible(False)
        outer.addWidget(self.progress_widget)

        # Results scroll
        self.scroll = QScrollArea()
        self.scroll.setWidgetResizable(True)
        self.scroll.setVisible(False)
        self.result_content = QWidget()
        self.result_layout = QVBoxLayout(self.result_content)
        self.scroll.setWidget(self.result_content)
        outer.addWidget(self.scroll, 1)

        # History
        self.history = HistoryWidget()
        self.history.selected.connect(self._load_history)
        self.history.compare.connect(self._load_history_compare)
        outer.addWidget(self.history)

    def _on_provider_change(self, idx):
        prov = self.provider_combo.currentText()
        self.custom_row.setVisible(prov == 'custom')
        p = LLM_PROVIDERS.get(prov, {})
        self.model_combo.blockSignals(True)
        self.model_combo.clear()
        self.model_combo.addItems(p.get('models') or [])
        self.model_combo.blockSignals(False)
        base = p.get('base', '')
        if base:
            self.base_url_input.setPlaceholderText(base + ' (可覆盖)')
        else:
            self.base_url_input.setPlaceholderText('https://.../v1')
        self.model_combo.setCurrentIndex(-1)

    def _on_pat_src_change(self, idx):
        self.mcp_row.setVisible(self.pat_src_combo.currentText() == 'mcp')

    def start_pipeline(self):
        name = self.input_field.text().strip()
        if not name:
            QMessageBox.warning(self, '提示', '请输入靶点名称'); return
        gene, mutation, raw = parse_target(name)

        provider = self.provider_combo.currentText()
        api_key = self.api_key_input.text().strip()
        model = self.model_combo.currentText().strip()
        base_url = self.base_url_input.text().strip()

        if provider == 'custom':
            if not base_url:
                QMessageBox.warning(self, '提示', '自定义模式请填写 Base URL'); return
        else:
            if not base_url:
                base_url = LLM_PROVIDERS.get(provider, {}).get('base', '')
        if not model:
            p = LLM_PROVIDERS.get(provider, {})
            models = p.get('models') or []
            model = models[0] if models else 'deepseek-chat'

        self.start_btn.setEnabled(False)
        self.input_field.setEnabled(False)
        self.stop_btn.setVisible(True)
        self.progress_widget.setVisible(True)
        self.scroll.setVisible(False)
        self.progress_bar.setValue(0)
        self.progress_label.setText('启动中...')

        self.worker = PipelineWorker(
            name, gene, mutation, provider, api_key, model, base_url,
            use_llm=self.chk_use_llm.isChecked(),
            src_papers=self.chk_papers.isChecked(),
            src_trials=self.chk_trials.isChecked(),
            src_target=self.chk_target.isChecked(),
            src_drugs=self.chk_drugs.isChecked(),
            src_patents=self.chk_patents.isChecked(),
            patent_src=self.pat_src_combo.currentText(),
            patent_key=self.pat_key_input.text().strip(),
            mcp_url=self.mcp_url_input.text().strip(),
            mcp_tool=self.mcp_tool_input.text().strip(),
            pat_deep=self.chk_pat_deep.isChecked(),
            year_range=self.year_combo.currentText(),
            ncbi_key=self.ncbi_key_input.text().strip(),
            intl_trials=self.chk_intl.isChecked(),
            ai_cites=self.chk_ai_cites.isChecked(),
            save_hist=self.chk_save_hist.isChecked(),
        )
        self.worker.progress.connect(self._on_progress)
        self.worker.finished.connect(self._on_finished)
        self.worker.error.connect(self._on_error)
        self.worker.start()

    def _stop_pipeline(self):
        if self.worker:
            self.worker.request_stop()
            self.stop_btn.setEnabled(False)
            self.progress_label.setText('正在停止...')
            self.start_btn.setEnabled(True)
            self.input_field.setEnabled(True)

    def _on_progress(self, value, text):
        self.progress_bar.setValue(value)
        self.progress_label.setText(text)

    def _on_finished(self, report):
        self.progress_bar.setValue(100)
        self.progress_label.setText('完成!')
        self.start_btn.setEnabled(True)
        self.input_field.setEnabled(True)
        self.stop_btn.setVisible(False)
        self._report = report
        if report is None:
            self.progress_widget.setVisible(False)
            QMessageBox.information(self, '已停止', '本次生成已取消。')
            return
        self._render_report(report)
        if self.chk_save_hist.isChecked():
            report_save(report)
            self.history.refresh()

    def _on_error(self, err_msg):
        self.start_btn.setEnabled(True)
        self.input_field.setEnabled(True)
        self.stop_btn.setVisible(False)
        self.progress_widget.setVisible(False)
        QMessageBox.critical(self, '错误', '生成过程中出现错误:' + NL + NL + err_msg)

    def _render_report(self, report):
        for i in reversed(range(self.result_layout.count())):
            w = self.result_layout.itemAt(i).widget()
            if w: w.deleteLater()
        self.scroll.setVisible(True)

        success = QLabel('报告生成完成')
        success.setStyleSheet('background:#d4edda;color:#155724;padding:10px;border-radius:6px;font-size:14px;font-weight:600')
        success.setAlignment(Qt.AlignCenter)
        self.result_layout.addWidget(success)

        cards = [('文献', len(report.papers)), ('临床', len(report.trials)),
                 ('药物', len(report.drugs)), ('分子', len(report.molecules)),
                 ('专利', len(report.patents))]
        metrics_row = QHBoxLayout()
        for label, count in cards:
            m = QLabel('<b>' + str(count) + '</b><br><span style="color:#7F8C8D">' + label + '</span>')
            m.setAlignment(Qt.AlignCenter)
            m.setStyleSheet('background:white;border:1px solid #e0e0e0;border-radius:8px;padding:12px;min-width:80px')
            metrics_row.addWidget(m)
        self.result_layout.addLayout(metrics_row)

        hints = []
        if report.increment_note:
            hints.append(report.increment_note)
        for f in report.cached_flags:
            hints.append(f + ' TTL 24h')
        if hints:
            hint = QLabel(' · '.join(hints))
            hint.setWordWrap(True)
            hint.setStyleSheet('background:#FFF9E6;color:#7D6608;padding:8px 12px;border-radius:6px;font-size:12px')
            self.result_layout.addWidget(hint)

        if report.errors:
            errb = QLabel('部分数据源失败:' + NL + NL.join('- ' + e for e in report.errors[:10]))
            errb.setWordWrap(True)
            errb.setStyleSheet('background:#FDEDEC;color:#922B21;padding:10px;border-radius:6px;font-size:12px')
            self.result_layout.addWidget(errb)

        dl_row = QHBoxLayout()
        for fmt, lbl, clr in [('md', '下载 MD', '#27AE60'), ('json', '下载 JSON', '#7F8C8D'),
                               ('html', '下载 HTML', '#D35400'), ('ppt', '下载 PPT', '#2E86C1'),
                               ('docx', '下载 Word', '#1B4F72')]:
            btn = QPushButton(lbl)
            btn.setStyleSheet('QPushButton{padding:8px 14px;font-size:12px;font-weight:600;background:' + clr + ';color:white;border:none;border-radius:4px}')
            btn.clicked.connect(lambda checked, f=fmt: self._save(f))
            dl_row.addWidget(btn)
        copybtn = QPushButton('复制 MD')
        copybtn.setStyleSheet('QPushButton{padding:8px 14px;font-size:12px;font-weight:600;background:#6C3483;color:white;border:none;border-radius:4px}')
        copybtn.clicked.connect(self._copy_markdown)
        dl_row.addWidget(copybtn)
        printbtn = QPushButton('打印/PDF')
        printbtn.setStyleSheet('QPushButton{padding:8px 14px;font-size:12px;font-weight:600;background:#17202A;color:white;border:none;border-radius:4px}')
        printbtn.clicked.connect(self._print_html)
        dl_row.addWidget(printbtn)
        savebtn = QPushButton('保存历史')
        savebtn.setStyleSheet('QPushButton{padding:8px 14px;font-size:12px;font-weight:600;background:#1ABC9C;color:white;border:none;border-radius:4px}')
        savebtn.clicked.connect(self._save_history)
        dl_row.addWidget(savebtn)
        self.result_layout.addLayout(dl_row)

        d = report.target_detail
        info = []
        if d:
            if d.gene_symbol: info.append('基因符号: ' + d.gene_symbol)
            if d.protein_class: info.append('蛋白类别: ' + d.protein_class)
            if d.tractability: info.append('成药性: ' + ', '.join(d.tractability[:3]))
            if d.related_diseases: info.append('相关疾病: ' + ', '.join(d.related_diseases[:5]))
            if d.pdb: info.append('PDB结构: ' + str(len(d.pdb)) + '个')
            if d.alphafold: info.append('AlphaFold: ' + str(len(d.alphafold)) + '个结构')
            if d.string_interactions: info.append('STRING互作: ' + str(len(d.string_interactions)) + '个')
            if d.kegg_pathways: info.append('KEGG通路: ' + str(len(d.kegg_pathways)) + '条')
            if d.pubchem: info.append('PubChem: ' + str(len(d.pubchem)) + '条')
            if d.kegg_disease: info.append('KEGG疾病: ' + ', '.join(d.kegg_disease[:3]))
            if d.kegg_drugs: info.append('KEGG药物: ' + ', '.join(d.kegg_drugs[:3]))
            if d.clinvar: info.append('ClinVar变异: ' + str(len(d.clinvar)) + '条')
            if d.clinpgx:
                info.append('ClinPGx药物基因组关联: ' + str(len(d.clinpgx)) + '条')
                for a in d.clinpgx[:12]:
                    bits = []
                    if a.drug_name: bits.append('药物: ' + a.drug_name)
                    if a.level: bits.append('Level: ' + a.level)
                    if a.disease: bits.append('疾病: ' + a.disease)
                    if a.association_type: bits.append('类型: ' + a.association_type)
                    if a.evidence: bits.append('证据: ' + a.evidence)
                    if bits: info.append('  - ' + ' | '.join(bits))
            if d.hpa: info.append('HPA表达: ' + (d.hpa.protein_tissue or d.hpa.rna_cancer or '有数据'))
            if d.gtex: info.append('GTEx组织: ' + str(len(d.gtex)) + '项')
            if d.sequence: info.append('蛋白序列: ' + str(d.sequence_length or len(d.sequence)) + ' aa')
        if info:
            self._render_section('靶点详情', html_text='<br>'.join(escape(x) for x in info))
        if d and d.sequence:
            fasta = ('&gt;%s|%s' % (escape(d.gene_symbol or self.input_field.text()), escape(d.uniprot_id or ''))) \
                + NL + escape(d.sequence)
            self._render_section('序列 FASTA', html_text='<pre style="white-space:pre-wrap;word-break:break-all">' + fasta + '</pre>')

        # 研究进展: AI 文本 + 年份分布条 + 文献列表
        yc = self._year_counts(report.papers)
        years = sorted(yc)[-10:] if yc else []
        bars = [(y, yc[y]) for y in years]
        rp = ''
        if report.research_progress:
            rp = escape(report.research_progress).replace(NL, '<br>')
        if report.key_findings:
            rp += '<br><br><b>核心发现</b><br>' + escape(report.key_findings).replace(NL, '<br>')
        if report.papers:
            rp += '<br><br><b>文献列表（%d 篇）</b><br>' % len(report.papers) \
                + '<br>'.join(self._paper_html(p, i) for i, p in enumerate(report.papers))
        if rp or bars:
            self._render_section('研究进展', bars=bars, html_text=rp)

        # 临床概况: AI 文本 + 阶段/状态分布条 + 试验列表
        ph = self._count_key(report.trials, 'phase', '未明确')
        st = self._count_key(report.trials, 'status', '未知')
        cbars = [('阶段: ' + k, v) for k, v in ph.items()] + [('状态: ' + k, v) for k, v in st.items()]
        cl = ''
        if report.clinical_landscape:
            cl = escape(report.clinical_landscape).replace(NL, '<br>')
        if report.trials:
            cl += '<br><br><b>试验列表（%d 项）</b><br>' % len(report.trials) \
                + '<br>'.join(self._trial_html(t, i) for i, t in enumerate(report.trials))
        if cl or cbars:
            self._render_section('临床概况', bars=cbars, html_text=cl)

        # 药物展望: AI 文本 + 药物管线 + ChEMBL 分子
        fo = ''
        if report.future_outlook:
            fo = escape(report.future_outlook).replace(NL, '<br>')
        if report.drugs:
            fo += '<br><br><b>药物管线（%d 个）</b><br>' % len(report.drugs) \
                + '<br>'.join(self._drug_html(x, i) for i, x in enumerate(report.drugs))
        if report.molecules:
            fo += '<br><br><b>靶向活性分子（ChEMBL，%d 个）</b><br>' % len(report.molecules) \
                + '<br>'.join(self._molecule_html(x, i) for i, x in enumerate(report.molecules))
        if fo:
            self._render_section('药物展望', html_text=fo)

        # 专利调研: AI 文本 + 专利列表 + 申请人布局
        pt = ''
        if report.patent_landscape:
            pt = escape(report.patent_landscape).replace(NL, '<br>')
        if report.patents:
            pt += '<br><br><b>专利列表（%d 件）</b><br>' % len(report.patents) \
                + '<br>'.join(self._patent_html(x, i) for i, x in enumerate(report.patents))
        if report.patent_insight and report.patent_insight.get('assignees'):
            entries = [(a, ps) for a, ps in report.patent_insight['assignees'].items() if ps]
            if entries:
                pt += '<br><br><b>🧭 主要申请人布局</b><br>' + '<br>'.join(
                    '<b>%s</b>（%d 件）：%s' % (escape(str(a)), len(ps), '；'.join(
                        '%s %s' % (escape(str(x.get('number', ''))), escape(str(x.get('title', '')))[:60])
                        for x in ps[:6])) for a, ps in entries)
        if pt:
            self._render_section('专利调研', html_text=pt)

        if report.web_summary:
            self._render_section('AI综合研判', html_text=escape(report.web_summary).replace(NL, '<br>'))

        self.result_layout.addStretch()

        fq = QGroupBox('追问')
        fq.setCheckable(True); fq.setChecked(True)
        fq.setStyleSheet('QGroupBox{font-weight:600;font-size:13px;padding:8px;border:1px solid #ddd;border-radius:6px;margin-top:4px}')
        fql = QVBoxLayout(fq)
        fr = QHBoxLayout()
        self.fq_input = QLineEdit()
        self.fq_input.setPlaceholderText('针对本报告提问，例如: 该靶点最近有哪些临床试验阶段进展?')
        fr.addWidget(self.fq_input, 1)
        fq_btn = QPushButton('提问')
        fq_btn.clicked.connect(lambda: self._ask_followup())
        self.fq_input.returnPressed.connect(self._ask_followup)
        fr.addWidget(fq_btn)
        fql.addLayout(fr)
        self.fq_output = QTextBrowser()
        self.fq_output.setOpenExternalLinks(True)
        fql.addWidget(self.fq_output)
        self.result_layout.addWidget(fq)

    def _render_section(self, title, bars=None, html_text=''):
        sec = QGroupBox(title)
        sec.setCheckable(True); sec.setChecked(False)
        sec.setStyleSheet('QGroupBox{font-weight:600;font-size:13px;padding:8px;border:1px solid #ddd;border-radius:6px;margin-top:4px}')
        ql = QVBoxLayout(sec)
        if bars:
            maxv = max(v for _, v in bars) or 1
            for label, v in bars:
                row = QHBoxLayout()
                lb = QLabel(label)
                lb.setFixedWidth(100)
                row.addWidget(lb)
                bar = QProgressBar()
                bar.setRange(0, maxv); bar.setValue(v)
                bar.setTextVisible(False); bar.setFixedHeight(14)
                bar.setStyleSheet('QProgressBar{border:1px solid #bdc3c7;border-radius:4px;background:#ECF0F1} QProgressBar::chunk{background:#2E86C1;border-radius:4px}')
                row.addWidget(bar, 1)
                cnt = QLabel(str(v))
                cnt.setFixedWidth(28)
                row.addWidget(cnt)
                ql.addLayout(row)
        if html_text:
            tb = QTextBrowser()
            tb.setOpenExternalLinks(True)
            tb.setHtml('<div style="font-size:12px;line-height:1.6">' + html_text + '</div>')
            tb.setMinimumHeight(60); tb.setMaximumHeight(300)
            ql.addWidget(tb)
        self.result_layout.addWidget(sec)

    @staticmethod
    def _year_counts(papers):
        m = {}
        for p in papers:
            y = str(getattr(p, 'year', '') or '')
            if re.fullmatch(r'\d{4}', y):
                m[y] = m.get(y, 0) + 1
        return m

    @staticmethod
    def _count_key(items, attr, default):
        m = {}
        for x in items:
            k = str(getattr(x, attr, '') or '').strip() or default
            m[k] = m.get(k, 0) + 1
        return m

    @staticmethod
    def _paper_html(p, i):
        bits = '<b>%d.</b> %s' % (i + 1, escape(p.title or ''))
        if p.journal:
            bits += ' — <i>%s</i>' % escape(p.journal)
        if p.year:
            bits += ' (%s)' % escape(p.year)
        if p.pmid:
            bits += ' <a href="https://pubmed.ncbi.nlm.nih.gov/%s/">[PubMed]</a>' % escape(p.pmid)
        elif p.doi:
            bits += ' <a href="https://doi.org/%s">[DOI]</a>' % escape(p.doi)
        return bits

    @staticmethod
    def _trial_html(t, i):
        href = _link_for_ref(t.nct_id)
        bits = '<b>%d.</b> %s %s' % (i + 1,
                                     ('<a href="%s">%s</a>' % (href, escape(t.nct_id))) if href != '#' else escape(t.nct_id),
                                     escape(t.title or ''))
        meta = ' / '.join(x for x in [t.phase, t.status,
                                      ('、'.join(t.conditions[:4]) if t.conditions else ''),
                                      t.source] if x)
        if meta:
            bits += '<br>&nbsp;&nbsp;<span style="color:#7F8C8D">%s</span>' % escape(meta)
        return bits

    @staticmethod
    def _drug_html(d, i):
        bits = '<b>%d.</b> %s' % (i + 1, escape(d.name or ''))
        meta = ' / '.join(x for x in [d.company, d.phase, d.mechanism_of_action, d.disease] if x)
        if meta:
            bits += ' — <span style="color:#7F8C8D">%s</span>' % escape(meta)
        return bits

    @staticmethod
    def _molecule_html(m, i):
        bits = '<b>%d.</b> %s' % (i + 1, escape(m.name or '—'))
        if m.chembl:
            bits += ' <a href="https://www.ebi.ac.uk/chembl/compound_report_card/%s/">%s</a>' % (escape(m.chembl), escape(m.chembl))
        if m.pchembl:
            bits += ' (pChEMBL %s)' % escape(m.pchembl)
        return bits

    @staticmethod
    def _patent_html(p, i):
        num = p.number or ''
        href = p.link or ('https://patents.google.com/patent/' + num + '/en' if num else '#')
        bits = '<b>%d.</b> <a href="%s">%s</a> — %s' % (i + 1, escape(href), escape(num or '—'), escape(p.title or ''))
        extra = []
        if p.assignee: extra.append(p.assignee)
        if p.year: extra.append(p.year)
        if p.cited_by: extra.append('被引%d' % p.cited_by)
        if p.claims: extra.append('约%d项权利要求' % p.claims)
        if p.application_no: extra.append('申请号 ' + p.application_no)
        if p.cn_fulltext:
            extra.append('<a href="%s">全文PDF</a>' % escape(p.cn_fulltext))
        elif p.cn_download:
            extra.append('<a href="%s">全文下载页</a>' % escape(p.cn_download))
        if extra:
            bits += '<br>&nbsp;&nbsp;<span style="color:#7F8C8D">%s</span>' % escape(' · '.join(extra))
        return bits

    def _ask_followup(self):
        q = self.fq_input.text().strip()
        if not q or not self._report:
            return
        provider = self.provider_combo.currentText()
        api_key = self.api_key_input.text().strip()
        model = self.model_combo.currentText().strip()
        base_url = self.base_url_input.text().strip()
        if provider != 'custom' and not base_url:
            base_url = LLM_PROVIDERS.get(provider, {}).get('base', '')
        if not model:
            p = LLM_PROVIDERS.get(provider, {})
            models = p.get('models') or []
            model = models[0] if models else 'deepseek-chat'
        self.fq_output.setHtml('<div style="color:#7F8C8D;font-size:12px">🔍 检索并生成中…</div>')
        self.fq_worker = FollowUpWorker(self._report, q, provider, api_key, model, base_url)
        self.fq_worker.result.connect(self._on_followup_result)
        self.fq_worker.start()

    def _on_followup_result(self, text, docs):
        if not text:
            if not docs:
                self.fq_output.setHtml('<div style="color:#7F8C8D;font-size:12px">未从当前报告检索到相关内容，请换种问法。</div>')
                return
            lines = ['<div style="font-size:12px;line-height:1.7">']
            for doc in docs:
                link = _link_for_ref(doc['id'])
                ref = ('<a href="' + link + '">' + doc['id'] + '</a>') if link and link != '#' else doc['id']
                title = (doc['text'] or '')[:120]
                lines.append('<b>[' + doc['kind'] + ']</b> ' + ref + ' — ' + title)
            lines.append('</div>')
            self.fq_output.setHtml(NL.join(lines))
            return
        body = escape(text).replace(NL, '<br>')
        refs_line = ' · '.join(
            '<a href="%s">[%d] %s %s</a>' % (_link_for_ref(doc['id']), i + 1, doc['kind'], escape(doc['id']))
            for i, doc in enumerate(docs))
        self.fq_output.setHtml('<div style="font-size:12px;line-height:1.7">' + body + '</div>'
                               + '<div style="color:#7F8C8D;font-size:11px;margin-top:6px">' + refs_line + '</div>')

    def _save_history(self):
        if not self._report:
            return
        rid = report_save(self._report, self._report.errors)
        self.history.refresh()
        QMessageBox.information(self, '已保存', '已保存到历史' + ('' if rid else '（保存失败）'))

    def _copy_markdown(self):
        if not self._report: return
        md = to_markdown(self._report)
        QApplication.clipboard().setText(md)
        QMessageBox.information(self, '已复制', 'Markdown 已复制到剪贴板。')

    def _print_html(self):
        if not self._report: return
        html = to_html(self._report)
        path = os.path.join(tempfile.gettempdir(), slug(self._report.target_name) + '_report.html')
        with open(path, 'w', encoding='utf-8') as f:
            f.write(html)
        try:
            if sys.platform.startswith('win'):
                os.startfile(path)
            elif sys.platform.startswith('darwin'):
                subprocess.Popen(['open', path])
            else:
                subprocess.Popen(['xdg-open', path])
        except Exception:
            QMessageBox.information(self, '已生成', 'HTML 已生成: ' + path)

    def _save(self, fmt):
        if not self._report: return
        r = self._report; name = slug(r.target_name)
        ext_map = {'md': ('.md', 'Markdown (*.md)'), 'json': ('.json', 'JSON (*.json)'),
                   'html': ('.html', 'HTML (*.html)'), 'ppt': ('.pptx', 'PPT (*.pptx)'),
                   'docx': ('.docx', 'Word (*.docx)')}
        if fmt not in ext_map: return
        ext, filt = ext_map[fmt]
        path, _ = QFileDialog.getSaveFileName(self, '保存文件', name + '_报告' + ext, filt)
        if not path: return
        QApplication.setOverrideCursor(Qt.WaitCursor)
        try:
            writers = {'md': lambda: to_markdown(r).encode('utf-8'),
                       'json': lambda: to_json(r).encode('utf-8'),
                       'html': lambda: to_html(r).encode('utf-8'),
                       'ppt': lambda: generate_ppt(r),
                       'docx': lambda: generate_docx(r)}
            data = writers[fmt]()
            with open(path, 'wb') as f: f.write(data)
            QMessageBox.information(self, '保存成功', '文件已保存至:' + NL + path)
        except Exception as e:
            QMessageBox.critical(self, '保存失败', str(e))
        finally:
            QApplication.restoreOverrideCursor()

    def _load_history(self, rid):
        report = report_load(rid)
        if report:
            self._report = report; self._render_report(report)
        else:
            QMessageBox.warning(self, '提示', '无法加载该历史记录')

    def _load_history_compare(self, rid_a, rid_b):
        ra = report_load(rid_a); rb = report_load(rid_b)
        if not ra or not rb:
            QMessageBox.warning(self, '提示', '无法加载对比记录')
            return
        def sig(r):
            return (len(r.papers), len(r.trials), len(r.drugs), len(r.molecules), len(r.patents))
        a, b = sig(ra), sig(rb)
        def dl(r):
            return {k: v for k, v in [('papers', len(r.papers)), ('trials', len(r.trials)),
                                      ('drugs', len(r.drugs)), ('molecules', len(r.molecules)),
                                      ('patents', len(r.patents))]}
        da, db = dl(ra), dl(rb)
        lines = ['<div style="font-size:12px;line-height:1.7">',
                 '<b>对比:</b> ' + escape(ra.target_name) + ' <i>(基准)</i>  vs  ' + escape(rb.target_name) + '<br><br>']
        for k, zh in [('papers', '文献'), ('trials', '临床'), ('drugs', '药物'),
                      ('molecules', '分子'), ('patents', '专利')]:
            d = db.get(k, 0) - da.get(k, 0)
            arrow = '▲' if d > 0 else ('▼' if d < 0 else '—')
            color = '#27AE60' if d > 0 else ('#C0392B' if d < 0 else '#7F8C8D')
            lines.append('%s: <span style="color:#7F8C8D">%s</span> → <span style="color:#2E86C1">%s</span> '
                         '<span style="color:%s">%s %+d</span><br>' % (zh, da.get(k, 0), db.get(k, 0), color, arrow, d))
        if ra.increment_note:
            lines.append('<br><i>基准增量:</i> ' + escape(ra.increment_note))
        if ra.errors or rb.errors:
            lines.append('<br><i>基准失败数据源:</i> ' + (escape('; '.join(ra.errors)) or '无') +
                         '<br><i>对比失败数据源:</i> ' + (escape('; '.join(rb.errors)) or '无'))
        lines.append('</div>')
        box = QMessageBox(self)
        box.setWindowTitle('历史对比')
        box.setTextFormat(Qt.RichText)
        box.setText(NL.join(lines))
        box.setStandardButtons(QMessageBox.Ok)
        box.exec_()


# ─── Main ──────────────────────────────────────────────────────────────────


def main():
    app = QApplication(sys.argv)
    app.setStyle('Fusion')
    font = QFont(); font.setPointSize(10); app.setFont(font)
    w = MainWindow(); w.show()
    sys.exit(app.exec_())


if __name__ == '__main__':
    main()
