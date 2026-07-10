import asyncio
import io
import os
import re
import sys
import xml.etree.ElementTree as ET
from datetime import datetime
from typing import Dict, List, Optional

from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QPushButton, QLineEdit, QLabel, QProgressBar, QGroupBox,
    QTextEdit, QTextBrowser, QScrollArea, QFrame, QMessageBox,
    QGridLayout, QFileDialog, QSizePolicy, QSpacerItem
)
from PyQt5.QtCore import Qt, QThread, pyqtSignal, QUrl
from PyQt5.QtGui import QFont, QDesktopServices

import httpx
from dotenv import load_dotenv
from pydantic import BaseModel
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── Configuration ──────────────────────────────────────────────────────────

load_dotenv()

PUBMED_EMAIL = os.getenv("PUBMED_EMAIL", "user@example.com")
PUBMED_SEARCH_URL = "https://eutils.ncbi.nlm.nih.gov/entrez/eutils/esearch.fcgi"
PUBMED_FETCH_URL = "https://eutils.ncbi.nlm.nih.gov/entrez/eutils/efetch.fcgi"
CT_API_BASE = "https://clinicaltrials.gov/api/v2"
OPEN_TARGETS_URL = "https://api.platform.opentargets.org/api/v4/graphql"

# ─── Retry helper ──────────────────────────────────────────────────────────


async def _request_with_retry(
    client: httpx.AsyncClient,
    method: str,
    url: str,
    max_retries: int = 3,
    base_delay: float = 1.5,
    **kwargs,
) -> httpx.Response:
    last_exc: Optional[Exception] = None
    for attempt in range(max_retries):
        try:
            if method == "GET":
                return await client.get(url, **kwargs)
            elif method == "POST":
                return await client.post(url, **kwargs)
            else:
                raise ValueError(f"Unsupported method: {method}")
        except (
            httpx.RemoteProtocolError,
            httpx.ConnectError,
            httpx.TimeoutException,
            httpx.ReadError,
            httpx.WriteError,
        ) as e:
            last_exc = e
            if attempt < max_retries - 1:
                delay = base_delay * (2**attempt)
                await asyncio.sleep(delay)
            continue
    raise last_exc  # type: ignore[misc]


# ─── ZhipuAI LLM ────────────────────────────────────────────────────────────

ZHIPUAI_API_KEY: Optional[str] = os.getenv("ZHIPUAI_API_KEY")
ZHIPUAI_URL = "https://open.bigmodel.cn/api/paas/v4/chat/completions"
ZHIPUAI_MODEL = "GLM-4.7-Flash"


async def _llm_chat(prompt: str, api_key: str) -> str:
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }
    payload = {
        "model": ZHIPUAI_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.6,
        "max_tokens": 131072,
        "stream": True,
        "thinking": {
            "type": "enabled"
        },
    }
    async with httpx.AsyncClient(timeout=120) as client:
        resp = await _request_with_retry(
            client, "POST", ZHIPUAI_URL,
            json=payload, headers=headers,
        )
        if resp.status_code != 200:
            return ""
        choices = resp.json().get("choices", [])
        if not choices:
            return ""
        return choices[0].get("message", {}).get("content", "")


_LLM_PROMPTS = {
    "target_overview": (
        "你是一个生物医药领域的专家。请根据以下靶点信息，生成一份专业的靶点概述（中文，200-500字），"
        "涵盖基因符号、蛋白类别、功能描述和相关疾病。\n\n{context}"
    ),
    "research_progress": (
        "你是一个生物医药领域的专家。以下是关于某靶点的文献数据（年份分布和标题列表），"
        "请分析研究趋势，指出近三年的研究热点和方向（中文，150-500字）。\n\n{context}"
    ),
    "clinical_landscape": (
        "你是一个临床研究专家。以下是关于某靶点的临床试验数据（阶段分布、适应症列表），"
        "请分析临床开发格局，指出当前阶段分布特征和主要探索方向（中文，150-500字）。\n\n{context}"
    ),
    "key_findings": (
        "你是一个生物医药领域的专家。请根据以下文献标题列表，提取最有代表性的3-5个研究发现，"
        "逐条列出并简要说明其意义（中文，200-800字）。\n\n{context}"
    ),
    "future_outlook": (
        "你是一个药物研发专家。以下是针对某靶点的药物管线数据（已获批药物和在研药物），"
        "请分析其未来发展趋势和值得关注的方向（中文，150-500字）。\n\n{context}"
    ),
    "web_summary": (
        "你是一个生物医药行业情报分析师。请利用联网搜索能力，针对以下靶点，"
        "搜集并整理最新全网关键情报（中文，400-800字），涵盖以下方面：\n"
        "1. 关键里程碑事件：重要的研究发现、学术突破、指南更新等\n"
        "2. 重大收购与合作：药企之间的并购、授权引进、战略合作等\n"
        "3. 临床试验重大进展：近期成功的 III 期结果、FDA 突破性疗法认定、"
        "重大失败或终止的试验及其影响\n"
        "4. 监管动态：FDA/EMA/NMPA 批准、加速审批、孤儿药认定等\n"
        "5. 竞争格局变化：新进入者、专利到期、生物类似药进展等\n\n"
        "请确保每条信息标注来源（如新闻标题或机构名称），"
        "并说明该信息对药物研发的潜在影响。\n\n靶点名称: {target_name}\n"
        "基因符号: {gene_symbol}\n蛋白类别: {protein_class}"
    ),
}


async def _llm_chat_with_search(prompt: str, api_key: str) -> str:
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }

    async def _try_payload(payload) -> str:
        async with httpx.AsyncClient(timeout=180) as client:
            resp = await _request_with_retry(
                client, "POST", ZHIPUAI_URL,
                json=payload, headers=headers,
                max_retries=2,
            )
            if resp.status_code != 200:
                return ""
            choices = resp.json().get("choices", [])
            if not choices:
                return ""
            return choices[0].get("message", {}).get("content", "") or ""

    # 1) Try with web_search flag (supported by GLM-4.7-Flash)
    text = await _try_payload({
        "model": ZHIPUAI_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.5,
        "max_tokens": 4096,
        "web_search": True,
    })
    if text:
        return text

    # 2) Fallback: tools parameter
    text = await _try_payload({
        "model": ZHIPUAI_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.5,
        "max_tokens": 4096,
        "tools": [{"type": "web_search", "web_search": {"enable": True}}],
    })
    if text:
        return text

    # 3) Final fallback: regular chat (model knowledge, no web search)
    return await _llm_chat(prompt, api_key)


# ─── Data Models ────────────────────────────────────────────────────────────


class Paper(BaseModel):
    pmid: str = ""
    title: str = ""
    authors: List[str] = []
    journal: str = ""
    year: str = ""
    doi: str = ""
    abstract: str = ""
    keywords: List[str] = []


class ClinicalTrial(BaseModel):
    nct_id: str = ""
    title: str = ""
    phase: str = ""
    status: str = ""
    conditions: List[str] = []
    interventions: List[str] = []
    sponsor: str = ""
    start_date: str = ""
    completion_date: str = ""
    brief_summary: str = ""
    enrollment: int = 0


class DrugInfo(BaseModel):
    name: str = ""
    company: str = ""
    phase: str = ""
    mechanism_of_action: str = ""
    disease: str = ""


class TargetDetail(BaseModel):
    target_name: str = ""
    gene_symbol: str = ""
    uniprot_id: str = ""
    protein_class: str = ""
    description: str = ""
    synonyms: List[str] = []
    related_diseases: List[str] = []


class ReportContent(BaseModel):
    target_name: str
    target_detail: Optional[TargetDetail] = None
    papers: List[Paper] = []
    trials: List[ClinicalTrial] = []
    drugs: List[DrugInfo] = []
    target_overview: str = ""
    research_progress: str = ""
    clinical_landscape: str = ""
    key_findings: str = ""
    future_outlook: str = ""
    web_summary: str = ""

# ─── PubMed Literature Search ──────────────────────────────────────────────


def _expand_query_terms(term: str) -> List[str]:
    seen = set()
    out = []
    def add(t):
        t = t.strip()
        if t and t not in seen:
            seen.add(t)
            out.append(t)

    add(term)
    add(term.upper())
    add(term.lower())
    add(term.capitalize())

    plain = re.sub(r'[\-–—.\' _]+', '', term)
    if plain != term:
        add(plain)
        add(plain.upper())
        add(plain.lower())

    h1 = re.sub(r'([a-zA-Z])(\d)', r'\1-\2', term)
    h2 = re.sub(r'(\d)([a-zA-Z])', r'\1-\2', term)
    if h1 != term:
        add(h1)
    if h2 != term:
        add(h2)

    for src, dst in [("2", "II"), ("1", "I")]:
        rv = term.replace(src, dst)
        if rv != term:
            add(rv)
            add(rv.upper())
            add(rv.lower())

    return out


async def search_papers(target: str, max_results: int = 40) -> List[Paper]:
    now_year = datetime.now().year
    terms = _expand_query_terms(target)
    query_parts = [f"({t}[Title/Abstract])" for t in terms]
    query = " OR ".join(query_parts)

    params = {
        "db": "pubmed",
        "term": f"({query}) AND (review[pt] OR clinical trial[pt] OR systematic review[pt])",
        "retmax": max_results, "retmode": "json", "sort": "relevance",
        "email": PUBMED_EMAIL,
        "mindate": str(now_year - 20), "maxdate": str(now_year), "datetype": "pdat",
    }
    async with httpx.AsyncClient(timeout=60) as client:
        resp = await _request_with_retry(client, "GET", PUBMED_SEARCH_URL, params=params)
        resp.raise_for_status()
        id_list = resp.json().get("esearchresult", {}).get("idlist", [])

        if not id_list:
            params["term"] = f"({query})"
            params["mindate"] = str(now_year - 20)
            params["maxdate"] = str(now_year)
            resp = await _request_with_retry(client, "GET", PUBMED_SEARCH_URL, params=params)
            resp.raise_for_status()
            id_list = resp.json().get("esearchresult", {}).get("idlist", [])

        if not id_list:
            return []

        return await fetch_papers(id_list)


async def fetch_papers(pmids: List[str]) -> List[Paper]:
    params = {
        "db": "pubmed", "id": ",".join(pmids),
        "retmode": "xml", "rettype": "abstract",
        "email": PUBMED_EMAIL,
    }
    async with httpx.AsyncClient(timeout=60) as client:
        resp = await _request_with_retry(client, "GET", PUBMED_FETCH_URL, params=params, max_retries=4)
        resp.raise_for_status()
        return _parse_pubmed_xml(resp.text)


def _parse_pubmed_xml(xml_str: str) -> List[Paper]:
    papers = []
    for article_elem in ET.fromstring(xml_str).findall(".//PubmedArticle"):
        try:
            paper = _parse_single_article(article_elem)
            if paper:
                papers.append(paper)
        except Exception:
            continue
    return papers


def _parse_single_article(article_elem) -> Optional[Paper]:
    medline = article_elem.find(".//MedlineCitation")
    article = medline.find(".//Article") if medline is not None else None
    if medline is None or article is None:
        return None

    title_elem = article.find("./ArticleTitle")
    title = "".join(title_elem.itertext()) if title_elem is not None else ""

    journal_elem = article.find("./Journal/Title")
    journal = "".join(journal_elem.itertext()) if journal_elem is not None else ""

    year_elem = article.find("./Journal/JournalIssue/PubDate/Year")
    year = "".join(year_elem.itertext()) if year_elem is not None else ""

    pmid_elem = medline.find("./PMID")
    pmid = "".join(pmid_elem.itertext()) if pmid_elem is not None else ""

    abstract = ""
    abstract_elem = article.find("./Abstract")
    if abstract_elem is not None:
        parts = []
        for elem in abstract_elem.iter():
            if elem.tag == "AbstractText":
                label = elem.get("Label", "")
                text = "".join(elem.itertext())
                parts.append(f"{label}: {text}" if label else text)
        abstract = " ".join(parts)

    authors = []
    author_list = article.find("./AuthorList")
    if author_list is not None:
        for author in author_list:
            last = author.find("./LastName")
            fore = author.find("./ForeName")
            if last is not None:
                name = last.text or ""
                if fore is not None and fore.text:
                    name += " " + fore.text
                authors.append(name)

    keywords = []
    kw_list = medline.find(".//KeywordList")
    if kw_list is not None:
        for kw in kw_list:
            if kw.text:
                keywords.append(kw.text)

    doi = ""
    for eid in article_elem.findall(".//ArticleIdList/ArticleId"):
        if eid.get("IdType") == "doi":
            doi = eid.text or ""
            break

    return Paper(
        pmid=pmid, title=title, authors=authors[:5],
        journal=journal, year=year, doi=doi,
        abstract=abstract[:2000], keywords=keywords,
    )

# ─── ClinicalTrials.gov API ────────────────────────────────────────────────


async def search_trials(target: str, max_results: int = 20) -> List[ClinicalTrial]:
    now_year = datetime.now().year
    terms = _expand_query_terms(target)
    query = " OR ".join(terms)
    params = {
        "query.term": query,
        "pageSize": min(max_results, 100), "format": "json",
        "fields": "NCTId|BriefTitle|OverallStatus|Phase|Condition|"
                  "InterventionType|InterventionName|LeadSponsorName|StartDate|CompletionDate|"
                  "BriefSummary|EnrollmentCount",
        "filter.overallStatus": "ACTIVE_NOT_RECRUITING|COMPLETED|RECRUITING|NOT_YET_RECRUITING|ENROLLING_BY_INVITATION|AVAILABLE",
    }
    async with httpx.AsyncClient(timeout=60) as client:
        resp = await _request_with_retry(client, "GET", f"{CT_API_BASE}/studies", params=params)
        if resp.status_code != 200:
            return []
        return [_parse_study(s) for s in resp.json().get("studies", [])]


def _parse_study(study: dict) -> ClinicalTrial:
    p = study.get("protocolSection", {})
    return ClinicalTrial(
        nct_id=p.get("identificationModule", {}).get("nctId", ""),
        title=p.get("identificationModule", {}).get("briefTitle", ""),
        status=p.get("statusModule", {}).get("overallStatus", ""),
        start_date=_fmt_date(p.get("statusModule", {}).get("startDate", {})),
        completion_date=_fmt_date(p.get("statusModule", {}).get("completionDate", {})),
        phase=(p.get("designModule", {}).get("phases") or [""])[0],
        conditions=p.get("conditionsModule", {}).get("conditions", []),
        interventions=[
            f"[{a.get('interventionType','')}] {a['interventionName']}"
            for a in p.get("armsInterventionsModule", {}).get("interventions", [])
            if a.get("interventionName")
        ],
        sponsor=(p.get("sponsorCollaboratorsModule", {}).get("leadSponsor") or {}).get("name", ""),
        brief_summary=(p.get("descriptionModule", {}).get("briefSummary") or "")[:1000],
        enrollment=(p.get("designModule", {}).get("enrollmentInfo") or {}).get("count", 0),
    )


def _fmt_date(d: dict) -> str:
    return "-".join(str(d[k]) for k in ["year", "month", "day"] if d.get(k))

# ─── Open Targets GraphQL API ──────────────────────────────────────────────


async def _get_ensembl_id(target_name: str) -> Optional[str]:
    q = """query($q:String!){search(queryString:$q,entityNames:["target"]){hits{id}}}"""
    async with httpx.AsyncClient(timeout=30) as client:
        resp = await _request_with_retry(client, "POST", OPEN_TARGETS_URL,
                                         json={"query": q, "variables": {"q": target_name}})
        hits = resp.json().get("data", {}).get("search", {}).get("hits", [])
        return hits[0]["id"] if hits else None


async def get_target_detail(target_name: str) -> TargetDetail:
    detail = TargetDetail(target_name=target_name)
    eid = await _get_ensembl_id(target_name)
    if not eid:
        return detail

    q = """query($id:String!){
      target(ensemblId:$id){
        approvedSymbol functionDescriptions
        nameSynonyms{label source}
        targetClass{label}
        associatedDiseases{rows{disease{name}}}
      }
    }"""
    async with httpx.AsyncClient(timeout=30) as client:
        resp = await _request_with_retry(client, "POST", OPEN_TARGETS_URL,
                                         json={"query": q, "variables": {"id": eid}})
        data = resp.json().get("data", {}).get("target", {})
        if not data:
            return detail

        detail.gene_symbol = data.get("approvedSymbol", "")
        descs = data.get("functionDescriptions") or []
        detail.description = descs[0] if descs else ""
        detail.synonyms = [
            s["label"] for s in (data.get("nameSynonyms") or [])
        ]
        cls = data.get("targetClass") or []
        detail.protein_class = cls[0].get("label", "") if cls else ""
        detail.related_diseases = [
            r["disease"]["name"]
            for r in (data.get("associatedDiseases") or {}).get("rows", [])[:10]
            if r.get("disease", {}).get("name")
        ]
    return detail


async def get_drugs(target_name: str) -> List[DrugInfo]:
    eid = await _get_ensembl_id(target_name)
    if not eid:
        return []

    q = """query($id:String!){
      target(ensemblId:$id){
        drugAndClinicalCandidates{
          rows{
            maxClinicalStage
            drug{
              name
              mechanismsOfAction{rows{mechanismOfAction actionType}}
            }
            diseases{diseaseFromSource disease{name}}
          }
        }
      }
    }"""
    async with httpx.AsyncClient(timeout=30) as client:
        resp = await _request_with_retry(client, "POST", OPEN_TARGETS_URL,
                                         json={"query": q, "variables": {"id": eid}})
        rows = (resp.json().get("data", {}).get("target", {})
                .get("drugAndClinicalCandidates", {})).get("rows", [])
        return [
            DrugInfo(
                name=r.get("drug", {}).get("name", ""),
                mechanism_of_action=(
                    (r.get("drug", {}).get("mechanismsOfAction", {})
                     .get("rows") or [{}])[0].get("mechanismOfAction", "")
                ),
                phase=r.get("maxClinicalStage", ""),
                disease=(
                    (r.get("diseases") or [{}])[0]
                     .get("disease", {}) or {}
                ).get("name", "")
                or (r.get("diseases") or [{}])[0].get("diseaseFromSource", ""),
            )
            for r in rows
        ]

async def _enrich_drugs_with_company(drugs: List[DrugInfo]) -> None:
    """Populate DrugInfo.company by searching each drug name on ClinicalTrials.gov."""
    names = [d.name for d in drugs if d.name]
    if not names:
        return
    query = " OR ".join(f'"{n}"' for n in names)
    params = {"query.term": query, "pageSize": min(len(names) * 5, 100),
              "format": "json",
              "fields": "NCTId|InterventionType|InterventionName|LeadSponsorName"}
    async with httpx.AsyncClient(timeout=30) as client:
        resp = await _request_with_retry(client, "GET", f"{CT_API_BASE}/studies", params=params)
        if resp.status_code != 200:
            return
        mapping = {}
        for s in resp.json().get("studies", []):
            p = s.get("protocolSection", {})
            sponsor = (p.get("sponsorCollaboratorsModule", {})
                       .get("leadSponsor") or {}).get("name", "")
            if not sponsor:
                continue
            for inv in (p.get("armsInterventionsModule", {})
                        .get("interventions") or []):
                inv_name = inv.get("name", "")
                if inv_name:
                    mapping.setdefault(inv_name.lower(), sponsor)
    for drug in drugs:
        key = drug.name.lower()
        if key in mapping:
            drug.company = mapping[key]
        else:
            for inv_name, sponsor in mapping.items():
                if key in inv_name or inv_name in key:
                    drug.company = sponsor
                    break

# ─── Rule-based Report Builder ─────────────────────────────────────────────


def build_report(
    target_name: str,
    target_detail: TargetDetail,
    papers: List[Paper],
    trials: List[ClinicalTrial],
    drugs: List[DrugInfo],
) -> ReportContent:
    report = ReportContent(
        target_name=target_name,
        target_detail=target_detail,
        papers=papers, trials=trials, drugs=drugs,
    )

    if target_detail:
        parts = [f"靶点名称: {target_name}"]
        if target_detail.gene_symbol:
            parts.append(f"基因符号: {target_detail.gene_symbol}")
        if target_detail.protein_class:
            parts.append(f"蛋白类别: {target_detail.protein_class}")
        if target_detail.description:
            parts.append(f"功能描述: {target_detail.description}")
        if target_detail.related_diseases:
            parts.append(f"相关疾病: {', '.join(target_detail.related_diseases)}")
        report.target_overview = "\n".join(parts)
    else:
        report.target_overview = f"靶点 {target_name} 的基本信息暂未从 Open Targets 获取到。"

    year_count = _count_by_year(papers)
    total = len(papers)
    if year_count:
        years = list(year_count.keys())
        trend = "上升" if len(years) >= 2 and year_count[years[-1]] > year_count[years[0]] else "波动"
        now_year = datetime.now().year
        recent = [str(y) for y in range(now_year - 2, now_year + 1)]
        report.research_progress = (
            f"共检索到 {total} 篇相关文献，覆盖 {len(year_count)} 个年份"
            f"（{years[0]}–{years[-1]}），整体呈 {trend} 趋势。"
            f"近三年（{recent[0]}–{recent[-1]}）共发表 {sum(year_count.get(y, 0) for y in recent)} 篇。"
        )
    else:
        report.research_progress = f"共检索到 {total} 篇相关文献。"

    phase_count = _count_trials_by_phase(trials)
    if phase_count:
        phases_str = "、".join(f"{k} {v}项" for k, v in sorted(phase_count.items()))
        report.clinical_landscape = (
            f"共 {len(trials)} 项临床试验，阶段分布：{phases_str}。"
            f"覆盖适应症：{_collect_conditions(trials)}。"
        )
    else:
        report.clinical_landscape = f"共 {len(trials)} 项临床试验。"

    top_papers = sorted(papers, key=lambda p: p.year or "0", reverse=True)[:5]
    findings = []
    for i, p in enumerate(top_papers, 1):
        findings.append(f"{i}. {p.title} — {p.journal} ({p.year})")
    report.key_findings = "\n".join(findings) if findings else "暂无文献数据。"

    if drugs:
        approved = [d for d in drugs if "approv" in d.phase.lower() or "上市" in d.phase]
        dev = [d for d in drugs if d not in approved]
        outlook_parts = [f"已有 {len(approved)} 个药物获批，{len(dev)} 个在研。"]
        if approved:
            outlook_parts.append(
                f"已获批药物：{', '.join(d.name for d in approved)}。")
        report.future_outlook = " ".join(outlook_parts)
    else:
        report.future_outlook = "暂未检索到针对该靶点的药物信息。"

    return report


def _count_by_year(papers) -> Dict[str, int]:
    years = {}
    for p in papers:
        if p.year and p.year.isdigit():
            years[p.year] = years.get(p.year, 0) + 1
    return dict(sorted(years.items()))


def _count_trials_by_phase(trials) -> Dict[str, int]:
    phases = {}
    for t in trials:
        p = t.phase or "未明确"
        phases[p] = phases.get(p, 0) + 1
    return phases


def _collect_conditions(trials) -> str:
    conds = set()
    for t in trials:
        for c in t.conditions:
            conds.add(c)
    return "、".join(list(conds)[:8]) if conds else "—"


# ─── LLM Report Enhancement ────────────────────────────────────────────────


async def _enhance_report_with_llm(
    api_key: str,
    report: ReportContent,
    progress_callback=None,
) -> ReportContent:
    """Override rule-based report sections with LLM-generated content."""

    def _build_context(field: str) -> str:
        parts = [f"靶点名称: {report.target_name}"]
        d = report.target_detail
        if d:
            if d.gene_symbol:
                parts.append(f"基因符号: {d.gene_symbol}")
            if d.protein_class:
                parts.append(f"蛋白类别: {d.protein_class}")
            if d.description:
                parts.append(f"功能描述: {d.description}")
            if d.related_diseases:
                parts.append(f"相关疾病: {', '.join(d.related_diseases)}")
            if d.synonyms:
                parts.append(f"别名: {'; '.join(d.synonyms[:10])}")

        if field == "research_progress" or field == "key_findings":
            parts.append(f"\n文献总数: {len(report.papers)}")
            year_count = _count_by_year(report.papers)
            if year_count:
                years_str = ", ".join(f"{y}({c}篇)" for y, c in year_count.items())
                parts.append(f"年份分布: {years_str}")
            if report.papers:
                parts.append("\n文献标题（最近10篇）:")
                for i, p in enumerate(report.papers[:10], 1):
                    parts.append(f"  {i}. {p.title} ({p.journal}, {p.year})")

        if field in ("clinical_landscape", "future_outlook"):
            parts.append(f"\n临床试验总数: {len(report.trials)}")
            phase_count = _count_trials_by_phase(report.trials)
            if phase_count:
                parts.append(f"阶段分布: {dict(phase_count)}")
            conds = _collect_conditions(report.trials)
            if conds and conds != "—":
                parts.append(f"适应症: {conds}")

        if field == "future_outlook":
            parts.append(f"\n药物总数: {len(report.drugs)}")
            for d in report.drugs:
                parts.append(f"  - {d.name} ({d.phase}), 公司: {d.company or '—'}, MoA: {d.mechanism_of_action}")

        return "\n".join(parts)

    sections = [
        ("target_overview", "target_overview"),
        ("research_progress", "research_progress"),
        ("clinical_landscape", "clinical_landscape"),
        ("key_findings", "key_findings"),
        ("future_outlook", "future_outlook"),
    ]

    for idx, (field, prompt_key) in enumerate(sections):
        context = _build_context(field)
        prompt = _LLM_PROMPTS[prompt_key].format(context=context)
        if progress_callback:
            progress_callback(70 + idx * 5, f"AI 正在分析 {field}...")
        text = await _llm_chat(prompt, api_key)
        if text:
            setattr(report, field, text)

    # ── AI networked intelligence ──
    d = report.target_detail
    web_prompt = _LLM_PROMPTS["web_summary"].format(
        target_name=report.target_name,
        gene_symbol=d.gene_symbol if d and d.gene_symbol else "—",
        protein_class=d.protein_class if d and d.protein_class else "—",
    )
    if progress_callback:
        progress_callback(96, "AI 正在联网搜索全网情报...")
    web_text = await _llm_chat_with_search(web_prompt, api_key)
    if web_text:
        report.web_summary = web_text

    return report


# ─── PPT Style Definitions ──────────────────────────────────────────────────

DARK_BLUE = RGBColor(0x1B, 0x4F, 0x72)
MID_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)
MED_GRAY = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FONT = "Microsoft YaHei"


def set_slide_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_bottom_bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), SLIDE_WIDTH, Inches(0.5))
    s.fill.solid()
    s.fill.fore_color.rgb = DARK_BLUE
    s.line.fill.background()


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tb


def add_multi_text(slide, left, top, width, height, items, font_size=14,
                   color=DARK_GRAY, line_spacing=1.5):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    for i, item in enumerate(items):
        p = tb.text_frame.paragraphs[0] if i == 0 else tb.text_frame.add_paragraph()
        p.text = item
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return tb


def add_slide_title(slide, title_text, subtitle_text=""):
    add_slide_bg(slide)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    add_textbox(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.8),
                title_text, font_size=28, bold=True, color=WHITE)
    if subtitle_text:
        add_textbox(slide, Inches(0.8), Inches(0.75), Inches(11), Inches(0.4),
                    subtitle_text, font_size=14, color=LIGHT_BLUE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(1.5), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_RED
    line.line.fill.background()
    add_bottom_bar(slide)


def add_slide_bg(slide):
    set_slide_bg(slide, WHITE)
    a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), SLIDE_HEIGHT)
    a.fill.solid()
    a.fill.fore_color.rgb = LIGHT_BLUE
    a.line.fill.background()


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    ts = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for j, h in enumerate(headers):
        c = table.cell(0, j)
        c.text = h
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT
            p.alignment = PP_ALIGN.CENTER
        c.fill.solid()
        c.fill.fore_color.rgb = DARK_BLUE

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = table.cell(i + 1, j)
            c.text = str(val)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK_GRAY
                p.font.name = FONT
            c.fill.solid()
            c.fill.fore_color.rgb = LIGHT_GRAY if i % 2 == 1 else WHITE
    return table

# ─── PPT Generator ─────────────────────────────────────────────────────────

_FONT_CANDIDATES = ["Microsoft YaHei", "SimHei", "WenQuanYi Micro Hei", "Noto Sans CJK SC",
                     "Noto Sans SC", "AR PL UMing CN", "DejaVu Sans"]
_available_fonts = {f.name for f in matplotlib.font_manager.fontManager.ttflist}
for _f in _FONT_CANDIDATES:
    if _f in _available_fonts:
        plt.rcParams["font.family"] = _f
        break
else:
    plt.rcParams["font.family"] = "sans-serif"
plt.rcParams["axes.unicode_minus"] = False


def _quiet_savefig(buf, **kwargs):
    """Save figure suppressing C-level libpng stderr warnings."""
    try:
        old = os.dup(2)
        null = os.open(os.devnull, os.O_WRONLY)
        os.dup2(null, 2)
        os.close(null)
        try:
            plt.savefig(buf, **kwargs)
        finally:
            os.dup2(old, 2)
            os.close(old)
    except Exception:
        plt.savefig(buf, **kwargs)


def _render_chart(years_dict: dict, title: str, chart_type: str = "bar") -> bytes:
    fig, ax = plt.subplots(figsize=(8, 4))
    if chart_type == "pie":
        ax.pie(years_dict.values(), labels=years_dict.keys(), autopct="%1.1f%%",
               colors=["#2E86C1", "#1B4F72", "#E74C3C", "#F39C12", "#27AE60"])
        ax.set_title(title, fontsize=14, pad=20)
    else:
        ax.bar(years_dict.keys(), years_dict.values(), color="#2E86C1",
               edgecolor="#1B4F72", linewidth=0.8)
        ax.set_title(title, fontsize=14, fontweight="bold", color="#1B4F72", pad=15)
        ax.set_xlabel("Year", fontsize=11)
        ax.set_ylabel("# Papers", fontsize=11)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.tick_params(colors="#2C3E50")

    buf = io.BytesIO()
    kwargs = {"format": "png", "dpi": 150, "bbox_inches": "tight",
              "pil_kwargs": {"icc_profile": None}}
    if chart_type == "bar":
        kwargs.update({"transparent": False, "facecolor": "white"})
    else:
        kwargs["transparent"] = True
    _quiet_savefig(buf, **kwargs)
    plt.close()
    buf.seek(0)
    return buf.getvalue()


def _slice(text: str, n: int = 80) -> str:
    return text if len(text) <= n else text[:n - 3] + "..."


def generate_ppt(report: ReportContent) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _title(prs, report)
    _toc(prs, report)
    _overview(prs, report)
    _lit_trend(prs, report)
    _key_papers(prs, report)
    _trial_overview(prs, report)
    _key_trials(prs, report)
    _drug_pipeline(prs, report)
    _milestones(prs, report)
    _web_intel(prs, report)
    _summary(prs, report)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT


def generate_docx(report: ReportContent) -> bytes:
    doc = Document()

    # ── styles ──
    style = doc.styles["Normal"]
    style.font.name = "Microsoft YaHei"
    style.font.size = Pt(10.5)
    style.paragraph_format.space_after = Pt(4)

    def add_heading(text, level=1):
        h = doc.add_heading(text, level=level)
        for run in h.runs:
            run.font.name = "Microsoft YaHei"
        return h

    def add_para(text, bold=False, italic=False, size=None):
        p = doc.add_paragraph()
        r = p.add_run(text)
        r.bold = bold
        r.italic = italic
        r.font.name = "Microsoft YaHei"
        if size:
            r.font.size = Pt(size)
        return p

    def add_table(headers, rows):
        t = doc.add_table(rows=1 + len(rows), cols=len(headers))
        t.style = "Table Grid"
        t.alignment = WD_TABLE_ALIGNMENT.CENTER
        for j, h in enumerate(headers):
            c = t.rows[0].cells[j]
            c.text = h
            for p in c.paragraphs:
                p.runs[0].bold = True if p.runs else False
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                t.rows[i + 1].cells[j].text = str(val)
        return t

    # ── Title ──
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(f"靶点 {report.target_name.upper()} 研究进展与临床分析报告")
    r.bold = True
    r.font.size = Pt(22)
    r.font.name = "Microsoft YaHei"

    add_para(f"生成日期: {datetime.now().strftime('%Y-%m-%d %H:%M')}", size=11)
    add_para("数据来源: PubMed · ClinicalTrials.gov · Open Targets", size=10)
    doc.add_page_break()

    # ── 1. Target Overview ──
    add_heading("一、靶点概述与背景", 1)
    d = report.target_detail
    if d:
        add_para(f"靶点名称: {report.target_name}")
        add_para(f"基因符号: {d.gene_symbol or '—'}")
        add_para(f"蛋白类别: {d.protein_class or '—'}")
        if d.synonyms:
            add_para(f"别名: {'; '.join(d.synonyms)}")
        if d.description:
            add_para(f"功能描述: {d.description}")
        if d.related_diseases:
            add_para(f"相关疾病: {'; '.join(d.related_diseases)}")
    if report.target_overview:
        for line in report.target_overview.split("\n"):
            add_para(line)
    doc.add_page_break()

    # ── 2. Literature ──
    add_heading("二、文献检索结果", 1)
    add_para(f"共检索到 {len(report.papers)} 篇相关文献", bold=True)
    if report.research_progress:
        add_para(report.research_progress)

    if report.papers:
        add_heading(f"全部文献列表（共 {len(report.papers)} 篇）", 2)
        for i, p in enumerate(report.papers, 1):
            add_heading(f"第 {i} 篇", 3)
            add_para(f"标题: {p.title}")
            if p.authors:
                add_para(f"作者: {'; '.join(p.authors)}")
            add_para(f"期刊: {p.journal}")
            add_para(f"年份: {p.year}")
            if p.pmid:
                add_para(f"PMID: {p.pmid}")
            if p.doi:
                add_para(f"DOI: {p.doi}")
            if p.keywords:
                add_para(f"关键词: {'; '.join(p.keywords)}")
            if p.abstract:
                add_para(f"摘要: {p.abstract}")
            doc.add_paragraph()  # spacing

    doc.add_page_break()

    # ── 3. Clinical Trials ──
    add_heading("三、临床试验数据", 1)
    add_para(f"共检索到 {len(report.trials)} 项相关临床试验", bold=True)
    if report.clinical_landscape:
        add_para(report.clinical_landscape)

    if report.trials:
        add_heading("全部临床试验列表", 2)
        headers = ["NCT ID", "标题", "阶段", "状态", "适应症", "干预措施",
                    "申办方", "开始日期", "完成日期", "入组人数", "简要摘要"]
        rows = []
        for t in report.trials:
            rows.append([
                t.nct_id, t.title, t.phase, t.status,
                ", ".join(t.conditions),
                ", ".join(t.interventions),
                t.sponsor, t.start_date, t.completion_date,
                str(t.enrollment), t.brief_summary,
            ])
        add_table(headers, rows)

    doc.add_page_break()

    # ── 4. Drug Pipeline ──
    add_heading("四、药物研发管线", 1)
    add_para(f"共检索到 {len(report.drugs)} 个靶向药物", bold=True)

    if report.drugs:
        add_heading("全部药物列表", 2)
        headers = ["药物名称", "公司", "作用机制", "临床阶段", "适应症"]
        rows = [[d.name, d.company or "—", d.mechanism_of_action, d.phase, d.disease]
                for d in report.drugs]
        add_table(headers, rows)

        approved = [d for d in report.drugs if "approv" in d.phase.lower()]
        if approved:
            add_para("")
            add_para(f"已获批药物: {'; '.join(d.name for d in approved)}", bold=True)

    doc.add_page_break()

    # ── 5. Key Milestones ──
    add_heading("五、关键里程碑", 1)
    year_count = {}
    for p in report.papers:
        if p.year and p.year.isdigit():
            year_count[p.year] = year_count.get(p.year, 0) + 1

    milestones = []
    for year in sorted(year_count.keys())[-5:]:
        milestones.append(f"▸ {year}年 — 发表 {year_count[year]} 篇相关研究文献")
    if milestones:
        for m in milestones:
            add_para(m)

    phases = sorted(set(t.phase for t in report.trials if t.phase))
    if phases:
        add_para(f"▸ 临床试验覆盖阶段: {', '.join(phases)}")

    active = sum(1 for t in report.trials
                 if t.status in ("ACTIVE", "RECRUITING", "NOT_YET_RECRUITING"))
    if active:
        add_para(f"▸ 目前有 {active} 项活跃临床试验正在进行")

    completed = sum(1 for t in report.trials if t.status == "COMPLETED")
    if completed:
        add_para(f"▸ 已完成 {completed} 项临床试验")

    approved = sum(1 for d in report.drugs if "approv" in d.phase.lower())
    if approved:
        approved_names = [d.name for d in report.drugs if "approv" in d.phase.lower()]
        add_para(f"▸ 已有 {approved} 个靶向药物获批上市: {'; '.join(approved_names)}")

    if report.web_summary:
        doc.add_page_break()
        add_heading("六、AI联网情报分析", 1)
        add_para("以下内容由AI通过联网搜索实时生成，仅供参考。", size=10, italic=True)
        for line in report.web_summary.split("\n"):
            line = line.strip()
            if line:
                add_para(line)

    doc.add_page_break()

    # ── 7. Summary ──
    add_heading("七、总结与展望", 1)
    if report.research_progress:
        add_para("【文献概况】", bold=True)
        add_para(report.research_progress)
    if report.future_outlook:
        add_para("【药物展望】", bold=True)
        add_para(report.future_outlook)

    add_para("")
    add_para(f"报告生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}", size=9)
    add_para("数据来源: PubMed · ClinicalTrials.gov · Open Targets", size=9)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _title(prs, r):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, DARK_BLUE)
    add_textbox(s, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                f"靶点 {r.target_name.upper()}\n研究进展与临床分析报告",
                font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
                f"生成日期: {datetime.now().strftime('%Y-%m-%d')}",
                font_size=16, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
                "数据来源: PubMed · ClinicalTrials.gov · Open Targets",
                font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    l = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.8), Inches(4), Inches(0.04))
    l.fill.solid()
    l.fill.fore_color.rgb = ACCENT_RED
    l.line.fill.background()


def _toc(prs, r):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "目  录", "CONTENTS")
    items = ["01  靶点概述与背景", "02  文献发表趋势", "03  核心文献解读",
             "04  临床试验概况", "05  关键临床试验", "06  药物研发管线",
             "07  关键里程碑", "08  总结与展望"]
    if r.web_summary:
        items.insert(-1, "09  AI联网情报")
    add_multi_text(s, Inches(2), Inches(1.8), Inches(9), Inches(5),
                   items,
                   font_size=22, color=DARK_BLUE, line_spacing=2.0)


def _overview(prs, r):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "靶点概述与背景", f"{r.target_name.upper()} 基本信息")
    d = r.target_detail
    L, R, T, W = Inches(0.8), Inches(7.0), Inches(1.6), Inches(5.5)

    if d:
        add_multi_text(s, L, T, W, Inches(3), [
            f"靶点名称: {r.target_name}",
            f"基因符号: {d.gene_symbol or '—'}",
            f"蛋白类别: {d.protein_class or '—'}",
            f"别    名: {', '.join(d.synonyms[:6]) if d.synonyms else '—'}",
        ], font_size=14, color=DARK_GRAY)

    add_multi_text(s, R, T, W, Inches(5),
                   (r.target_overview or "暂无靶点概述信息。").split("\n"),
                   font_size=13, color=DARK_GRAY)

    if d and d.related_diseases:
        add_textbox(s, L, Inches(4.5), W, Inches(2),
                    "相关疾病: " + ", ".join(d.related_diseases[:8]),
                    font_size=12, color=MED_GRAY)


def _lit_trend(prs, r):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "文献发表趋势", "近年研究热度与分布")

    years = _count_by_year(r.papers)
    if years:
        img = _render_chart(years, f"{r.target_name.upper()} - Publication Trend")
        s.shapes.add_picture(io.BytesIO(img), Inches(0.8), Inches(1.6), Inches(7.5), Inches(4.5))

    if r.research_progress:
        add_textbox(s, Inches(8.8), Inches(1.8), Inches(4), Inches(4.5),
                    r.research_progress, font_size=12, color=DARK_GRAY)


def _add_table_slides(prs, items, title_base, subtitle_base, headers, row_fn, col_widths, max_rows=12):
    if not items:
        return
    chunks = [items[i:i + max_rows] for i in range(0, len(items), max_rows)]
    for idx, chunk in enumerate(chunks):
        title = title_base if idx == 0 else f"{title_base}（续{idx}）"
        subtitle = subtitle_base if idx == 0 else f"{subtitle_base}（续）"
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_title(s, title, subtitle)
        rows = [row_fn(i + 1 + idx * max_rows, item) for i, item in enumerate(chunk)]
        add_table(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.5),
                  headers, rows, col_widths)


def _key_papers(prs, r):
    papers = r.papers
    has_extra = bool(r.key_findings)
    if not papers:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_title(s, "核心文献解读", "暂无文献数据")
        return
    chunks = [papers[i:i + 12] for i in range(0, len(papers), 12)]
    for idx, chunk in enumerate(chunks):
        title = "核心文献解读" if idx == 0 else f"核心文献解读（续{idx}）"
        sub = f"共检索到 {len(r.papers)} 篇相关文献" if idx == 0 else ""
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_title(s, title, sub)
        rows = [[str(i + 1 + idx * 12), _slice(p.title, 60), _slice(p.journal, 25), p.year]
                for i, p in enumerate(chunk)]
        add_table(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(3.5),
                  ["#", "标题", "期刊", "年份"], rows,
                  [Inches(0.5), Inches(6.5), Inches(3), Inches(0.8)])
        if idx == 0 and r.key_findings:
            add_textbox(s, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.5),
                        f"代表性文献:\n{r.key_findings[:400]}", font_size=12, color=DARK_BLUE)


def _trial_overview(prs, r):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "临床试验概况", f"共 {len(r.trials)} 项相关临床试验")

    phases = _count_trials_by_phase(r.trials)
    if phases:
        img = _render_chart(phases, f"{r.target_name.upper()} - Trial Phase Distribution", "pie")
        s.shapes.add_picture(io.BytesIO(img), Inches(0.8), Inches(1.6), Inches(5), Inches(4.5))

    add_textbox(s, Inches(6.5), Inches(1.8), Inches(5.5), Inches(4.5),
                r.clinical_landscape or "暂无临床试验总结信息。",
                font_size=13, color=DARK_GRAY)


def _key_trials(prs, r):
    _add_table_slides(prs, r.trials, "关键临床试验一览", "代表性临床试验详情",
                      ["NCT ID", "Phase", "Status", "适应症", "干预措施"],
                      lambda i, t: [t.nct_id, t.phase, t.status,
                                    _slice(", ".join(t.conditions[:2]), 25),
                                    _slice(", ".join(t.interventions[:1]), 25)],
                      [Inches(1.8), Inches(1.0), Inches(1.5), Inches(3.5), Inches(3.5)],
                      max_rows=12)


def _drug_pipeline(prs, r):
    if r.drugs:
        _add_table_slides(prs, r.drugs, "药物研发管线", "靶向药物的研发现状",
                          ["药物名称", "公司", "作用机制", "临床阶段", "适应症"],
                          lambda i, d: [d.name, _slice(d.company or "—", 18),
                                        _slice(d.mechanism_of_action, 25), d.phase,
                                        _slice(d.disease, 18)],
                          [Inches(2.0), Inches(1.8), Inches(3.5), Inches(1.5), Inches(2.5)],
                          max_rows=12)
    else:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_title(s, "药物研发管线", "靶向药物的研发现状")
        add_textbox(s, Inches(2), Inches(3), Inches(9), Inches(1.5),
                    "暂未从 Open Targets 获取到药物信息。\n"
                    "该靶点可能暂无已知靶向药物，或数据尚未收录。",
                    font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


def _milestones(prs, r):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "关键里程碑", f"{r.target_name.upper()} 研究发展历程")

    m = []
    for year, count in list(_count_by_year(r.papers).items())[-5:]:
        m.append(f"▸ {year}年 — 发表 {count} 篇相关研究文献")

    phases = sorted(set(t.phase for t in r.trials if t.phase))
    if phases:
        m.append(f"▸ 临床试验覆盖阶段: {', '.join(phases)}")

    active = sum(1 for t in r.trials if t.status in ("ACTIVE", "RECRUITING", "NOT_YET_RECRUITING"))
    if active:
        m.append(f"▸ 目前有 {active} 项活跃临床试验正在进行")

    completed = sum(1 for t in r.trials if t.status == "COMPLETED")
    if completed:
        m.append(f"▸ 已完成 {completed} 项临床试验")

    approved = sum(1 for d in r.drugs if "approv" in d.phase.lower() or "上市" in d.phase)
    if approved:
        m.append(f"▸ 已有 {approved} 个靶向药物获批上市")

    if not m:
        m = ["暂无足够的里程碑数据"]

    add_multi_text(s, Inches(1.5), Inches(1.8), Inches(10), Inches(4.5),
                   m, font_size=16, color=DARK_BLUE, line_spacing=2.0)

    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.65), Inches(10), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MID_BLUE
    bar.line.fill.background()


def _web_intel(prs, r):
    if not r.web_summary:
        return
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "AI联网情报分析", "网络聚合 · 实时信息")
    text = r.web_summary.replace("\n", "\n\n")
    add_textbox(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(6.0),
                text, font_size=12, color=DARK_GRAY)


def _summary(prs, r):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(s, "总结与展望", "Conclusion & Outlook")

    add_textbox(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(2.5),
                f"【文献概况】\n{r.research_progress[:300]}",
                font_size=13, color=DARK_GRAY)
    add_textbox(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.0),
                f"【药物展望】\n{r.future_outlook[:300]}",
                font_size=13, color=DARK_BLUE)
    add_textbox(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.5),
                f"报告生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}  |  "
                f"数据来源: PubMed · ClinicalTrials.gov · Open Targets",
                font_size=10, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# ─── PyQt5 App ─────────────────────────────────────────────────────────────


class PipelineWorker(QThread):
    progress = pyqtSignal(int, str)
    finished = pyqtSignal(object, bytes, bytes, list, list)
    error = pyqtSignal(str)

    def __init__(self, target_name, use_llm=False, api_key=""):
        super().__init__()
        self.target_name = target_name
        self.use_llm = use_llm
        self.api_key = api_key

    def _progress(self, value, text):
        self.progress.emit(value, text)

    def run(self):
        try:
            async def _run():
                self.progress.emit(10, "正在搜索相关文献 (PubMed)...")
                papers = await search_papers(self.target_name)

                self.progress.emit(25, "正在查询临床试验 (ClinicalTrials.gov)...")
                trials = await search_trials(self.target_name)

                self.progress.emit(40, "正在获取靶点信息 (Open Targets)...")
                target_detail = await get_target_detail(self.target_name)

                self.progress.emit(50, "正在查询药物信息 (Open Targets)...")
                drugs = await get_drugs(self.target_name)

                self.progress.emit(60, "正在匹配药物公司与临床数据...")
                await _enrich_drugs_with_company(drugs)

                self.progress.emit(65, "正在生成基础报告...")
                report = build_report(self.target_name, target_detail, papers, trials, drugs)

                if self.use_llm and self.api_key:
                    report = await _enhance_report_with_llm(
                        self.api_key, report, progress_callback=self._progress,
                    )

                self.progress.emit(85, "正在生成 PPT...")
                ppt_bytes = generate_ppt(report)

                self.progress.emit(92, "正在生成 Word 文档...")
                docx_bytes = generate_docx(report)

                self.progress.emit(100, "完成!")
                self.finished.emit(report, ppt_bytes, docx_bytes, papers, trials)

            asyncio.run(_run())
        except Exception as e:
            import traceback
            self.error.emit(f"{str(e)}\n\n{traceback.format_exc()}")


class CollapsibleSection(QWidget):
    def __init__(self, title, parent=None):
        super().__init__(parent)
        self.toggle_btn = QPushButton(f"▶ {title}")
        self.toggle_btn.setCheckable(True)
        self.toggle_btn.setChecked(False)
        self.toggle_btn.clicked.connect(self._toggle)
        self.toggle_btn.setStyleSheet("""
            QPushButton {
                text-align: left; padding: 8px 12px;
                background: #f0f2f6; border: 1px solid #ddd;
                border-radius: 6px; font-size: 14px;
            }
            QPushButton:checked {
                background: #e8f0fe; border-color: #2E86C1;
            }
        """)
        self.content = QWidget()
        self.content.setVisible(False)
        self.content_layout = QVBoxLayout(self.content)
        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(2)
        layout.addWidget(self.toggle_btn)
        layout.addWidget(self.content)

    def _toggle(self):
        checked = self.toggle_btn.isChecked()
        self.content.setVisible(checked)
        prefix = "▼" if checked else "▶"
        self.toggle_btn.setText(f"{prefix} {self.toggle_btn.text()[2:]}")

    def setContent(self, widget):
        while self.content_layout.count():
            item = self.content_layout.takeAt(0)
            w = item.widget()
            if w:
                w.deleteLater()
        self.content_layout.addWidget(widget)


class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("文献总结 · PPT & Word 生成器")
        self.setMinimumSize(700, 600)
        self.resize(900, 760)

        central = QWidget()
        self.setCentralWidget(central)
        outer = QVBoxLayout(central)
        outer.setContentsMargins(30, 20, 30, 20)

        # ── Title ──
        title = QLabel("📚  文献总结 · PPT & Word 生成器")
        title.setStyleSheet("""
            font-size: 22px; font-weight: 700; color: #1B4F72;
            padding-bottom: 6px; border-bottom: 3px solid #2E86C1;
        """)
        title.setAlignment(Qt.AlignCenter)
        outer.addWidget(title)

        subtitle = QLabel("输入靶点名称，自动检索文献、临床试验、药物信息，一键生成专业报告")
        subtitle.setStyleSheet("color: #7F8C8D; font-size: 13px;")
        subtitle.setAlignment(Qt.AlignCenter)
        outer.addWidget(subtitle)
        outer.addSpacing(12)

        # ── Input ──
        input_row = QHBoxLayout()
        input_label = QLabel("🔬 靶点名称:")
        input_label.setStyleSheet("font-size: 14px; font-weight: 600;")
        self.input_field = QLineEdit()
        self.input_field.setPlaceholderText("例如: PD-1, EGFR, HER2, KRAS, CD19...")
        self.input_field.setStyleSheet("""
            QLineEdit { padding: 8px 12px; font-size: 14px;
                        border: 2px solid #ddd; border-radius: 6px; }
            QLineEdit:focus { border-color: #2E86C1; }
        """)
        self.start_btn = QPushButton("🚀 开始生成")
        self.start_btn.setCursor(Qt.PointingHandCursor)
        self.start_btn.setStyleSheet("""
            QPushButton {
                padding: 8px 24px; font-size: 14px; font-weight: 600;
                background: #2E86C1; color: white; border: none; border-radius: 6px;
            }
            QPushButton:hover { background: #2874A6; }
            QPushButton:disabled { background: #BDC3C7; }
        """)
        self.start_btn.clicked.connect(self.start_pipeline)
        input_row.addWidget(input_label)
        input_row.addWidget(self.input_field, 1)
        input_row.addWidget(self.start_btn)
        outer.addLayout(input_row)
        outer.addSpacing(6)

        # ── LLM options ──
        llm_row = QHBoxLayout()
        self.llm_check = QPushButton("☐ AI 智能分析")
        self.llm_check.setCheckable(True)
        self.llm_check.setChecked(bool(ZHIPUAI_API_KEY))
        self.llm_check.setCursor(Qt.PointingHandCursor)
        self.llm_check.setStyleSheet("""
            QPushButton {
                text-align: left; padding: 6px 12px; font-size: 13px;
                background: #f8f9fa; border: 1px solid #ddd; border-radius: 6px;
            }
            QPushButton:checked {
                background: #e8f0fe; border-color: #2E86C1;
                font-weight: 600;
            }
        """)
        self.llm_check.toggled.connect(self._toggle_llm_input)
        self.llm_key = QLineEdit()
        self.llm_key.setPlaceholderText("智谱 API Key（留空则使用 .env 中的 ZHIPUAI_API_KEY）")
        self.llm_key.setText(ZHIPUAI_API_KEY or "")
        self.llm_key.setStyleSheet("""
            QLineEdit { padding: 6px 10px; font-size: 12px;
                        border: 1px solid #ddd; border-radius: 6px; }
            QLineEdit:focus { border-color: #2E86C1; }
        """)
        self.llm_key.setEchoMode(QLineEdit.Password)
        self.llm_key.setVisible(self.llm_check.isChecked())
        llm_row.addWidget(self.llm_check)
        llm_row.addWidget(self.llm_key, 1)
        outer.addLayout(llm_row)
        outer.addSpacing(10)

        # ── Progress ──
        self.progress_bar = QProgressBar()
        self.progress_bar.setValue(0)
        self.progress_bar.setTextVisible(True)
        self.progress_bar.setFixedHeight(26)
        self.progress_bar.setStyleSheet("""
            QProgressBar {
                border: 1px solid #ddd; border-radius: 6px;
                text-align: center; font-size: 12px; background: #f0f0f0;
            }
            QProgressBar::chunk { background: #2E86C1; border-radius: 5px; }
        """)
        self.progress_label = QLabel("")
        self.progress_label.setStyleSheet("color: #555; font-size: 13px;")
        self.progress_label.setAlignment(Qt.AlignCenter)
        self.progress_widget = QWidget()
        prog_l = QVBoxLayout(self.progress_widget)
        prog_l.setContentsMargins(0, 0, 0, 0)
        prog_l.addWidget(self.progress_bar)
        prog_l.addWidget(self.progress_label)
        self.progress_widget.setVisible(False)
        outer.addWidget(self.progress_widget)
        outer.addSpacing(10)

        # ── Scroll area for results ──
        self.results_scroll = QScrollArea()
        self.results_scroll.setWidgetResizable(True)
        self.results_scroll.setVisible(False)
        self.results_scroll.setStyleSheet("QScrollArea { border: none; }")
        self.results_content = QWidget()
        self.results_layout = QVBoxLayout(self.results_content)
        self.results_scroll.setWidget(self.results_content)
        outer.addWidget(self.results_scroll, 1)

        # ── Footer ──
        footer = QLabel(
            '数据来源: <a href="https://pubmed.ncbi.nlm.nih.gov/" style="color:#2E86C1;">PubMed</a>'
            ' · <a href="https://clinicaltrials.gov/" style="color:#2E86C1;">ClinicalTrials.gov</a>'
            ' · <a href="https://www.opentargets.org/" style="color:#2E86C1;">Open Targets</a>'
            ' &nbsp;|&nbsp; PPT 引擎: python-pptx'
        )
        footer.setOpenExternalLinks(True)
        footer.setAlignment(Qt.AlignCenter)
        footer.setStyleSheet("color: #7F8C8D; font-size: 12px; padding: 10px 0;")
        outer.addWidget(footer)

    # ── slots ──

    def _toggle_llm_input(self, checked):
        self.llm_key.setVisible(checked)
        prefix = "☑" if checked else "☐"
        self.llm_check.setText(f"{prefix} AI 智能分析")
        if checked and not self.llm_key.text() and not ZHIPUAI_API_KEY:
            QMessageBox.information(self, "提示",
                "请在输入框中填写智谱 API Key，\n"
                "或在 .env 文件中设置 ZHIPUAI_API_KEY。")
            self.llm_key.setFocus()

    def start_pipeline(self):
        target_name = self.input_field.text().strip()
        if not target_name:
            QMessageBox.warning(self, "提示", "请输入靶点名称")
            return

        use_llm = self.llm_check.isChecked()
        api_key = self.llm_key.text().strip() or ZHIPUAI_API_KEY or ""
        if use_llm and not api_key:
            QMessageBox.warning(self, "提示", "请填写智谱 API Key 或在 .env 中设置 ZHIPUAI_API_KEY")
            return

        self.start_btn.setEnabled(False)
        self.input_field.setEnabled(False)
        self.progress_widget.setVisible(True)
        self.results_scroll.setVisible(False)
        self.progress_bar.setValue(0)
        self.progress_label.setText("初始化...")
        self._clear_results()

        self.worker = PipelineWorker(target_name, use_llm=use_llm, api_key=api_key)
        self.worker.progress.connect(self._on_progress)
        self.worker.finished.connect(self._on_finished)
        self.worker.error.connect(self._on_error)
        self.worker.start()

    def _clear_results(self):
        for i in reversed(range(self.results_layout.count())):
            item = self.results_layout.itemAt(i)
            w = item.widget()
            if w:
                w.deleteLater()

    def _on_progress(self, value, text):
        self.progress_bar.setValue(value)
        self.progress_label.setText(text)

    def _on_finished(self, report, ppt_bytes, docx_bytes, papers, trials):
        self.progress_bar.setValue(100)
        self.progress_label.setText("完成!")
        self.start_btn.setEnabled(True)
        self.input_field.setEnabled(True)

        self._report = report
        self._ppt_bytes = ppt_bytes
        self._docx_bytes = docx_bytes
        self._papers = papers
        self._trials = trials

        self._build_results(report, len(papers), len(trials), len(report.drugs))
        self.results_scroll.setVisible(True)
        self.results_scroll.verticalScrollBar().setValue(0)

    def _on_error(self, err_msg):
        self.start_btn.setEnabled(True)
        self.input_field.setEnabled(True)
        self.progress_widget.setVisible(False)
        QMessageBox.critical(self, "错误", f"生成过程中出现错误:\n{err_msg}")

    # ── results UI ──

    def _build_results(self, report, n_papers, n_trials, n_drugs):
        success = QLabel("✅ 报告生成完成！  成功获取数据，PPT 和 Word 文档已生成！")
        success.setStyleSheet("""
            background: #d4edda; color: #155724; padding: 12px;
            border-radius: 8px; font-size: 14px; font-weight: 600;
        """)
        success.setAlignment(Qt.AlignCenter)
        self.results_layout.addWidget(success)

        # metrics
        metrics_w = QWidget()
        metrics_l = QHBoxLayout(metrics_w)
        metrics_l.setContentsMargins(0, 10, 0, 10)
        for label, cnt in [
            ("📄 文献数", n_papers),
            ("🧪 临床试验", n_trials),
            ("💊 相关药物", n_drugs),
        ]:
            m = QLabel(f"<b>{label}</b><br><span style='font-size:28px'>{cnt}</span>")
            m.setAlignment(Qt.AlignCenter)
            m.setTextFormat(Qt.RichText)
            m.setStyleSheet("""
                QLabel {
                    background: white; border: 1px solid #e0e0e0;
                    border-radius: 8px; padding: 14px; min-width: 110px;
                }
            """)
            metrics_l.addWidget(m)
        self.results_layout.addWidget(metrics_w)

        # download buttons
        dl_w = QWidget()
        dl_l = QHBoxLayout(dl_w)
        dl_l.setContentsMargins(0, 4, 0, 10)
        ppt_btn = QPushButton("📥 下载 PPT 报告")
        ppt_btn.setCursor(Qt.PointingHandCursor)
        ppt_btn.setStyleSheet("""
            QPushButton {
                padding: 10px 16px; font-size: 14px; font-weight: 600;
                background: #2E86C1; color: white; border: none; border-radius: 6px;
            }
            QPushButton:hover { background: #2874A6; }
        """)
        ppt_btn.clicked.connect(lambda: self._save_file("ppt"))

        docx_btn = QPushButton("📄 下载 Word 文档（完整数据）")
        docx_btn.setCursor(Qt.PointingHandCursor)
        docx_btn.setStyleSheet("""
            QPushButton {
                padding: 10px 16px; font-size: 14px; font-weight: 600;
                background: #27AE60; color: white; border: none; border-radius: 6px;
            }
            QPushButton:hover { background: #229954; }
        """)
        docx_btn.clicked.connect(lambda: self._save_file("docx"))
        dl_l.addWidget(ppt_btn)
        dl_l.addWidget(docx_btn)
        self.results_layout.addWidget(dl_w)

        # report summary
        parts = []
        if report.target_overview:
            parts.append(("靶点概述", report.target_overview))
        if report.research_progress:
            parts.append(("研究进展", report.research_progress))
        if report.clinical_landscape:
            parts.append(("临床概况", report.clinical_landscape))
        if report.key_findings:
            parts.append(("代表性文献", report.key_findings))
        if report.future_outlook:
            parts.append(("药物展望", report.future_outlook))

        if parts:
            sec = CollapsibleSection("📖 查看报告摘要")
            html_parts = "\n".join(
                f"<h3 style='margin-top:12px;'>{title}</h3><p>{text}</p>"
                for title, text in parts
            )
            tb = QTextBrowser()
            tb.setHtml(f"<div style='font-size:13px; line-height:1.7;'>{html_parts}</div>")
            tb.setOpenExternalLinks(True)
            tb.setMinimumHeight(120)
            sec.setContent(tb)
            self.results_layout.addWidget(sec)

        # paper list
        sec2 = CollapsibleSection("📄 查看原始文献列表")
        papers_html = ""
        for i, p in enumerate(self._papers[:20], 1):
            papers_html += (
                f"<p><b>{i}. {p.title}</b><br>"
                f"<span style='color:#666;'>{p.journal} ({p.year}) | PMID: {p.pmid}</span>"
            )
            if p.abstract:
                papers_html += f"<br><span style='color:#888;'>{p.abstract[:200]}...</span>"
            papers_html += "</p>"
        tb2 = QTextBrowser()
        tb2.setHtml(f"<div style='font-size:13px;'>{papers_html}</div>")
        tb2.setOpenExternalLinks(True)
        tb2.setMinimumHeight(120)
        sec2.setContent(tb2)
        self.results_layout.addWidget(sec2)

        # web intelligence
        if self._report.web_summary:
            sec3 = CollapsibleSection("🌐 AI联网情报分析（含实时搜索）")
            tb3 = QTextBrowser()
            html = "<div style='font-size:13px; line-height:1.7;'>"
            for line in self._report.web_summary.split("\n"):
                line = line.strip()
                if line:
                    html += f"<p>{line}</p>"
            html += "</div>"
            tb3.setHtml(html)
            tb3.setOpenExternalLinks(True)
            tb3.setMinimumHeight(150)
            sec3.setContent(tb3)
            self.results_layout.addWidget(sec3)

        self.results_layout.addStretch()

    def _save_file(self, ext):
        ext_map = {
            "ppt": (f"{self._report.target_name.upper()}_文献总结报告.pptx",
                    "PPT 文件 (*.pptx)"),
            "docx": (f"{self._report.target_name.upper()}_文献总结报告.docx",
                     "Word 文件 (*.docx)"),
        }
        name, filt = ext_map[ext]
        data = self._ppt_bytes if ext == "ppt" else self._docx_bytes
        path, _ = QFileDialog.getSaveFileName(self, "保存文件", name, filt)
        if path:
            with open(path, "wb") as f:
                f.write(data)
            QMessageBox.information(self, "保存成功",
                                    f"文件已保存至:\n{path}")


def main():
    app = QApplication(sys.argv)
    app.setStyle("Fusion")
    font = QFont()
    font.setPointSize(10)
    app.setFont(font)
    w = MainWindow()
    w.show()
    sys.exit(app.exec_())


if __name__ == "__main__":
    main()
